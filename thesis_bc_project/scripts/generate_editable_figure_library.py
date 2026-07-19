#!/usr/bin/env python3
"""Generate the Gate V0 editable PowerPoint-native figure library.

The generator reads canonical TSV/log inputs under ``raw_data`` and ``result``.
It never imports the existing F1--F7 raster/vector outputs as slide content.

Run from ``thesis_bc_project`` with the temporary dependencies used by Gate V0::

    PYTHONPATH=/tmp/gate_v0_editable_deps python3 scripts/generate_editable_figure_library.py
"""

from __future__ import annotations

import csv
import io
import math
import re
import statistics
import sys
import zipfile
from collections import defaultdict
from pathlib import Path
from xml.etree import ElementTree as ET

from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
RAW = ROOT / "raw_data"
RESULT = ROOT / "result"
OUT = ROOT / "docs" / "thesis" / "figures" / "editable"
PPTX_PATH = OUT / "thesis_figure_library.pptx"
XLSX_PATH = OUT / "figure_data.xlsx"
MANIFEST_PATH = OUT / "FIGURE_MANIFEST.tsv"
README_PATH = OUT / "README.md"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Arial"

# Color-blind-friendly, stable semantic palette.
C = {
    "gpu": "0072B2",
    "cpu": "7F8C8D",
    "host": "A7B0B5",
    "um": "009E73",
    "pure": "E69F00",
    "chunked": "2E8B57",
    "pathmerge": "7B2CBF",
    "failure": "C23B22",
    "neutral": "343A40",
    "light": "F4F6F8",
    "white": "FFFFFF",
    "pale_blue": "DDEFF8",
    "pale_teal": "DDF3EC",
    "pale_orange": "FCECCB",
    "pale_green": "E1F2E8",
    "pale_purple": "EEE4F6",
    "pale_red": "F7DEDA",
}

GRAPH_ORDER = ["email-EuAll", "roadNet-PA", "roadNet-TX", "roadNet-CA"]
GRAPH_SHORT = {"email-EuAll": "email", "roadNet-PA": "PA", "roadNet-TX": "TX", "roadNet-CA": "CA"}

SOURCES = {
    "t2": "result/tables/thesis/T2_main_performance.tsv",
    "comparison": "result/main_performance/proposed_vs_pathmerge/comparison.tsv",
    "t3": "result/tables/thesis/T3_ablation_summary.tsv",
    "t4": "result/tables/thesis/T4_memory_scalability.tsv",
    "memory": "result/memory_scalability/corrected_325557/feasibility_boundary.tsv",
    "ablation_old3": "raw_data/ablation/synthetic/job_2354994_20260710/ablation_results.tsv",
    "ablation_corrected": "raw_data/corrected_325557/job_2406254/ablation_results.tsv",
    "kernel_pa": "raw_data/tuning/kernel_selection/roadNet-PA/gpu_opt_forced/job_2354329_20260710/kernel_selection_results.tsv",
    "kernel_tx": "raw_data/tuning/kernel_selection/roadNet-TX/gpu_opt_forced/job_2354330_20260710/kernel_selection_results.tsv",
    "design": "docs/thesis/04_method_design.md",
    "main_source": "result/main_performance/proposed_variants/SOURCE.md",
}


def read_tsv(rel: str) -> list[dict[str, str]]:
    with (ROOT / rel).open(newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f, delimiter="\t"))


def median(xs):
    return float(statistics.median(list(xs)))


def sample_sd(xs):
    vals = list(xs)
    return float(statistics.stdev(vals)) if len(vals) >= 2 else None


def raw_gpu_path(graph: str) -> str:
    return f"raw_data/main_performance/proposed_variants/{graph}/_run/job_2357334_20260711/results.tsv"


def raw_sweep_path(graph: str) -> str:
    return f"raw_data/tuning/pathmerge/{graph}/pathmerge_bc/job_multi_20260710/pathmerge_sweep_results.tsv"


def raw_sweep_log(graph: str) -> str:
    return f"raw_data/tuning/pathmerge/{graph}/pathmerge_bc/job_multi_20260710/pathmerge_sweep.log"


def load_runtime() -> list[dict]:
    t2 = {r["Graph"]: r for r in read_tsv(SOURCES["t2"])}
    comp = {r["Graph"]: r for r in read_tsv(SOURCES["comparison"])}
    legacy = read_tsv("raw_data/main_performance/seven_implementations/legacy_partial/large/no_gpu_opt/job_notrecorded_legacy/results_no_gpu_opt.tsv")
    out = []
    for graph in GRAPH_ORDER:
        tuned_batch = int(comp[graph]["tuned_batch"].removeprefix("b"))
        gpu_batch = int(t2[graph]["GPU_Opt Batch"].removeprefix("b"))
        pvals = [float(r["Time_sec"]) for r in read_tsv(raw_gpu_path(graph)) if r["Implementation"] == "GPU_Opt"]
        if graph in ("email-EuAll", "roadNet-CA"):
            pmrows = [r for r in read_tsv(raw_sweep_path(graph)) if r["Config"] == f"PathMerge_b{tuned_batch}"]
            pmsource = raw_sweep_path(graph)
        else:
            pmrows = [r for r in legacy if r["Graph"] == graph and r["Implementation"] == "PathMerge_BC"]
            pmsource = "raw_data/main_performance/seven_implementations/legacy_partial/large/no_gpu_opt/job_notrecorded_legacy/results_no_gpu_opt.tsv"
        pmvals = [float(r["Time_sec"]) for r in pmrows]
        pmed, pmmed = median(pvals), median(pmvals)
        speedup = pmmed / pmed
        assert round(pmed, 2) == float(t2[graph]["GPU_Opt Median (s)"])
        assert round(pmmed, 2) == float(t2[graph]["PathMerge Median (s)"])
        assert round(speedup, 2) == float(comp[graph]["Speedup"])
        assert comp[graph]["tuned_batch"] == f"b{tuned_batch}"
        out.append({
            "Graph": graph,
            "GPU_Opt_Median_s": pmed,
            "GPU_Opt_Sample_SD_s": sample_sd(pvals),
            "GPU_Opt_N": len(pvals),
            "GPU_Opt_Batch": gpu_batch,
            "PathMerge_Median_s": pmmed,
            "PathMerge_Sample_SD_s": sample_sd(pmvals),
            "PathMerge_N": len(pmvals),
            "PathMerge_Tuned_Batch": tuned_batch,
            "Speedup": speedup,
            "GPU_Opt_Source": raw_gpu_path(graph),
            "PathMerge_Source": pmsource,
        })
    return out


def load_sweep() -> list[dict]:
    out = []
    for graph in GRAPH_ORDER:
        grouped = defaultdict(list)
        for row in read_tsv(raw_sweep_path(graph)):
            requested = int(row["Config"].replace("PathMerge_b", ""))
            grouped[requested].append(float(row["Time_sec"]))
        effective = {}
        log_path = ROOT / raw_sweep_log(graph)
        if log_path.exists():
            text = log_path.read_text(encoding="utf-8", errors="replace")
            for req, eff in re.findall(r"WARNING: batch_size=(\d+).*?clamping to (\d+)", text):
                effective[int(req)] = int(eff)
        for requested in sorted(grouped):
            vals = grouped[requested]
            out.append({
                "Graph": graph,
                "Requested_Batch": requested,
                "Effective_Batch": effective.get(requested, requested),
                "Clamped": "yes" if requested in effective else "no",
                "Median_Runtime_s": median(vals),
                "Sample_SD_s": sample_sd(vals),
                "N": len(vals),
                "CanonicalSource": raw_sweep_path(graph),
            })
    email_clamp = [r for r in out if r["Graph"] == "email-EuAll" and r["Requested_Batch"] == 8192][0]
    assert email_clamp["Effective_Batch"] == 7393
    return out


def parse_cfg(label: str):
    m = re.search(r"H([01])_?W([01])_?A([01])", label)
    if not m:
        raise ValueError(f"Unexpected ablation config: {label}")
    return tuple(int(v) for v in m.groups())


def graph_effects(rows: list[dict]) -> dict[str, float]:
    grouped = defaultdict(list)
    for r in rows:
        grouped[parse_cfg(r["Config"])].append(float(r["Time_sec"]))
    assert len(grouped) == 8 and len({len(v) for v in grouped.values()}) == 1
    meds = {k: median(v) for k, v in grouped.items()}
    out = {}
    for pos, factor in enumerate(("H", "W", "A")):
        ratios = []
        for cfg0 in sorted(k for k in meds if k[pos] == 0):
            cfg1 = list(cfg0)
            cfg1[pos] = 1
            ratios.append(meds[cfg0] / meds[tuple(cfg1)])
        out[factor] = math.exp(sum(math.log(x) for x in ratios) / len(ratios))
    return out


def load_ablation() -> list[dict]:
    old_rows = read_tsv(SOURCES["ablation_old3"])
    corrected_rows = read_tsv(SOURCES["ablation_corrected"])
    names = ["benchmark_7000_41459", "benchmark_11023_62184", "56438_300801"]
    effects = {}
    for name in names:
        effects[name] = graph_effects([r for r in old_rows if r["Graph"] == name])
    corrected_name = "325557_3216152_corrected_v1"
    effects[corrected_name] = graph_effects(corrected_rows)
    aggregate = {f: math.exp(sum(math.log(effects[g][f]) for g in names + [corrected_name]) / 4) for f in ("H", "W", "A")}
    out = []
    source_for = {g: SOURCES["ablation_old3"] for g in names}
    source_for[corrected_name] = SOURCES["ablation_corrected"]
    for graph in names + [corrected_name]:
        for factor in ("H", "W", "A"):
            out.append({"Graph": graph, "Factor": factor, "Main_Effect": effects[graph][factor], "Aggregate": "no", "CanonicalSource": source_for[graph]})
    for factor in ("H", "W", "A"):
        out.append({"Graph": "Synthetic-4 aggregate", "Factor": factor, "Main_Effect": aggregate[factor], "Aggregate": "mixed-checkpoint", "CanonicalSource": f"{SOURCES['ablation_old3']};{SOURCES['ablation_corrected']}"})
    assert [round(effects[corrected_name][f], 4) for f in ("H", "W", "A")] == [1.4767, 1.1012, 1.5563]
    assert [round(aggregate[f], 4) for f in ("H", "W", "A")] == [1.6787, 1.0661, 1.3914]
    return out


def load_memory() -> list[dict]:
    out = []
    for r in read_tsv(SOURCES["memory"]):
        observed = r["Observed"]
        if observed == "SUCCESS":
            status = "Success"
        elif r["OOMEvidenceClass"] == "cuda_oom":
            status = "CUDA OOM"
        elif r["RunnerExit"] == "137":
            status = "Cgroup host-memory OOM kill"
        else:
            raise ValueError(f"Unclassified memory result: {r}")
        out.append({
            "Config": r["Config"], "Implementation": r["Implementation"],
            "Requested_Batch": int(r["RequestedBatch"]), "Status": status,
            "Runtime_s": None if r["RuntimeSec"] == "not_recorded" else float(r["RuntimeSec"]),
            "Runner_Exit": int(r["RunnerExit"]), "CanonicalSource": SOURCES["memory"],
        })
    assert [(r["Config"], r["Status"]) for r in out] == [
        ("pure_b4096", "Success"), ("pure_b8192", "CUDA OOM"),
        ("um_b10240", "Success"), ("um_b12288", "Cgroup host-memory OOM kill"),
        ("chunked_b16384", "Success")]
    return out


def load_kernel() -> list[dict]:
    out = []
    for graph, key in (("roadNet-PA", "kernel_pa"), ("roadNet-TX", "kernel_tx")):
        rows = read_tsv(SOURCES[key])
        for kernel in ("shared", "block"):
            vals = [float(r["Time_sec"]) for r in rows if r["Kernel"] == kernel]
            out.append({"Graph": graph, "Kernel": kernel, "Median_Runtime_s": median(vals), "Sample_SD_s": sample_sd(vals), "N": len(vals), "CanonicalSource": SOURCES[key]})
    for graph, expected in (("roadNet-PA", 1.52), ("roadNet-TX", 1.66)):
        sh = next(r for r in out if r["Graph"] == graph and r["Kernel"] == "shared")
        bl = next(r for r in out if r["Graph"] == graph and r["Kernel"] == "block")
        assert round(sh["Median_Runtime_s"] / bl["Median_Runtime_s"], 2) == expected
    return out


def load_phase() -> list[dict]:
    out = []
    for graph in GRAPH_ORDER:
        rel = f"raw_data/main_performance/proposed_variants/{graph}/_run/job_2357334_20260711/phase_timing.log"
        samples, pending, active = [], None, False
        for line in (ROOT / rel).read_text(encoding="utf-8", errors="replace").splitlines():
            s = line.strip()
            if s.startswith("Running: GPU_Opt on"):
                active = True
            elif s.startswith("Running: GPU_Opt_Pure"):
                active = False
            elif active and s.startswith("> [GPU Phase]"):
                bfs = float(s.split("BFS wall=")[1].split(" s")[0])
                backward = float(s.split("Backward wall=")[1].split(" s")[0])
                pending = (bfs, backward)
            elif active and s.startswith("> Elapse time"):
                total = float(s.split("=")[1].strip())
                assert pending is not None
                bfs, backward = pending
                samples.append((bfs, backward, total - bfs - backward, total))
                active, pending = False, None
        assert samples
        for idx, phase in enumerate(("BFS", "Backward", "Other")):
            out.append({"Graph": graph, "Phase": phase, "Median_Component_s": median(v[idx] for v in samples), "N": len(samples), "Measurement": "full-duration wall-clock run", "CanonicalSource": rel})
    return out


def load_design_constants() -> dict[str, int]:
    design = (ROOT / SOURCES["design"]).read_text(encoding="utf-8")
    source = (ROOT / SOURCES["main_source"]).read_text(encoding="utf-8")
    alpha_match = re.search(r"alpha\s*=\s*(\d+)", design)
    beta_match = re.search(r"beta\s*=\s*(\d+)", design)
    ns_match = re.search(r"NS_eff\s*=\s*(\d+)", source)
    if not (alpha_match and beta_match and ns_match):
        raise ValueError("Cannot locate alpha, beta, or NS_eff in canonical design sources")
    return {"alpha": int(alpha_match.group(1)), "beta": int(beta_match.group(1)), "NS_eff": int(ns_match.group(1))}


def load_all():
    runtime = load_runtime()
    return {
        "runtime": runtime,
        "speedup": [{"Graph": r["Graph"], "Speedup": r["Speedup"], "GPU_Opt_Median_s": r["GPU_Opt_Median_s"], "PathMerge_Median_s": r["PathMerge_Median_s"], "CanonicalSource": f"{r['GPU_Opt_Source']};{r['PathMerge_Source']}"} for r in runtime],
        "sweep": load_sweep(), "ablation": load_ablation(), "memory": load_memory(),
        "kernel": load_kernel(), "phase": load_phase(), "design": load_design_constants(),
    }


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


# --- Text metrics -----------------------------------------------------------
# Helvetica / Helvetica-Bold advance widths (1/1000 em) for printable ASCII 32
# to 126, in code-point order. Arial is metrically compatible over this set, so
# these drive the layout guard that keeps every authored line on one rendered
# line inside its shape. Nothing here is measured on a font file, so the guard
# stays deterministic on a node without any font stack installed.
_ASCII = "".join(chr(c) for c in range(32, 127))
_ADV_REGULAR = [
    278, 278, 355, 556, 556, 889, 667, 191, 333, 333, 389, 584, 278, 333, 278, 278,
    556, 556, 556, 556, 556, 556, 556, 556, 556, 556, 278, 278, 584, 584, 584, 556,
    1015, 667, 667, 722, 722, 667, 611, 778, 722, 278, 500, 667, 556, 833, 722, 778,
    667, 778, 722, 667, 611, 722, 667, 944, 667, 667, 611, 278, 278, 278, 469, 556,
    333, 556, 556, 500, 556, 556, 278, 556, 556, 222, 222, 500, 222, 833, 556, 556,
    556, 556, 333, 500, 278, 556, 500, 722, 500, 500, 500, 334, 260, 334, 584,
]
_ADV_BOLD = [
    278, 333, 474, 556, 556, 889, 722, 238, 333, 333, 389, 584, 278, 333, 278, 278,
    556, 556, 556, 556, 556, 556, 556, 556, 556, 556, 333, 333, 584, 584, 584, 611,
    975, 722, 722, 722, 722, 667, 611, 778, 722, 278, 556, 722, 611, 833, 722, 778,
    667, 778, 722, 667, 611, 722, 667, 944, 667, 667, 611, 333, 278, 333, 584, 556,
    333, 556, 611, 556, 611, 556, 333, 611, 611, 278, 278, 556, 278, 889, 611, 611,
    611, 611, 389, 556, 333, 611, 556, 778, 556, 556, 500, 389, 280, 389, 584,
]
ADVANCE = {False: dict(zip(_ASCII, _ADV_REGULAR)), True: dict(zip(_ASCII, _ADV_BOLD))}

TEXT_INSET_IN = 0.12  # set_text() writes 0.06 in left + 0.06 in right margins.
LINE_HEIGHT = 1.21  # Single-spaced Arial line advance as a multiple of point size.


def text_width_in(text: str, size_pt: float, bold: bool) -> float:
    table = ADVANCE[bool(bold)]
    return sum(table[ch] for ch in text) / 1000.0 * size_pt / 72.0


def text_block_in(text: str, size_pt: float, bold: bool) -> tuple[float, float]:
    """Return the (width, height) in inches that ``text`` needs unwrapped."""
    lines = re.split(r"\r\n?|\n|\v", text)
    width = max((text_width_in(line, size_pt, bold) for line in lines), default=0.0)
    return width + TEXT_INSET_IN, len(lines) * size_pt * LINE_HEIGHT / 72.0 + 0.06


def set_text(shape, text: str, size=16, color="neutral", bold=False, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    # ``python-pptx`` preserves LF assigned directly to ``run.text`` as a raw
    # newline inside ``a:t``.  PowerPoint's interoperable soft-line-break form
    # is a sibling ``a:br`` element, so author each visible line as its own run
    # and insert the break explicitly.  All runs retain identical formatting.
    for index, line_text in enumerate(text.split("\n")):
        if index:
            p.add_line_break()
        run = p.add_run()
        run.text = line_text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(C[color])
    return shape


def text_box(slide, x, y, w, h, text, size=16, color="neutral", bold=False, align=PP_ALIGN.CENTER, name=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    return set_text(shape, text, size, color, bold, align)


def box(slide, x, y, w, h, text, fill="light", line="neutral", size=16, bold=False, radius=True, name=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(C[fill])
    shape.line.color.rgb = rgb(C[line]); shape.line.width = Pt(1.5)
    if name:
        shape.name = name
    return set_text(shape, text, size, "neutral", bold)


def circle(slide, x, y, d, text, fill="pale_blue", line="gpu", size=16, name=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(C[fill])
    s.line.color.rgb = rgb(C[line]); s.line.width = Pt(1.5)
    if name: s.name = name
    return set_text(s, text, size, "neutral", True)


def _connector(slide, x1, y1, x2, y2, color, width, dashed, ends, name):
    """Straight connector from (x1, y1) to (x2, y2).

    ``add_connector`` normalises the geometry to a top-left origin plus
    ``flipH``/``flipV``, so the local path always runs begin -> end. ``a:tailEnd``
    is therefore the arrowhead at (x2, y2) and ``a:headEnd`` the one at
    (x1, y1); the child order required by ``CT_LineProperties`` is headEnd
    before tailEnd.
    """
    s = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    s.line.color.rgb = rgb(C[color]); s.line.width = Pt(width)
    if dashed: s.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    ln = s._element.spPr.get_or_add_ln()
    for tag in ends:
        end = OxmlElement(f"a:{tag}"); end.set("type", "triangle"); end.set("w", "sm"); end.set("len", "sm")
        ln.append(end)
    if name: s.name = name
    return s


def line(slide, x1, y1, x2, y2, color="neutral", width=2.0, dashed=False, name=None):
    """Undirected connector: no arrowhead on either end."""
    return _connector(slide, x1, y1, x2, y2, color, width, dashed, (), name)


def arrow(slide, x1, y1, x2, y2, color="neutral", width=2.0, dashed=False, name=None):
    """Directed connector: the arrowhead sits on the (x2, y2) end."""
    return _connector(slide, x1, y1, x2, y2, color, width, dashed, ("tailEnd",), name)


def bi_arrow(slide, x1, y1, x2, y2, color="neutral", width=2.0, dashed=False, name=None):
    """Bidirectional connector: an arrowhead on both ends."""
    return _connector(slide, x1, y1, x2, y2, color, width, dashed, ("headEnd", "tailEnd"), name)


def label_chip(slide, x, y, w, text, fill, line, size=12):
    return box(slide, x, y, w, 0.36, text, fill, line, size, False, True)


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(C["white"])
    return slide


def style_chart(chart, value_title=None, category_title=None, legend=True):
    chart.has_title = False
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = FONT; chart.legend.font.size = Pt(14)
    va = chart.value_axis
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = rgb("D7DCE0")
    va.tick_labels.font.name = FONT; va.tick_labels.font.size = Pt(14)
    if value_title:
        va.has_title = True; va.axis_title.text_frame.text = value_title
        for p in va.axis_title.text_frame.paragraphs:
            for r in p.runs: r.font.name = FONT; r.font.size = Pt(14)
    ca = chart.category_axis
    ca.tick_labels.font.name = FONT; ca.tick_labels.font.size = Pt(14)
    if category_title:
        ca.has_title = True; ca.axis_title.text_frame.text = category_title
        for p in ca.axis_title.text_frame.paragraphs:
            for r in p.runs: r.font.name = FONT; r.font.size = Pt(14)


def add_chart(slide, chart_type, x, y, w, h, categories, series, colors, value_title, legend=True, data_labels=False, number_format="0.00"):
    data = CategoryChartData(); data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    chart = slide.shapes.add_chart(chart_type, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    style_chart(chart, value_title=value_title, legend=legend)
    for idx, ser in enumerate(chart.series):
        ser.format.fill.solid(); ser.format.fill.fore_color.rgb = rgb(C[colors[idx]])
        ser.format.line.color.rgb = rgb(C[colors[idx]])
        if data_labels:
            # ``show_value`` is what writes c:showVal; a series has no
            # ``has_data_labels`` property, so assigning one only creates an unused
            # Python attribute and leaves the default c:showVal="0" in place, which
            # is why the formatted labels never rendered.
            ser.data_labels.show_value = True
            ser.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            ser.data_labels.font.name = FONT; ser.data_labels.font.size = Pt(14)
            ser.data_labels.number_format = number_format
            ser.data_labels.number_format_is_linked = False
    return chart


CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"


def _c(tag: str) -> str:
    return f"{{{CHART_NS}}}{tag}"


def set_log_scale(axis, base=10, minimum=None, maximum=None):
    """Give ``axis`` a genuine base-``base`` logarithmic scale.

    python-pptx has no ``log_base`` property, so assigning one only creates an
    unused Python attribute and leaves the axis linear. ``c:logBase`` is the
    first child of ``c:scaling``, so inserting at index 0 keeps the schema
    sequence valid alongside the ``c:max``/``c:min`` python-pptx writes.
    """
    if minimum is not None:
        axis.minimum_scale = minimum
    if maximum is not None:
        axis.maximum_scale = maximum
    scaling = axis._element.find(_c("scaling"))
    for stale in scaling.findall(_c("logBase")):
        scaling.remove(stale)
    el = OxmlElement("c:logBase"); el.set("val", str(base))
    scaling.insert(0, el)


def _plot_area(chart):
    return chart._chartSpace.find(_c("chart")).find(_c("plotArea"))


def series_to_line(chart, index, color, width=1.75, dashed=True):
    """Move series ``index`` out of the bar plot into a sibling line plot.

    A reference level drawn as a chart series stays pinned to its value when the
    chart is resized, which a floating shape cannot do. python-pptx has no combo
    chart API, so the series element is relocated by hand; the embedded workbook
    is untouched and still holds both series.
    """
    plot_area = _plot_area(chart)
    bar = plot_area.find(_c("barChart"))
    ser = bar.findall(_c("ser"))[index]
    bar.remove(ser)

    # CT_LineSer orders children idx, order, tx, spPr, marker, ..., cat, val, smooth.
    # Data labels are dropped outright rather than merely hidden: the bar-series
    # dLbls carries dLblPos="outEnd", which is not a legal position on a line
    # series and makes PowerPoint offer to repair the file.
    for stale in ser.findall(_c("invertIfNegative")) + ser.findall(_c("dLbls")):
        ser.remove(stale)
    spPr = ser.find(_c("spPr"))
    marker = OxmlElement("c:marker")
    symbol = OxmlElement("c:symbol"); symbol.set("val", "none"); marker.append(symbol)
    spPr.addnext(marker)
    smooth = OxmlElement("c:smooth"); smooth.set("val", "0"); ser.append(smooth)

    line_chart = OxmlElement("c:lineChart")
    grouping = OxmlElement("c:grouping"); grouping.set("val", "standard"); line_chart.append(grouping)
    vary = OxmlElement("c:varyColors"); vary.set("val", "0"); line_chart.append(vary)
    line_chart.append(ser)
    no_marker = OxmlElement("c:marker"); no_marker.set("val", "0"); line_chart.append(no_marker)
    for ax_id in bar.findall(_c("axId")):
        copy = OxmlElement("c:axId"); copy.set("val", ax_id.get("val")); line_chart.append(copy)
    bar.addnext(line_chart)

    ln = OxmlElement("a:ln"); ln.set("w", str(int(width * 12700)))
    fill = OxmlElement("a:solidFill")
    clr = OxmlElement("a:srgbClr"); clr.set("val", C[color]); fill.append(clr); ln.append(fill)
    if dashed:
        dash = OxmlElement("a:prstDash"); dash.set("val", "dash"); ln.append(dash)
    for stale in spPr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}ln"):
        spPr.remove(stale)
    for stale in spPr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill"):
        spPr.remove(stale)
    spPr.append(ln)
    return ser


AXIS_ID_BASE = 200000000


def normalize_axis_ids(chart, base):
    """Rewrite this chart's axis identifiers as small positive integers.

    python-pptx derives axis ids from ``hash()``, which is free to be negative;
    ``c:axId``/``c:crossAx`` are ``xsd:unsignedInt``, so those files are invalid
    even though PowerPoint tolerates them. Remapping keeps every reference
    consistent because definitions and references are rewritten together.
    """
    chart_space = chart._chartSpace
    axis_tags = {_c("catAx"), _c("valAx"), _c("dateAx"), _c("serAx")}
    mapping: dict[str, str] = {}
    for element in chart_space.iter():
        if element.tag in axis_tags:
            old = element.find(_c("axId")).get("val")
            if old not in mapping:
                mapping[old] = str(base + len(mapping))
    for element in chart_space.iter():
        if element.tag in (_c("axId"), _c("crossAx")):
            value = element.get("val")
            if value in mapping:
                element.set("val", mapping[value])
    return mapping


def slide_01(prs):
    s = blank_slide(prs)
    box(s, 0.45, 2.55, 1.75, 1.25, "Problem", "light", "neutral", 22, True, name="Problem")
    box(s, 2.75, 2.25, 2.65, 1.85, "Proposed GPU\nFramework", "pale_blue", "gpu", 22, True, name="Framework")
    arrow(s, 2.2, 3.18, 2.75, 3.18, "gpu", 2.5, name="Problem_to_Framework")
    labels = [("Performance\nEvaluation", "pale_teal", "um"), ("Ablation\nAnalysis", "pale_purple", "pathmerge"), ("Memory\nScalability", "pale_orange", "pure"), ("Correctness\nValidation", "light", "neutral")]
    ys = [0.45, 2.15, 3.85, 5.55]
    for (txt, fill, outline), y in zip(labels, ys):
        box(s, 7.25, y, 3.05, 1.15, txt, fill, outline, 20, True, name=txt.replace("\n", "_"))
        arrow(s, 5.4, 3.18, 7.25, y + 0.58, outline, 2.0, name=f"Framework_to_{txt.split()[0]}")
    text_box(s, 10.7, 2.45, 2.1, 1.45, "Evidence for\nRQ1--RQ4", 20, "neutral", True, name="Evidence")
    for y in ys:
        arrow(s, 10.3, y + 0.58, 10.7, 3.18, "neutral", 1.2, True)


def slide_02(prs):
    s = blank_slide(prs)
    # Each label is pre-broken into lines that fit unwrapped, so PowerPoint never
    # splits a word (an earlier layout dropped the last "n" of "Accumulation").
    steps = ["Source\nInitialization", "BFS\nTraversal", "Distance and\nPath Counting", "Reverse\nDependency\nAccumulation", "BC\nAccumulation", "Next\nSource"]
    fills = ["light", "pale_blue", "pale_teal", "pale_purple", "pale_orange", "light"]
    outlines = ["neutral", "gpu", "um", "pathmerge", "pure", "neutral"]
    widths = [1.90, 1.90, 2.00, 2.05, 2.00, 1.90]
    sizes = [17, 17, 16, 16, 17, 17]
    gap, first_x, last_x = 0.24, 0.0, 0.0
    x = 0.22
    for i, (label, fill, outline, w, size) in enumerate(zip(steps, fills, outlines, widths, sizes)):
        box(s, x, 2.55, w, 1.55, label, fill, outline, size, True, name=f"Step_{i+1}")
        if i == 0:
            first_x = x + w / 2
        if i == len(steps) - 1:
            last_x = x + w / 2
        else:
            arrow(s, x + w, 3.33, x + w + gap, 3.33, "gpu", 2.2, name=f"Flow_{i+1}")
        x += w + gap
    # Feedback loop: only the final leg carries the arrowhead, into Step 1.
    line(s, last_x, 4.10, last_x, 5.25, "neutral", 1.7, True, name="Loop_Down")
    line(s, last_x, 5.25, first_x, 5.25, "neutral", 1.7, True, name="Loop_Back")
    arrow(s, first_x, 5.25, first_x, 4.10, "neutral", 1.7, True, name="Loop_To_Step_1")
    text_box(s, 4.6, 5.35, 4.1, 0.55, "Repeat for every source vertex", 16, "neutral", True)


def slide_03(prs):
    s = blank_slide(prs)
    box(s, 0.45, 0.55, 4.05, 4.85, "", "light", "cpu", 16, name="Grace_CPU_Region")
    text_box(s, 0.75, 0.75, 3.45, 0.55, "Grace CPU", 22, "neutral", True)
    box(s, 0.85, 1.55, 3.25, 2.55, "Host Memory", "white", "cpu", 20, True, name="Host_Memory")
    box(s, 8.85, 0.55, 4.05, 4.85, "", "pale_blue", "gpu", 16, name="Hopper_GPU_Region")
    text_box(s, 9.15, 0.75, 3.45, 0.55, "Hopper GPU", 22, "gpu", True)
    box(s, 9.25, 1.55, 3.25, 2.55, "HBM3", "white", "gpu", 20, True, name="HBM3")
    line(s, 4.5, 2.35, 8.85, 2.35, "gpu", 6.0, name="NVLink_C2C")
    text_box(s, 5.25, 1.55, 2.85, 0.55, "NVLink-C2C", 20, "gpu", True)
    box(s, 5.1, 2.75, 3.1, 1.35, "Unified Memory\nManaged Allocation", "pale_teal", "um", 18, True, name="Unified_Memory")
    # Managed pages migrate in both directions; a single-headed arrow would read
    # as a one-way host-to-device copy, which is not what UM placement does.
    bi_arrow(s, 5.1, 3.42, 4.1, 3.42, "um", 2.3, name="Migration_Host")
    bi_arrow(s, 8.2, 3.42, 9.25, 3.42, "um", 2.3, name="Migration_HBM")
    text_box(s, 4.65, 4.2, 4.0, 0.5, "Page placement and data migration", 14, "neutral", False)
    label_chip(s, 0.65, 5.85, 2.45, "Graph / static data", "light", "neutral", 14)
    label_chip(s, 3.25, 5.85, 2.15, "BC output", "pale_purple", "pathmerge", 14)
    label_chip(s, 5.55, 5.85, 2.45, "Source-local state", "pale_blue", "gpu", 14)
    label_chip(s, 8.15, 5.85, 3.25, "Batch-dependent working set", "pale_orange", "pure", 14)
    text_box(s, 1.2, 6.45, 11.0, 0.55, "Managed memory addresses working-set placement; the input graph file size is not the capacity claim.", 14, "neutral", True)


def slide_04(prs):
    s = blank_slide(prs)
    box(s, 0.35, 2.65, 1.55, 1.15, "Source\nBatch", "light", "neutral", 18, True)
    box(s, 2.35, 0.9, 2.0, 1.1, "Stream 0", "pale_blue", "gpu", 20, True)
    box(s, 2.35, 4.45, 2.0, 1.1, "Stream 1", "pale_blue", "gpu", 20, True)
    arrow(s, 1.9, 3.23, 2.35, 1.45, "gpu", 2.0, name="Batch_to_Stream0")
    arrow(s, 1.9, 3.23, 2.35, 5.0, "gpu", 2.0, name="Batch_to_Stream1")
    for y, idx in ((0.72, 0), (4.27, 1)):
        box(s, 4.85, y, 2.0, 1.45, "One Block\nper Source", "light", "gpu", 18, True, name=f"BlockMapping_{idx}")
        arrow(s, 4.35, y + 0.73, 4.85, y + 0.73, "gpu", 2.0, name=f"Stream{idx}_to_Block")
        box(s, 7.3, y, 1.75, 1.45, "Hybrid\nBFS", "pale_teal", "um", 18, True)
        arrow(s, 6.85, y + 0.73, 7.3, y + 0.73, "um", 2.0, name=f"Block{idx}_to_BFS")
        box(s, 9.5, y, 2.05, 1.45, "Dependency\nAccumulation", "pale_purple", "pathmerge", 17, True)
        arrow(s, 9.05, y + 0.73, 9.5, y + 0.73, "pathmerge", 2.0, name=f"BFS{idx}_to_Dependency")
    box(s, 5.1, 2.65, 3.3, 1.15, "Global BC\nAccumulation", "pale_orange", "pure", 20, True)
    # Routed off the near edge of each Dependency box; a diagonal from the far
    # right edge would be drawn straight back across the box it starts from.
    arrow(s, 10.5, 2.17, 8.4, 2.72, "pure", 2.0, name="Dependency0_to_GlobalBC")
    arrow(s, 10.5, 4.27, 8.4, 3.73, "pure", 2.0, name="Dependency1_to_GlobalBC")
    text_box(s, 2.15, 6.35, 9.05, 0.55, "Two streams share the execution framework while retaining independent source-local buffers.", 14, "neutral", True)


def slide_05(prs, design):
    s = blank_slide(prs)
    text_box(s, 0.4, 0.35, 2.2, 0.5, "Source vertices", 20, "neutral", True)
    for i in range(8):
        circle(s, 0.45 + i * 0.62, 1.0, 0.48, str(i), "light", "neutral", 14, f"Source_{i}")
    line(s, 0.45, 1.75, 5.27, 1.75, "neutral", 1.5)
    text_box(s, 0.75, 1.82, 4.25, 0.55, "Requested batch: source grouping", 16, "neutral", True)
    box(s, 5.7, 0.72, 2.0, 1.35, "Effective\nBatch", "pale_teal", "um", 20, True)
    arrow(s, 5.27, 1.25, 5.7, 1.25, "um", 2.2, name="Requested_to_Effective")
    box(s, 8.25, 0.35, 1.75, 0.85, "Stream 0", "pale_blue", "gpu", 18, True)
    box(s, 8.25, 1.55, 1.75, 0.85, "Stream 1", "pale_blue", "gpu", 18, True)
    arrow(s, 7.7, 1.38, 8.25, 0.78, "gpu", 2.0, name="Effective_to_Stream0")
    arrow(s, 7.7, 1.38, 8.25, 1.98, "gpu", 2.0, name="Effective_to_Stream1")
    text_box(s, 10.4, 0.78, 2.3, 0.95, f"NS_eff = {design['NS_eff']}", 22, "gpu", True)
    for row, y in enumerate((3.15, 5.05)):
        text_box(s, 0.45, y, 1.35, 0.6, f"Stream {row}", 18, "gpu", True)
        for j in range(3):
            box(s, 2.0 + j * 1.55, y - 0.05, 1.25, 0.75, f"Block {j}", "pale_blue", "gpu", 16, True)
            box(s, 2.0 + j * 1.55, y + 0.9, 1.25, 0.65, "State", "light", "neutral", 14, False)
            arrow(s, 2.62 + j * 1.55, y + 0.70, 2.62 + j * 1.55, y + 0.90, "neutral", 1.4)
    # One CSR object shared by both streams. Per-stream copies of this box would
    # imply the graph is replicated or partitioned per batch; it is neither.
    box(s, 7.55, 4.05, 2.45, 1.55, "Shared full graph\n(static CSR)", "light", "neutral", 18, True, name="Shared_CSR")
    arrow(s, 6.45, 3.48, 7.55, 4.45, "neutral", 1.7, True, name="Stream0_reads_CSR")
    arrow(s, 6.45, 5.38, 7.55, 5.20, "neutral", 1.7, True, name="Stream1_reads_CSR")
    text_box(s, 7.30, 5.70, 2.95, 0.6, "Both streams read\nthe same CSR", 14, "neutral", True)
    text_box(s, 10.20, 4.15, 3.00, 1.35, "Batching groups sources;\nit does not partition\nthe graph.", 16, "failure", True)


def slide_06(prs, design):
    s = blank_slide(prs)
    box(s, 0.6, 2.5, 2.45, 1.5, "Top-Down BFS", "pale_blue", "gpu", 22, True)
    box(s, 5.45, 2.5, 2.35, 1.5, "Frontier", "light", "neutral", 22, True)
    box(s, 10.35, 2.5, 2.45, 1.5, "Bottom-Up BFS", "pale_teal", "um", 22, True)
    # Upper track reads left to right, lower track right to left.
    arrow(s, 3.05, 2.85, 5.45, 2.85, "gpu", 2.5, name="TopDown_to_Frontier")
    arrow(s, 7.8, 2.85, 10.35, 2.85, "um", 2.5, name="Frontier_to_BottomUp")
    arrow(s, 10.35, 3.65, 7.8, 3.65, "um", 2.5, name="BottomUp_to_Frontier")
    arrow(s, 5.45, 3.65, 3.05, 3.65, "gpu", 2.5, name="Frontier_to_TopDown")
    text_box(s, 3.15, 1.65, 2.15, 0.7, "Switch when\nm_f > m_u / alpha", 16, "gpu", True)
    text_box(s, 8.0, 4.25, 2.15, 0.7, "Return when\n|Q| < n / beta", 16, "um", True)
    box(s, 4.25, 5.35, 4.85, 1.1, f"alpha = {design['alpha']}     beta = {design['beta']}", "light", "neutral", 20, True)
    label_chip(s, 0.75, 5.45, 2.7, "m_f: frontier edge work", "pale_blue", "gpu", 14)
    label_chip(s, 9.85, 5.45, 2.7, "m_u: unvisited edge work", "pale_teal", "um", 14)
    text_box(s, 2.65, 6.7, 8.0, 0.45, "Both states are GPU traversal modes; this is not CPU--GPU hybrid execution.", 14, "failure", True)


def slide_07(prs):
    s = blank_slide(prs)
    text_box(s, 1.55, 0.45, 9.9, 0.5, "Time", 20, "neutral", True)
    # Time runs left to right, so the axis arrowhead belongs on the right end.
    arrow(s, 1.55, 1.08, 12.55, 1.08, "neutral", 2.0, name="Time_Axis")
    for y, label in ((2.0, "Stream 0"), (4.55, "Stream 1")):
        text_box(s, 0.25, y, 1.2, 0.55, label, 18, "gpu", True)
        line(s, 1.55, y + 0.28, 12.45, y + 0.28, "neutral", 1.0)
    # Segment widths hold every label on one line at 14 pt; stream 1 is the same
    # sequence staggered by 0.70 in, which is what produces the overlap window.
    segments = [(1.45, "Initialization", "light", "neutral"), (2.60, "BFS", "pale_blue", "gpu"),
                (2.20, "Dependency", "pale_purple", "pathmerge"), (0.90, "Sync", "pale_red", "failure"),
                (1.55, "Buffer Reuse", "light", "neutral")]
    starts = {}
    for row, (y, offset) in enumerate(((1.75, 0.0), (4.30, 0.70))):
        x = 1.65 + offset
        for w, label, fill, outline in segments:
            box(s, x, y, w, 1.05, label, fill, outline, 14, True, False)
            starts[(row, label)] = (x, x + w)
            x += w + 0.10
    overlap_start = starts[(1, "BFS")][0]
    overlap_end = starts[(0, "Dependency")][1]
    box(s, overlap_start, 3.25, overlap_end - overlap_start, 0.55, "Overlap", "pale_green", "chunked", 16, True, False)
    line(s, overlap_start, 3.10, overlap_start, 3.80, "chunked", 1.5, True)
    line(s, overlap_end, 3.10, overlap_end, 3.80, "chunked", 1.5, True)
    text_box(s, 2.1, 6.35, 9.7, 0.55, "Synchronization protects buffer reuse while staggered work overlaps across streams.", 14, "neutral", True)


def slide_08(prs):
    s = blank_slide(prs)
    box(s, 1.25, 0.45, 10.85, 1.0, "Common GPU Execution Framework", "pale_blue", "gpu", 22, True)
    arrow(s, 3.0, 1.45, 3.0, 2.2, "gpu", 2.0, name="Framework_to_UM")
    arrow(s, 6.67, 1.45, 6.67, 2.2, "gpu", 2.0, name="Framework_to_Pure")
    arrow(s, 10.35, 1.45, 10.35, 2.2, "gpu", 2.0, name="Framework_to_Chunked")
    variants = [
        (0.55, "GPU_Opt", "Unified Memory", "Managed allocation\nHost/HBM placement", "pale_teal", "um"),
        (4.22, "GPU_Opt_Pure", "Device-Only Memory", "Device allocation\nCapacity-bound\nworking set", "pale_orange", "pure"),
        (7.9, "GPU_Opt_Pure_Chunked", "Source Sub-Batching", "Limits resident\nsource state\nFull graph reused", "pale_green", "chunked"),
    ]
    for x, impl, label, desc, fill, outline in variants:
        box(s, x, 2.2, 3.45, 3.25, "", fill, outline, 16, name=impl)
        text_box(s, x + 0.15, 2.45, 3.15, 0.55, impl, 18 if len(impl) < 16 else 16, "neutral", True)
        text_box(s, x + 0.2, 3.2, 3.05, 0.65, label, 19, outline, True)
        text_box(s, x + 0.2, 4.05, 3.05, 1.1, desc, 15, "neutral", False)
    text_box(s, 1.05, 6.0, 11.2, 0.8, "Memory-management variants of one framework -- not three independent proposals.\nChunked partitions the source set, not the input graph.", 18, "neutral", True)


def slide_09(prs, data):
    s = blank_slide(prs)
    cats = [GRAPH_SHORT[r["Graph"]] for r in data]
    values = [r["GPU_Opt_Median_s"] for r in data] + [r["PathMerge_Median_s"] for r in data]
    assert all(v > 0 for v in values), "a logarithmic value axis requires strictly positive data"
    chart = add_chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, 0.65, 0.35, 12.0, 5.95, cats,
                      [("GPU_Opt", [r["GPU_Opt_Median_s"] for r in data]), ("Tuned PathMerge", [r["PathMerge_Median_s"] for r in data])],
                      ["gpu", "pathmerge"], "Median runtime (s, log scale)", True, False)
    # 10 to 10000 keeps every median (30.81 s to 3079.72 s) inside three decades,
    # so email and the road graphs stay readable on the same axis.
    assert min(values) >= 10 and max(values) <= 10000
    set_log_scale(chart.value_axis, 10, minimum=10, maximum=10000)
    text_box(s, 0.85, 6.45, 11.65, 0.55, "Bars show raw-trial medians; sample SD, trial count, and tuned batch are retained in F09_Runtime.", 14, "neutral", True)


def slide_10(prs, data):
    s = blank_slide(prs)
    vals = [r["Speedup"] for r in data]
    # The parity level is a chart series, not a floating shape: a shape has to be
    # placed with a guess at the plot-area inset (the previous one landed near
    # 1.2), whereas a series stays on Y = 1.0 through any resize.
    chart = add_chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, 0.75, 0.35, 11.85, 5.95,
                      [GRAPH_SHORT[r["Graph"]] for r in data],
                      [("Speedup", vals), ("Parity 1.0x", [1.0] * len(vals))],
                      ["gpu", "neutral"], "Speedup over Tuned PathMerge (x)", True, True, '0.00"x"')
    chart.value_axis.minimum_scale = 0; chart.value_axis.maximum_scale = 3.6
    series_to_line(chart, 1, "neutral", 1.8, True)
    shown = " / ".join(f"{r['Speedup']:.2f}x" for r in data)
    text_box(s, 1.25, 6.45, 10.8, 0.55, f"Values are recomputed from unrounded raw medians: {shown}.", 14, "neutral", True)


def slide_11(prs, data):
    s = blank_slide(prs)
    batches = sorted({r["Requested_Batch"] for r in data})
    series = []
    for graph in GRAPH_ORDER:
        lookup = {r["Requested_Batch"]: r["Median_Runtime_s"] for r in data if r["Graph"] == graph}
        series.append((GRAPH_SHORT[graph], [lookup.get(b) for b in batches]))
    chart = add_chart(s, XL_CHART_TYPE.LINE_MARKERS, 0.55, 0.35, 12.2, 5.8,
                      [str(b) for b in batches], series, ["gpu", "um", "pure", "pathmerge"], "Median runtime (s)", True, False)
    chart.category_axis.has_title = True; chart.category_axis.axis_title.text_frame.text = "Requested batch"
    for p in chart.category_axis.axis_title.text_frame.paragraphs:
        for r in p.runs: r.font.name = FONT; r.font.size = Pt(14)
    clamp = next(r for r in data if r["Graph"] == "email-EuAll" and r["Clamped"] == "yes")
    text_box(s, 0.85, 6.25, 11.7, 0.75, f"Requested b{clamp['Requested_Batch']} is effective b{clamp['Effective_Batch']} for email-EuAll. Historical malformed 325557 is excluded; it is not a current graph result.", 14, "failure", True)


def slide_12(prs, data):
    s = blank_slide(prs)
    names = ["7000", "11023", "56438", "325557 corrected", "Synthetic-4"]
    lookup = {(r["Graph"], r["Factor"]): r["Main_Effect"] for r in data}
    graph_names = ["benchmark_7000_41459", "benchmark_11023_62184", "56438_300801", "325557_3216152_corrected_v1", "Synthetic-4 aggregate"]
    series = [(f, [lookup[(g, f)] for g in graph_names]) for f in ("H", "W", "A")]
    # A 1.0x reference level makes "no effect" legible, which matters most for W.
    # It is a constant series, so it neither adds data nor moves the axis range.
    series.append(("No effect (1.0x)", [1.0] * len(names)))
    chart = add_chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, 0.55, 0.35, 12.2, 5.85, names, series, ["gpu", "chunked", "pathmerge", "neutral"], "Main-effect speedup (x)", True, False)
    chart.value_axis.minimum_scale = 0.8; chart.value_axis.maximum_scale = 2.1
    series_to_line(chart, 3, "neutral", 1.5, True)
    text_box(s, 0.7, 6.35, 12.0, 0.65, "H: Hybrid BFS    W: Warp-Cooperative Accumulation    A: Dual Streams    |    Synthetic-4 is a mixed-checkpoint aggregate.", 14, "neutral", True)


def slide_13(prs, data):
    s = blank_slide(prs)
    xmin, xmax, x0, x1 = 0, max(r["Requested_Batch"] for r in data) * 1.1, 2.65, 12.25
    rows = [("GPU_Opt_Pure", "Pure", "pure", 1.45), ("GPU_Opt", "UM", "um", 3.15), ("GPU_Opt_Pure_Chunked", "Chunked", "chunked", 4.85)]
    by_impl = defaultdict(list)
    for r in data: by_impl[r["Implementation"]].append(r)
    arrow(s, x0, 6.15, x1, 6.15, "neutral", 1.8, name="Batch_Axis")
    for tick in [0] + sorted({r["Requested_Batch"] for r in data}):
        x = x0 + (tick - xmin) / (xmax - xmin) * (x1 - x0)
        line(s, x, 6.05, x, 6.25, "neutral", 1.0)
        text_box(s, x - 0.45, 6.3, 0.9, 0.35, str(tick), 12, "neutral", False)
    text_box(s, 5.25, 6.7, 4.25, 0.4, "Requested batch (tested points only)", 14, "neutral", True)
    for impl, label, col, y in rows:
        text_box(s, 0.35, y - 0.25, 2.0, 0.5, label, 20, col, True, PP_ALIGN.RIGHT)
        line(s, x0, y, x1, y, "neutral", 1.0, True)
        for r in by_impl[impl]:
            x = x0 + r["Requested_Batch"] / xmax * (x1 - x0)
            if r["Status"] == "Success":
                circle(s, x - 0.18, y - 0.18, 0.36, "", f"pale_{'orange' if col == 'pure' else 'teal' if col == 'um' else 'green'}", col, 12)
                text_box(s, x - 0.7, y - 0.80, 1.4, 0.65, f"Success\nb{r['Requested_Batch']}", 12, col, True)
            else:
                line(s, x - 0.18, y - 0.18, x + 0.18, y + 0.18, "failure", 3.0)
                line(s, x - 0.18, y + 0.18, x + 0.18, y - 0.18, "failure", 3.0)
                txt = "CUDA OOM" if r["Status"] == "CUDA OOM" else "Host/cgroup\nOOM kill"
                text_box(s, x - 0.75, y + 0.28, 1.5, 0.7, txt, 12, "failure", True)
    text_box(s, 0.65, 0.05, 12.0, 0.45, "Failures are categorical outcomes, not zero-second runtimes; each boundary point is one targeted feasibility run.", 14, "failure", True)


def slide_14(prs, data):
    s = blank_slide(prs)
    lookup = {(r["Graph"], r["Kernel"]): r for r in data}
    series = [("Shared kernel", [lookup[(g, "shared")]["Median_Runtime_s"] for g in ("roadNet-PA", "roadNet-TX")]),
              ("Block kernel", [lookup[(g, "block")]["Median_Runtime_s"] for g in ("roadNet-PA", "roadNet-TX")])]
    chart = add_chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, 0.75, 0.35, 11.85, 5.85, ["roadNet-PA", "roadNet-TX"], series, ["cpu", "gpu"], "Median runtime (s)", True, False)
    pa_speedup = lookup[("roadNet-PA", "shared")]["Median_Runtime_s"] / lookup[("roadNet-PA", "block")]["Median_Runtime_s"]
    tx_speedup = lookup[("roadNet-TX", "shared")]["Median_Runtime_s"] / lookup[("roadNet-TX", "block")]["Median_Runtime_s"]
    text_box(s, 2.35, 1.05, 2.0, 0.5, f"Block: {pa_speedup:.2f}x", 18, "gpu", True)
    text_box(s, 8.45, 1.05, 2.0, 0.5, f"Block: {tx_speedup:.2f}x", 18, "gpu", True)
    text_box(s, 0.9, 6.35, 11.5, 0.65, "Measured only on roadNet-PA and roadNet-TX. No unmeasured-graph rule or avg_deg < 5 selector is inferred.", 14, "failure", True)


def slide_15(prs, data):
    s = blank_slide(prs)
    lookup = {(r["Graph"], r["Phase"]): r["Median_Component_s"] for r in data}
    series = [(phase, [lookup[(g, phase)] for g in GRAPH_ORDER]) for phase in ("BFS", "Backward", "Other")]
    chart = add_chart(s, XL_CHART_TYPE.BAR_STACKED, 0.55, 0.35, 12.15, 5.9, [GRAPH_SHORT[g] for g in GRAPH_ORDER], series, ["gpu", "pathmerge", "cpu"], "Median component time (s)", True, False)
    text_box(s, 0.75, 6.35, 11.8, 0.65, "Complete b512 wall-clock runs: Other = total - BFS - Backward. These are not partial Nsight trace totals.", 14, "neutral", True)


CONCEPT_ROWS = {
    "F01_Overview": [("Problem", "Problem", "docs/thesis/03_chapter_outline.md", "Research-to-evidence flow"), ("Framework", "Proposed GPU Framework", "docs/thesis/04_method_design.md", "Common implementation framework"), ("Evaluation", "Performance / Ablation / Memory / Correctness", "docs/thesis/01_research_questions.md", "RQ1--RQ4 evidence")],
    "F02_Brandes": [("Step", x, "docs/thesis/writing/japanese/02_background.md", "Brandes flow") for x in ["Source Initialization", "BFS Traversal", "Distance and Path Counting", "Reverse Dependency Accumulation", "BC Accumulation", "Next Source"]],
    "F03_GH200": [("Component", x, "docs/thesis/04_method_design.md", "GH200 memory concept") for x in ["Grace CPU", "Host Memory", "NVLink-C2C", "Hopper GPU", "HBM3", "Unified Memory", "Managed Allocation", "Data Migration"]],
    "F04_Framework": [("Stage", x, "docs/thesis/04_method_design.md", "GPU execution framework") for x in ["Source Batch", "Stream 0 / Stream 1", "One Block per Source", "Hybrid BFS", "Dependency Accumulation", "Global BC Accumulation"]],
    "F05_BatchMapping": [("Concept", x, "docs/thesis/04_method_design.md;result/main_performance/proposed_variants/SOURCE.md", "Source grouping, not graph partitioning") for x in ["source vertices", "requested batch", "effective batch", "blocks", "per-source state", "two streams", "NS_eff = {NS_eff}"]],
    "F06_HybridBFS": [("State", x, "docs/thesis/04_method_design.md", "GPU traversal mode") for x in ["Top-Down BFS", "Bottom-Up BFS", "Frontier", "m_f", "m_u", "alpha = {alpha}", "beta = {beta}", "switching conditions"]],
    "F07_DualStream": [("Timeline", x, "docs/thesis/04_method_design.md", "Editable timing concept") for x in ["Stream 0", "Stream 1", "Initialization", "BFS", "Dependency Accumulation", "Buffer Reuse", "Synchronization", "Overlap"]],
    "F08_MemoryVariants": [("Variant", x, "docs/thesis/04_method_design.md", "One framework, three memory variants") for x in ["GPU_Opt -- Unified Memory", "GPU_Opt_Pure -- Device-Only Memory", "GPU_Opt_Pure_Chunked -- Source Sub-Batching"]],
}


def write_workbook(all_data):
    wb = Workbook(); wb.remove(wb.active)
    for name, rows in CONCEPT_ROWS.items():
        ws = wb.create_sheet(name); ws.append(["ElementType", "DisplayLabel", "CanonicalSource", "Notes"])
        for row in rows:
            ws.append(tuple(value.format(**all_data["design"]) if isinstance(value, str) else value for value in row))
    data_map = {
        "F09_Runtime": all_data["runtime"], "F10_Speedup": all_data["speedup"],
        "F11_PathMergeSweep": all_data["sweep"], "F12_Ablation": all_data["ablation"],
        "F13_Memory": all_data["memory"], "F14_Kernel": all_data["kernel"], "F15_Phase": all_data["phase"],
    }
    for name, rows in data_map.items():
        ws = wb.create_sheet(name)
        headers = list(rows[0].keys()); ws.append(headers)
        for row in rows: ws.append([row[h] for h in headers])
    for ws in wb.worksheets:
        ws.freeze_panes = "A2"; ws.auto_filter.ref = ws.dimensions
        for cell in ws[1]:
            cell.font = Font(name=FONT, bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor=C["neutral"])
            cell.alignment = Alignment(vertical="center")
        for col in range(1, ws.max_column + 1):
            width = max(len(str(ws.cell(row=r, column=col).value or "")) for r in range(1, min(ws.max_row, 40) + 1))
            ws.column_dimensions[get_column_letter(col)].width = min(max(width + 2, 12), 70)
        for row in ws.iter_rows(min_row=2):
            for cell in row:
                cell.font = Font(name=FONT, size=10)
                cell.alignment = Alignment(vertical="top", wrap_text=True)
    wb.save(XLSX_PATH)


FIGURES = [
    ("F01", 1, "Thesis Overview", "Connect the problem, framework, and evaluation evidence", "1", "Overview / motivation", "docs/thesis/03_chapter_outline.md", "Native shapes; connectors; text", "Editable shape diagram", "Ready", "Integrated as Figure 1.1 (Section 1.3)"),
    ("F02", 2, "Brandes Algorithm Flow", "Explain the per-source Brandes stages", "2", "Algorithm background", "docs/thesis/writing/japanese/02_background.md", "Native shapes; connectors; text", "Editable shape diagram", "Ready", "Loop is explicit"),
    ("F03", 3, "GH200 Memory Hierarchy", "Separate compute, memory placement, and working-set concepts", "2", "Architecture and memory", "docs/thesis/04_method_design.md", "Native shapes; connectors; text", "Editable shape diagram", "Ready", "Integrated as Figure 2.2 (Section 2.5); referenced again from Chapters 4 and 8"),
    ("F04", 4, "Overall GPU Execution Framework", "Show the central batched dual-stream GPU pipeline", "4", "Method overview", "docs/thesis/04_method_design.md", "Native shapes; connectors; text", "Editable shape diagram", "Ready", "Common framework"),
    ("F05", 5, "Batch-to-Source Mapping", "Clarify source batching and one-block-per-source mapping", "4", "Batch semantics", "docs/thesis/04_method_design.md;result/main_performance/proposed_variants/SOURCE.md", "Native shapes; connectors; text", "Editable shape diagram", "Ready", "NS_eff=2 in capacity"),
    ("F06", 6, "Hybrid BFS State Transition", "Explain direction switching inside GPU BFS", "4", "Hybrid BFS", "docs/thesis/04_method_design.md", "Native shapes; connectors; text", "Editable shape diagram", "Ready", "Not CPU--GPU hybrid"),
    ("F07", 7, "Dual-Stream Timeline", "Explain staggered execution and buffer reuse", "4", "Streams / overlap", "docs/thesis/04_method_design.md", "Native shapes; connectors; text", "Editable shape diagram", "Ready", "Conceptual timeline"),
    ("F08", 8, "Memory Management Variants", "Place UM, Pure, and Chunked above one execution framework", "4/8", "Memory variants", "docs/thesis/04_method_design.md", "Native shapes; connectors; text", "Editable shape diagram", "Ready", "Chunked is source sub-batching"),
    ("F09", 9, "Main Runtime Comparison", "Compare GPU_Opt with tuned PathMerge on four graphs", "6", "Headline performance", f"{SOURCES['t2']};{SOURCES['comparison']};raw_data/main_performance/proposed_variants/*", "Native PowerPoint chart; text", "Editable native chart", "Ready", "Raw medians on a base-10 log value axis (10--10000 s); sample SD retained in workbook"),
    ("F10", 10, "Speedup over Tuned PathMerge", "Show graph-specific speedup over the tuned comparator", "6", "Headline speedup", f"{SOURCES['comparison']};raw_data/main_performance/proposed_variants/*", "Native PowerPoint chart; parity line series; text", "Editable native chart", "Ready", "Raw-median recomputation; parity is a constant 1.0 line series, not a floating shape"),
    ("F11", 11, "PathMerge Batch Sweep", "Show requested-batch sensitivity for current graphs", "6 / Appendix B", "Tuning", "raw_data/tuning/pathmerge/*/pathmerge_bc/job_multi_20260710/pathmerge_sweep_results.tsv", "Native PowerPoint chart; text", "Editable native chart", "Ready", "Historical malformed 325557 excluded"),
    ("F12", 12, "Ablation Contributions", "Show H/W/A main effects by graph and mixed aggregate", "7 / Appendix C", "Ablation", f"{SOURCES['ablation_old3']};{SOURCES['ablation_corrected']};{SOURCES['t3']}", "Native PowerPoint chart; text", "Editable native chart", "Ready", "Corrected 325557; mixed-checkpoint aggregate; 1.0x no-effect reference line series"),
    ("F13", 13, "Memory Scalability", "Show tested success and distinct failure boundaries", "8", "Memory feasibility", f"{SOURCES['memory']};{SOURCES['t4']}", "Native shapes; lines; markers; text", "Editable shape chart", "Ready", "Failures are not zero runtime"),
    ("F14", 14, "Shared vs Block Kernel", "Report forced-kernel comparison on measured road graphs", "7", "Kernel selection", f"{SOURCES['kernel_pa']};{SOURCES['kernel_tx']}", "Native PowerPoint chart; text", "Editable native chart", "Ready", "No selector generalization"),
    ("F15", 15, "Phase Breakdown", "Describe complete-run BFS/backward/other components", "6/7", "Performance diagnosis", "raw_data/main_performance/proposed_variants/*/_run/job_2357334_20260711/phase_timing.log", "Native PowerPoint chart; text", "Editable native chart", "Ready", "Not partial Nsight trace totals"),
]


# --------------------------------------------------------------------------- #
# Figure ID crosswalk.
#
# The editable library IDs F01--F15 and the canonical result figure IDs F1--F7
# (result/figures/thesis/FIGURE_MANIFEST.tsv) are SEPARATE namespaces: F01 is not
# F1. Neither set is renamed; this table records how they relate.
#
#   ThesisFigureNumber      the number the reader sees in the thesis body
#   CanonicalResultFigureID the result-figure ID this slide re-draws, when one
#                           exists; `not_applicable` for the conceptual figures,
#                           which have no measured-data counterpart
#   ExportedAssets          publication assets derived from this slide
#
# Conceptual figures F01--F08 are exported by scripts/export_conceptual_figures.py
# into docs/thesis/figures/exported/. Chart slides F09--F15 are editable
# presentation copies; the artifacts published in the thesis body are the
# canonical result figures generated by scripts/generate_thesis_artifacts.py.
# Every assignment below was checked against the slide's own visible text, the
# receiving chapter section, and result/figures/thesis/FIGURE_MANIFEST.tsv.
# --------------------------------------------------------------------------- #
EXPORT_DIR = "docs/thesis/figures/exported"
RESULT_FIG_DIR = "result/figures/thesis"

CROSSWALK = {
    # FigureID: (ThesisFigureNumber, CanonicalResultFigureID, ExportedAssets)
    "F01": ("1.1", "not_applicable", f"{EXPORT_DIR}/figure_1_1_thesis_overview.{{svg,pdf,png}}"),
    "F02": ("2.1", "not_applicable", f"{EXPORT_DIR}/figure_2_1_brandes_algorithm.{{svg,pdf,png}}"),
    "F03": ("2.2", "not_applicable", f"{EXPORT_DIR}/figure_2_2_gh200_memory_hierarchy.{{svg,pdf,png}}"),
    "F04": ("4.1", "not_applicable", f"{EXPORT_DIR}/figure_4_1_gpu_execution_framework.{{svg,pdf,png}}"),
    "F05": ("4.2", "not_applicable", f"{EXPORT_DIR}/figure_4_2_batch_source_mapping.{{svg,pdf,png}}"),
    "F06": ("4.3", "not_applicable", f"{EXPORT_DIR}/figure_4_3_hybrid_bfs.{{svg,pdf,png}}"),
    "F07": ("4.4", "not_applicable", f"{EXPORT_DIR}/figure_4_4_dual_stream_timeline.{{svg,pdf,png}}"),
    "F08": ("4.5", "not_applicable", f"{EXPORT_DIR}/figure_4_5_memory_management_variants.{{svg,pdf,png}}"),
    "F09": ("6.1", "F1", f"{RESULT_FIG_DIR}/main_runtime_comparison.{{pdf,png,svg}}"),
    "F10": ("6.2", "F2", f"{RESULT_FIG_DIR}/main_speedup_over_tuned_pathmerge.{{pdf,png,svg}}"),
    "F11": ("6.3", "F3", f"{RESULT_FIG_DIR}/pathmerge_batch_sweep.{{pdf,png,svg}}"),
    "F12": ("7.1", "F4", f"{RESULT_FIG_DIR}/ablation_contributions.{{pdf,png,svg}}"),
    "F13": ("8.1", "F5", f"{RESULT_FIG_DIR}/memory_scalability_325557.{{pdf,png,svg}}"),
    "F14": ("7.2", "F6", f"{RESULT_FIG_DIR}/shared_vs_block_kernel.{{pdf,png,svg}}"),
    "F15": ("7.3", "F7", f"{RESULT_FIG_DIR}/phase_breakdown.{{pdf,png,svg}}"),
}


def write_manifest():
    header = ["LibraryFigureID", "Namespace", "SlideNumber", "ThesisFigureNumber",
              "CanonicalResultFigureID", "Title", "Purpose", "ThesisChapter",
              "PresentationUse", "CanonicalData", "EditableSource", "Generator",
              "EditableObjectType", "ExportedAssets", "ExportFormats", "Status", "Notes"]
    with MANIFEST_PATH.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f, delimiter="\t", lineterminator="\n"); w.writerow(header)
        for fid, sn, title, purpose, chapter, use, canonical, obj_desc, editable, status, notes in FIGURES:
            thesis_no, result_id, assets = CROSSWALK[fid]
            formats = "PPTX;SVG;PDF;PNG" if result_id == "not_applicable" else "PPTX"
            w.writerow([fid, "editable_library", sn, thesis_no, result_id, title,
                        purpose, chapter, use, canonical,
                        "docs/thesis/figures/editable/thesis_figure_library.pptx",
                        "scripts/generate_editable_figure_library.py", editable,
                        assets, formats, status, notes])


def write_readme():
    rows = "\n".join(
        f"| {sn} | {fid} | {title} | Figure {CROSSWALK[fid][0]} | "
        f"{CROSSWALK[fid][1]} |" for fid, sn, title, *_ in FIGURES)
    README_PATH.write_text(f"""# Editable Thesis and Presentation Figure Library

このディレクトリは、修士論文と発表資料で再利用する **1 slide = 1 figure** の編集用ライブラリである。caption と slide title は図へ焼き込まず、論文・story deck 側で管理する。既存の正式 F1--F7 は置換しない。

## Figure ID namespace（重要）

本ライブラリの ID **`F01`--`F15`** と、`result/figures/thesis/` の canonical result
figure ID **`F1`--`F7`** は **別の namespace** である。

> **`F01` は `F1` ではない。** 桁数が異なるだけの別体系であり、対応関係は
> `FIGURE_MANIFEST.tsv` の `CanonicalResultFigureID` 列だけが定義する。

- `Namespace` 列は本ライブラリの全行で `editable_library` である。
- `ThesisFigureNumber` 列が、読者が本文で見る番号（Figure 1.1 など）である。
  library ID は本文に出さない。
- conceptual 図 (`F01`--`F08`) には対応する result figure が無いため
  `CanonicalResultFigureID` は `not_applicable` とする。
- chart 図 (`F09`--`F15`) は canonical result figure の編集用複製であり、
  論文本体へ掲載するのは `scripts/generate_thesis_artifacts.py` が生成する
  canonical result figure の方である。

## Files

- `thesis_figure_library.pptx`: 16:9、白背景、英語のみの図。概念図は native shape / connector / text、データ図は native chart または editable shape chart。
- `figure_data.xlsx`: 図別worksheet。データ図はcanonical TSV/logから取得した未丸め値、conceptual図はlabel・source・notesを保持する。
- `FIGURE_MANIFEST.tsv`: 図番号、目的、章、canonical source、編集形式の対応表。
- `../../../../scripts/generate_editable_figure_library.py`: 再生成器。

## Regeneration

repository内に成果物以外の依存物を置かないGate V0実行例:

```bash
cd thesis_bc_project
PYTHONPATH=/tmp/gate_v0_editable_deps python3 scripts/generate_editable_figure_library.py
```

必要libraryは `python-pptx`, `openpyxl`, `lxml`。generatorは `raw_data/`, `result/`, `docs/thesis/` のcanonical sourceを読み、主要値の照合に失敗すると停止する。

PPTXとXLSXはzip entry timestampと `docProps/core.xml` の日時を固定値へ正規化してから書き出すため、同一interpreterでの再生成はbyte-identicalになる。

## Layout and OOXML invariants

generatorは出力後に次を自己検証し、違反があれば停止する。

- 方向付きconnectorのarrowheadは終点側 (`a:tailEnd`)。無方向は `line()`、双方向は `bi_arrow()` を使い分ける。GH200のhost/HBM間migrationは双方向で描く。
- 全labelはgenerator入力の `\n` を標準DrawingMLの `<a:br/>` へ変換し、各行がbox内幅へ収まる (Arial互換advance widthで検査)。`<a:t>` 内のliteral LFやPowerPoint側の再折返しには依存しない。
- chartの `c:axId` / `c:crossAx` は正のunsigned integerで、chart内の参照が一致し、chart間で衝突しない。
- slide外objectなし、raster画像なし、12 pt未満の表示文字なし、non-ASCII表示文字なし。

Slide 9のvalue axisは実際の底10対数軸 (`c:logBase`, 10--10000 s)。Slide 10のparity 1.0xとSlide 12の1.0x基準線は、floating shapeではなく値1.0の定数line seriesで、chartをresizeしてもY=1.0へ追従する。Slide 5のCSRは1個のみ配置し、両streamから参照する構図とする (batchはsourceを分割するのであってgraphを分割しない)。

## Editing in PowerPoint

PowerPointでPPTXを開き、図形、connector、text box、chart、legend、axis、data labelを個別に選択して移動・再着色・文字編集できる。chartは右クリックして **Edit Data** を選ぶとembedded workbookを編集できる。Slide 13の成功点・OOM記号はnative shape/lineであり個別編集する。

boxを縮めたりfont sizeを上げたりすると上記の改行不変条件は崩れるため、恒久的な変更はgeneratorへ反映して再生成し、自己検証を通すこと。

論文または発表へ図を移す場合は、対象slideで `Ctrl+A` → `Ctrl+C` を行い、宛先へ **Keep Source Formatting** で貼り付ける。captionはPPTXに含めず、論文側のcaption/label機構で付与する。

## Slide map

`Figure` 列が本文の番号、`Result ID` 列が canonical result figure との対応である
（`not_applicable` は対応する result figure が存在しない conceptual 図）。

| Slide | Library ID | Figure library title | Figure | Result ID |
|---:|---|---|---|---|
{rows}

## Canonical data policy

主要データは `result/tables/thesis/T2_main_performance.tsv`, `T3_ablation_summary.tsv`, `T4_memory_scalability.tsv` を監査表として照合し、計算には対応するraw trial TSV/logを使う。PathMerge sweepは `raw_data/tuning/pathmerge/`、kernel比較は `raw_data/tuning/kernel_selection/`、phase breakdownは `raw_data/main_performance/proposed_variants/*/_run/job_2357334_20260711/phase_timing.log` がsourceである。修正版325557だけをcurrent resultとして使う。

`figure_data.xlsx` やPowerPointのembedded chart dataを手作業で変更すると、canonical sourceからの再現性は失われる。PowerPoint上の図形編集もgenerator出力との差分になる。再生成すると手作業変更は上書きされるため、必要な変更は別名保存するかgeneratorへ反映する。

## Export

conceptual 図 (`F01`--`F08`) の組版用 asset は、本PPTXを唯一の編集正本として
`scripts/export_conceptual_figures.py` が機械的に生成する。exporterはPPTX package
を直接読み、slideのshape・connector・text runをSVGへ描き起こすため、ライブラリと
export assetが食い違うことはない。

```bash
cd thesis_bc_project
PYTHONPATH=/tmp/gate_v0_editable_deps39 python3 scripts/export_conceptual_figures.py
```

出力は `../exported/` に `figure_<章>_<番号>_<slug>.{{svg,pdf,png}}` として置かれ、
SVG/PDFはvector（PDFはfont subset埋込み）、PNGは300 dpiのreview用である。詳細と
tool版数は `../exported/README.md` を参照。

chart 図 (`F09`--`F15`) はexport対象ではない。論文本体へ掲載するのは
`scripts/generate_thesis_artifacts.py` が生成する canonical result figure
(`result/figures/thesis/`) である。

PowerPointから手動で書き出す場合は、対象slideの全objectを選択して右クリック
**Save as Picture** (SVG/PNG)、PDFは **File > Export > Create PDF/XPS** を使う。
LibreOfficeでは `soffice --headless --convert-to pdf thesis_figure_library.pptx`
が使えるが、font/connector/chartの差を目視確認すること。手動書き出しは正式assetの
生成手段ではなく、確認用途に限る。
""", encoding="utf-8")


def build_pptx(all_data):
    prs = Presentation(); prs.slide_width = Inches(SLIDE_W); prs.slide_height = Inches(SLIDE_H)
    slide_01(prs); slide_02(prs); slide_03(prs); slide_04(prs); slide_05(prs, all_data["design"]); slide_06(prs, all_data["design"]); slide_07(prs); slide_08(prs)
    slide_09(prs, all_data["runtime"]); slide_10(prs, all_data["speedup"]); slide_11(prs, all_data["sweep"])
    slide_12(prs, all_data["ablation"]); slide_13(prs, all_data["memory"]); slide_14(prs, all_data["kernel"]); slide_15(prs, all_data["phase"])
    assert len(prs.slides) == 15
    axis_ids = {}
    index = 0
    for slide_number, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if getattr(shape, "has_chart", False):
                axis_ids[slide_number] = normalize_axis_ids(shape.chart, AXIS_ID_BASE + index * 10)
                index += 1
    prs.save(PPTX_PATH)
    return axis_ids


FIXED_ZIP_TIME = (1980, 1, 1, 0, 0, 0)
FIXED_TIMESTAMP = "2026-07-18T00:00:00Z"
_CORE_DATE_RE = re.compile(
    r"(<dcterms:(?:created|modified)\b[^>]*>)[^<]*(</dcterms:(?:created|modified)>)")


def normalize_package(blob: bytes) -> bytes:
    """Rebuild an OOXML package with every non-deterministic field pinned.

    Zip member timestamps come from the clock at write time, and both openpyxl
    and the xlsxwriter-backed chart workbooks stamp ``docProps/core.xml`` with
    "now". Neither carries meaning here, so both are frozen to keep regeneration
    byte-identical. Entry order, compression, and payloads are untouched, and
    embedded workbooks are normalized recursively.
    """
    source = zipfile.ZipFile(io.BytesIO(blob))
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as out:
        for item in source.infolist():
            payload = source.read(item.filename)
            if item.filename == "docProps/core.xml":
                payload = _CORE_DATE_RE.sub(rf"\g<1>{FIXED_TIMESTAMP}\g<2>", payload.decode("utf-8")).encode("utf-8")
            elif item.filename.lower().endswith(".xlsx"):
                payload = normalize_package(payload)
            info = zipfile.ZipInfo(item.filename, date_time=FIXED_ZIP_TIME)
            info.compress_type = item.compress_type
            info.external_attr = item.external_attr
            info.create_system = 0
            out.writestr(info, payload)
    return buffer.getvalue()


def normalize_in_place(path: Path) -> None:
    path.write_bytes(normalize_package(path.read_bytes()))


def check_text_fits(prs) -> list[tuple]:
    """Report any authored line that would not render inside its own shape.

    Every label is written pre-broken on ``\\n``; if a line is wider than the
    usable box width PowerPoint rewraps it mid-word, which is how "Initialization"
    and "BC Accumulation" broke in the first pass. Shapes are measured from the
    saved file so the check sees what PowerPoint will.
    """
    problems = []
    for slide_number, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
                continue
            runs = [r for p in shape.text_frame.paragraphs for r in p.runs if r.text]
            if not runs:
                continue
            size = max(r.font.size.pt for r in runs if r.font.size)
            bold = any(r.font.bold for r in runs)
            need_w, need_h = text_block_in(shape.text, size, bold)
            have_w, have_h = shape.width / 914400, shape.height / 914400
            if need_w > have_w + 1e-6:
                problems.append((slide_number, shape.name, "width", round(need_w, 3), round(have_w, 3), shape.text))
            if need_h > have_h + 1e-6:
                problems.append((slide_number, shape.name, "height", round(need_h, 3), round(have_h, 3), shape.text))
    return problems


def check_arrowheads() -> dict[str, int]:
    """Confirm no connector carries a lone ``a:headEnd``.

    A head-only arrowhead points back at the connector's begin point, which is
    the inverted-direction defect this generator previously shipped on every
    directed connector. Bidirectional connectors carry both ends and are fine.
    """
    counts = {"undirected": 0, "directed": 0, "bidirectional": 0}
    inverted = []
    with zipfile.ZipFile(PPTX_PATH) as zf:
        for name in sorted(n for n in zf.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)):
            text = zf.read(name).decode("utf-8")
            for cxn in re.findall(r"<p:cxnSp>.*?</p:cxnSp>", text, re.S):
                head, tail = "a:headEnd" in cxn, "a:tailEnd" in cxn
                if head and tail:
                    counts["bidirectional"] += 1
                elif tail:
                    counts["directed"] += 1
                elif head:
                    inverted.append((name, re.search(r'name="([^"]*)"', cxn).group(1)))
                else:
                    counts["undirected"] += 1
    assert not inverted, f"connectors with a reversed arrowhead: {inverted}"
    return counts


def check_drawingml_line_breaks() -> dict[str, int]:
    """Require standard DrawingML breaks instead of literal LF in ``a:t``."""
    namespace = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    literal_lf = []
    breaks = 0
    paragraph_boundaries = 0
    with zipfile.ZipFile(PPTX_PATH) as zf:
        slide_names = sorted(
            (n for n in zf.namelist()
             if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)),
            key=lambda n: int(re.search(r"slide(\d+)", n).group(1)),
        )
        for name in slide_names:
            root = ET.fromstring(zf.read(name))
            for text in root.findall(".//a:t", namespace):
                if text.text and "\n" in text.text:
                    literal_lf.append((name, text.text))
            breaks += len(root.findall(".//a:br", namespace))
            for body in root.findall(".//a:p/..", namespace):
                paragraphs = body.findall("a:p", namespace)
                paragraph_boundaries += max(0, len(paragraphs) - 1)
    assert not literal_lf, f"literal LF found inside a:t: {literal_lf}"
    return {
        "literal_lf_in_a_t": len(literal_lf),
        "drawingml_breaks": breaks,
        "paragraph_boundaries": paragraph_boundaries,
    }


def check_axis_ids() -> tuple[list, dict]:
    """Confirm no chart part carries a negative or dangling axis reference."""
    negatives, per_chart = [], {}
    with zipfile.ZipFile(PPTX_PATH) as zf:
        for name in sorted(n for n in zf.namelist() if re.fullmatch(r"ppt/charts/chart\d+\.xml", n)):
            text = zf.read(name).decode("utf-8")
            ids = [int(v) for v in re.findall(r'<c:axId val="(-?\d+)"/>', text)]
            cross = [int(v) for v in re.findall(r'<c:crossAx val="(-?\d+)"/>', text)]
            negatives += [(name, v) for v in ids + cross if v < 0]
            defined = set(ids)
            assert set(cross) <= defined, f"{name}: crossAx points at an undefined axis"
            assert len(defined) == 2, f"{name}: expected one category and one value axis, got {sorted(defined)}"
            per_chart[name] = sorted(defined)
    all_ids = [v for values in per_chart.values() for v in values]
    assert len(all_ids) == len(set(all_ids)), "axis ids collide across charts"
    return negatives, per_chart


def validate_outputs() -> dict:
    prs = Presentation(PPTX_PATH)
    assert len(prs.slides) == 15
    slide_stats, english_errors, bounds_errors, font_errors = [], [], [], []
    total_shapes = total_text = total_charts = total_pictures = 0
    for sn, slide in enumerate(prs.slides, 1):
        texts, charts, pictures = 0, 0, 0
        for shape in slide.shapes:
            total_shapes += 1
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
                bounds_errors.append((sn, shape.name))
            if getattr(shape, "has_chart", False): charts += 1
            if shape.shape_type == 13: pictures += 1
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                texts += 1
                non_ascii = [ch for ch in shape.text if ord(ch) > 127]
                if non_ascii: english_errors.append((sn, shape.name, "".join(sorted(set(non_ascii)))))
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text and run.font.size and run.font.size.pt < 12:
                            font_errors.append((sn, shape.name, run.font.size.pt))
        assert texts > 0, f"slide {sn}: no native text"
        assert len(slide.shapes) - pictures > 1, f"slide {sn}: not a multi-object editable figure"
        slide_stats.append({"Slide": sn, "Objects": len(slide.shapes), "TextObjects": texts, "Charts": charts, "Pictures": pictures})
        total_text += texts; total_charts += charts; total_pictures += pictures
    assert not english_errors, english_errors
    assert not bounds_errors, bounds_errors
    assert not font_errors, font_errors
    overflow_errors = check_text_fits(prs)
    assert not overflow_errors, overflow_errors
    negative_axis_ids, axis_ids = check_axis_ids()
    assert not negative_axis_ids, negative_axis_ids
    connectors = check_arrowheads()
    line_breaks = check_drawingml_line_breaks()
    assert total_pictures == 0
    assert total_charts == 6
    with zipfile.ZipFile(PPTX_PATH) as zf:
        names = zf.namelist()
        chart_xml = [n for n in names if re.fullmatch(r"ppt/charts/chart\d+\.xml", n)]
        embeddings = [n for n in names if n.startswith("ppt/embeddings/")]
        raster = [n for n in names if n.startswith("ppt/media/") and n.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff"))]
        largest_name, largest_size = max(((n, zf.getinfo(n).file_size) for n in names), key=lambda x: x[1])
        largest_embedded_name, largest_embedded_size = max(((n, zf.getinfo(n).file_size) for n in embeddings), key=lambda x: x[1])
        assert len(chart_xml) == 6 and len(embeddings) >= 6
        assert not raster
    wb = __import__("openpyxl").load_workbook(XLSX_PATH, read_only=True, data_only=False)
    expected_sheets = [f"F{i:02d}_{name}" for i, name in enumerate(["Overview", "Brandes", "GH200", "Framework", "BatchMapping", "HybridBFS", "DualStream", "MemoryVariants", "Runtime", "Speedup", "PathMergeSweep", "Ablation", "Memory", "Kernel", "Phase"], 1)]
    assert wb.sheetnames == expected_sheets, wb.sheetnames
    with MANIFEST_PATH.open(encoding="utf-8") as f:
        manifest_count = sum(1 for _ in csv.DictReader(f, delimiter="\t"))
    assert manifest_count == 15
    return {"slides": len(prs.slides), "objects": total_shapes, "text_objects": total_text, "charts": total_charts, "pictures": total_pictures, "raster_only_slides": 0, "shape_diagrams": 9, "worksheets": wb.sheetnames, "manifest": manifest_count, "largest_asset": largest_name, "largest_asset_bytes": largest_size, "largest_embedded_asset": largest_embedded_name, "largest_embedded_asset_bytes": largest_embedded_size, "slide_stats": slide_stats, "text_overflow": len(overflow_errors), "negative_axis_ids": len(negative_axis_ids), "axis_ids": axis_ids, "connectors": connectors, "line_breaks": line_breaks}


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    data = load_all()
    build_pptx(data); write_workbook(data); write_manifest(); write_readme()
    normalize_in_place(PPTX_PATH); normalize_in_place(XLSX_PATH)
    stats = validate_outputs()
    print("Gate V0.1 editable figure library generated and validated")
    for key in ("slides", "objects", "text_objects", "charts", "pictures", "raster_only_slides", "shape_diagrams", "manifest", "text_overflow", "negative_axis_ids", "largest_asset", "largest_asset_bytes", "largest_embedded_asset", "largest_embedded_asset_bytes"):
        print(f"{key}\t{stats[key]}")
    print("connectors\t" + " ".join(f"{k}={v}" for k, v in stats["connectors"].items()))
    print("line_breaks\t" + " ".join(f"{k}={v}" for k, v in stats["line_breaks"].items()))
    print("worksheets\t" + ",".join(stats["worksheets"]))
    for name, ids in stats["axis_ids"].items():
        print(f"axis_ids\t{name}\t{ids[0]},{ids[1]}")
    for row in stats["slide_stats"]:
        print("slide\t" + "\t".join(f"{k}={v}" for k, v in row.items()))


if __name__ == "__main__":
    try:
        main()
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        raise

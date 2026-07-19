#!/usr/bin/env python3
"""Generate the Gate V1.2.1 results-first master's-thesis presentation deck.

Every character the audience can see is English: slide titles, body text,
callouts, captions, table cells, chart categories, series names, axis titles, and
diagram nodes. Japanese survives only in the speaker notes, which carry a
complete English script and a complete Japanese script per slide so the talk can
be delivered in either language without reading both scripts back to back.

Canonical values are recomputed from ``raw_data``/``result`` through the Gate V0
figure-library loaders, so this deck and the figure library cannot drift apart:
the shared helpers, palette, font, and package normalization are imported rather
than copied.

Run from ``thesis_bc_project``::

    PYTHONPATH=<deps> python3 scripts/generate_thesis_presentation.py
"""

from __future__ import annotations

import csv
import io
import re
import sys
import zipfile
from pathlib import Path
from typing import NamedTuple

SCRIPTS = Path(__file__).resolve().parent
if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import generate_editable_figure_library as lib
from generate_editable_figure_library import (  # shared design system
    ADVANCE,
    AXIS_ID_BASE,
    C,
    FONT,
    GRAPH_ORDER,
    GRAPH_SHORT,
    TEXT_INSET_IN,
    LINE_HEIGHT,
    add_chart,
    arrow,
    blank_slide,
    box,
    circle,
    label_chip,
    line,
    load_all,
    normalize_axis_ids,
    normalize_in_place,
    rgb,
    series_to_line,
    set_log_scale,
    text_box,
)

ROOT = SCRIPTS.parent
OUT = ROOT / "docs" / "thesis" / "presentation"
EDITABLE = OUT / "editable"
PPTX_PATH = EDITABLE / "master_thesis_presentation_v1.pptx"
MANIFEST_PATH = OUT / "PRESENTATION_MANIFEST.tsv"
README_PATH = OUT / "README.md"
PLAN_PATH = OUT / "presentation_plan.md"
NOTES_PATH = OUT / "speaker_notes_bilingual.md"

DECK_TITLE_PT = 32   # deck title floor is 30 pt
TITLE_PT = 28        # slide title floor
BODY_PT = 20         # body floor
FIGURE_MIN_PT = 16   # chart / table / diagram floor
FOOTNOTE_PT = 14     # used sparingly; nothing is smaller than this

THESIS_TITLE_EN = (
    "Design and Evaluation of a Batch-Based\n"
    "GPU Execution Framework for\n"
    "Betweenness Centrality on GH200"
)
THESIS_TITLE_FLAT = THESIS_TITLE_EN.replace("\n", " ")
RESULT_HEADLINE = (
    "Main Result\n"
    "1.31–3.17× faster than the tuned external comparator"
)
TBF = "[TO BE FILLED]"

# Visible slides are English, but a few typographic characters outside ASCII are
# deliberate: an en dash for numeric ranges, an em dash for the historical
# banner, and a multiplication sign for speedups. Everything else visible must be
# ASCII, and no CJK may appear on a slide at all.
VISIBLE_NON_ASCII = {"–": 556, "—": 1000, "×": 584, "·": 278, "≈": 889}
CJK_RE = re.compile(r"[\u3040-\u30ff\u3400-\u9fff]")

# Claim-boundary sentences. These are the wordings the deck is required to state
# verbatim, so they are authored once and asserted back out of the saved package
# rather than being retyped into each slide body.
WORKING_SET_SENTENCE = "Input file size is not working-set size."
NOT_PARTITION_SENTENCE = "Source batching groups source vertices. It does not split the graph."
COMPARATOR_SENTENCE = "PathMerge is a third-party external comparator. It is not ground truth."
NO_GENERALIZE_SENTENCE = "I do not generalize this result to PathMerge in general."
FAILURE_CATEGORICAL_SENTENCE = "Failures are not zero-second runtimes."
CHUNKED_BOUND_SENTENCE = "A tested upper bound is not an unlimited capacity claim."
TIER_A_SENTENCE = "Tier A uses an independent Sequential CPU reference."
TIER_B_SENTENCE = ("Tier B compares implementation paths and is not an independent "
                   "ground-truth evaluation.")
TOLERANCE_SENTENCE = "All 13 comparisons passed the mixed tolerance, but none was byte-identical."
SCOPE_SENTENCE = "I do not generalize beyond these conditions."
CORRECTNESS_ROLE_SENTENCE = ("Correctness is treated as required validation. "
                             "It is not presented as the main research result.")
HISTORICAL_BANNER = "HISTORICAL EVIDENCE \u2014 NOT PART OF THE CURRENT RESULTS"

# Slide number -> sentences that must appear verbatim on that slide, ignoring the
# line breaks introduced purely to fit a shape.
REQUIRED_SENTENCES = {
    7: [WORKING_SET_SENTENCE, NOT_PARTITION_SENTENCE],
    11: [COMPARATOR_SENTENCE, NO_GENERALIZE_SENTENCE],
    13: [FAILURE_CATEGORICAL_SENTENCE, CHUNKED_BOUND_SENTENCE],
    14: [SCOPE_SENTENCE],
    22: [TIER_B_SENTENCE, CORRECTNESS_ROLE_SENTENCE],
    23: [HISTORICAL_BANNER, CORRECTNESS_ROLE_SENTENCE],
}


# --- Text metrics -----------------------------------------------------------
def _is_fullwidth(ch: str) -> bool:
    o = ord(ch)
    return (
        0x1100 <= o <= 0x115F
        or 0x2E80 <= o <= 0xA4CF
        or 0xAC00 <= o <= 0xD7A3
        or 0xF900 <= o <= 0xFAFF
        or 0xFE30 <= o <= 0xFE6F
        or 0xFF00 <= o <= 0xFF60
        or 0xFFE0 <= o <= 0xFFE6
    )


def text_width_in(text: str, size_pt: float, bold: bool) -> float:
    """Width of one unwrapped line in the deck's Arial metrics.

    Visible text is English, so the Arial-compatible advance table the figure
    library already uses covers nearly everything; the few deliberate non-ASCII
    typographic characters carry their own Arial advances. Nothing is measured
    from an installed font, so the guard stays deterministic on a node with no
    font stack.
    """
    table = ADVANCE[bool(bold)]
    total = 0
    for ch in text:
        if ch in table:
            total += table[ch]
        elif ch in VISIBLE_NON_ASCII:
            total += VISIBLE_NON_ASCII[ch]
        elif _is_fullwidth(ch):
            total += 1000
        else:
            total += 500
    return total / 1000.0 * size_pt / 72.0


def text_block_in(text: str, size_pt: float, bold: bool) -> tuple[float, float]:
    lines = text.split("\n")
    width = max((text_width_in(l, size_pt, bold) for l in lines), default=0.0)
    return width + TEXT_INSET_IN, len(lines) * size_pt * LINE_HEIGHT / 72.0 + 0.06


def tx(slide, x, y, w, h, text, size=BODY_PT, color="neutral", bold=False,
       align=PP_ALIGN.LEFT, name="Body", valign=MSO_ANCHOR.TOP):
    """Left-aligned, top-anchored narrative text box in the deck's Latin font.

    ``lib.text_box`` centres both axes, which suits figure labels but not slide
    prose, so the narrative path keeps its own anchoring and naming while sharing
    the same font and palette.
    """
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    lib.set_text(shape, text, size, color, bold, align, valign)
    return shape


# A slide title is authored as its canonical one-line string; the three that do
# not fit the band at the 28 pt floor declare their own break here. Shrinking the
# title instead would breach the floor, and relying on PowerPoint's auto-wrap
# would put the break outside the text-fit guard's view.
TITLE_BREAKS = {
    "Brandes' Algorithm Reduces the Cost of Exact BC":
        "Brandes' Algorithm Reduces the Cost\nof Exact BC",
    "The Proposal Is an Integrated GPU Execution Framework":
        "The Proposal Is an Integrated GPU\nExecution Framework",
    "Source Batching Creates a Batch-Dependent Working Set":
        "Source Batching Creates a Batch-Dependent\nWorking Set",
    "Three Memory-Management Variants Share One Framework":
        "Three Memory-Management Variants Share\nOne Framework",
    "Hybrid BFS and Dual Streams Gave the Largest Observed Effects":
        "Hybrid BFS and Dual Streams Gave\nthe Largest Observed Effects",
    "The Framework Improved Performance and Expanded the Tested Batch Range":
        "The Framework Improved Performance\nand Expanded the Tested Batch Range",
}


def slide_frame(prs, title: str, accent: str = "gpu"):
    """Blank white slide carrying the standard English title band."""
    display = TITLE_BREAKS.get(title, title)
    assert display.replace("\n", " ") == title, f"title break altered wording: {title}"
    lines = display.count("\n") + 1
    assert lines <= 2, f"title needs more than two lines: {title}"
    # Two-line titles start higher and push the accent rule down; every slide
    # body begins at 1.55 in or lower, so the rule still clears the content.
    top, height, rule_y = (0.30, 0.78, 1.16) if lines == 1 else (0.22, 1.02, 1.30)
    s = blank_slide(prs)
    tx(s, 0.55, top, 12.25, height, display,
       TITLE_PT, "neutral", True, PP_ALIGN.LEFT, "Title")
    rule = s.shapes.add_shape(lib.MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(rule_y), Inches(12.25), Inches(0.045))
    rule.fill.solid()
    rule.fill.fore_color.rgb = rgb(C[accent])
    rule.line.fill.background()
    rule.name = "TitleRule"
    return s


def note(slide, text, color="neutral", y=6.72, size=FIGURE_MIN_PT, name="Note"):
    return tx(slide, 0.55, y, 12.25, 0.55, text, size, color, True, PP_ALIGN.LEFT, name)


def bullets(slide, items, x=0.9, y=1.75, w=11.6, gap=0.92, size=BODY_PT,
            accent="gpu", prefix="Bullet"):
    for i, text in enumerate(items):
        top = y + i * gap
        chip = slide.shapes.add_shape(lib.MSO_SHAPE.RECTANGLE, Inches(x), Inches(top + 0.10), Inches(0.09), Inches(0.42))
        chip.fill.solid()
        chip.fill.fore_color.rgb = rgb(C[accent])
        chip.line.fill.background()
        chip.name = f"{prefix}Mark_{i+1}"
        tx(slide, x + 0.32, top, w, min(0.80, gap), text, size, "neutral", False,
           PP_ALIGN.LEFT, f"{prefix}_{i+1}")


def table(slide, x, y, w, h, rows, col_widths=None, size=FIGURE_MIN_PT, name="Table",
          header_fill="neutral", first_col_bold=True):
    """Native, individually editable PowerPoint table. Interior is English only."""
    n_rows, n_cols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    tbl = shape.table
    tbl.first_row = True
    tbl.horz_banding = True
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)
    for r, row in enumerate(rows):
        tbl.rows[r].height = Inches(h / n_rows)
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = value
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            for run in para.runs:
                run.font.name = FONT
                run.font.size = Pt(size)
                run.font.bold = r == 0 or (first_col_bold and c == 0)
                run.font.color.rgb = rgb(C["white"] if r == 0 else C["neutral"])
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(C[header_fill])
    return shape


def bump_chart_fonts(chart, size=FIGURE_MIN_PT):
    """Raise every chart-interior font to the presentation floor.

    ``style_chart`` targets the thesis figure library, where 14 pt is legible at
    page scale. A projected slide needs the same floor as the rest of the deck,
    so legend, ticks, axis titles, and data labels are re-set here instead of
    forking the shared styler.
    """
    if chart.has_legend:
        chart.legend.font.size = Pt(size)
        chart.legend.font.name = FONT
    for axis in (chart.value_axis, chart.category_axis):
        axis.tick_labels.font.size = Pt(size)
        axis.tick_labels.font.name = FONT
        if axis.has_title:
            for p in axis.axis_title.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(size)
                    r.font.name = FONT
    for ser in chart.series:
        try:
            ser.data_labels.font.size = Pt(size)
            ser.data_labels.font.name = FONT
        except (AttributeError, ValueError):
            pass
    return chart


# --- Main deck --------------------------------------------------------------
def slide_title(prs):
    s = blank_slide(prs)
    band = s.shapes.add_shape(lib.MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(13.333), Inches(0.34))
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(C["gpu"])
    band.line.fill.background()
    band.name = "TitleBand"
    tx(s, 0.65, 0.58, 12.05, 1.88, THESIS_TITLE_EN, DECK_TITLE_PT,
       "neutral", True, PP_ALIGN.CENTER, "ThesisTitle", MSO_ANCHOR.MIDDLE)
    tx(s, 0.85, 2.48, 11.6, 0.38, "Master's Thesis Presentation",
       18, "gpu", True, PP_ALIGN.CENTER, "TitleSubtitle")
    box(s, 1.20, 2.98, 10.93, 1.20, RESULT_HEADLINE,
        "pale_blue", "gpu", 22, True, name="MainResult")
    tx(s, 1.20, 4.22, 10.93, 0.45,
       "One GH200 GPU · Four evaluated graphs · Median-to-median comparison",
       FIGURE_MIN_PT, "neutral", True, PP_ALIGN.CENTER, "ResultScope")
    rule = s.shapes.add_shape(lib.MSO_SHAPE.RECTANGLE, Inches(0.88), Inches(4.82), Inches(3.1), Inches(0.04))
    rule.fill.solid()
    rule.fill.fore_color.rgb = rgb(C["gpu"])
    rule.line.fill.background()
    rule.name = "TitleAccent"
    tx(s, 0.85, 5.02, 6.4, 1.72,
       f"Name: {TBF}\nAffiliation: {TBF}\nSupervisor: {TBF}\nDate: {TBF}",
       16, "neutral", False, PP_ALIGN.LEFT, "TitleMeta")


def slide_agenda(prs):
    s = slide_frame(prs, "Agenda")
    items = [
        "1. What Is Betweenness Centrality?",
        "2. How Does Brandes' Algorithm Work?",
        "3. How Do I Run Brandes on the GPU?",
        "4. How Do I Evaluate the Framework?",
        "5. What Results Do I Get?",
        "6. What Are the Limits?",
    ]
    for i, item in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.75 + col * 6.15
        y = 1.62 + row * 1.30
        box(s, x, y, 5.65, 0.94, item, "light" if col == 0 else "pale_blue",
            "gpu", 19, True, name=f"Agenda_{i + 1}")
    box(s, 1.35, 5.72, 10.65, 1.02,
        "I first explain BC and Brandes.\nI then show the GPU design, results, and limits.",
        "pale_teal", "um", 19, True, name="AgendaSummary")


def slide_bc_graph(prs):
    s = slide_frame(prs, "Betweenness Centrality Finds Important Bridge Vertices")
    centers = [1.25, 3.35, 5.55, 7.75, 9.85]
    y, d = 2.10, 1.05
    for i in range(4):
        line(s, centers[i] + d, y + d / 2, centers[i + 1], y + d / 2,
             "neutral", 3.0, name=f"BC_Edge_{i + 1}")
    labels = ["A", "B", "C", "D", "E"]
    values = [0, 3, 4, 3, 0]
    for i, (label, value) in enumerate(zip(labels, values)):
        fill = "pale_orange" if label == "C" else "pale_blue"
        outline = "pure" if label == "C" else "gpu"
        circle(s, centers[i], y, d, label, fill, outline, 24, f"BC_Node_{label}")
        text_box(s, centers[i] - 0.28, 3.30, 1.61, 0.42,
                 f"{label}: BC = {value}", FIGURE_MIN_PT, outline,
                 label == "C", name=f"BC_Value_{label}")
    box(s, 0.75, 4.05, 5.85, 1.10,
        "Vertex C has the highest BC.\nC lies on the most shortest paths.",
        "pale_orange", "pure", 20, True, name="BCMessage")
    box(s, 6.90, 4.05, 5.70, 1.25,
        "BC helps find important transfer points.\nExamples include communication,\ntransport, and social networks.",
        "pale_teal", "um", 18, True, name="BCUse")
    note(s, "The graph is a simple path. C is the central bridge between both sides.", "gpu", 6.10)


def slide_brandes(prs):
    s = slide_frame(prs, "Brandes' Algorithm Reduces the Cost of Exact BC")
    tx(s, 0.70, 1.46, 12.0, 0.78,
       "Direct path enumeration is expensive.\nBrandes reuses shortest-path information.",
       BODY_PT, "neutral", True, PP_ALIGN.CENTER, "BrandesIntro")
    steps = ["Choose One\nSource", "Run BFS", "Count Shortest\nPaths",
             "Accumulate\nDependencies", "Update BC"]
    widths = [2.00, 1.55, 2.25, 2.35, 1.75]
    x = 0.60
    for i, (label, width) in enumerate(zip(steps, widths)):
        box(s, x, 2.55, width, 1.02, label, "pale_blue", "gpu",
            FIGURE_MIN_PT, True, name=f"BrandesStep_{i + 1}")
        if i < len(steps) - 1:
            arrow(s, x + width, 3.06, x + width + 0.28, 3.06,
                  "gpu", 2.0, name=f"BrandesFlow_{i + 1}")
        x += width + 0.28
    box(s, 0.75, 4.18, 4.30, 1.18,
        "Unweighted exact BC:\nO(|V|(|V| + |E|))",
        "light", "neutral", 22, True, name="BrandesComplexity")
    box(s, 5.45, 4.18, 7.10, 1.18,
        "Brandes improves the algorithmic cost.\nI do not change Brandes' mathematical algorithm.",
        "pale_orange", "pure", 19, True, name="BrandesBoundary")
    note(s, "This study improves GPU execution. It does not improve Brandes' asymptotic complexity.", "failure", 6.12)


def slide_gpu_contribution(prs):
    s = slide_frame(prs, "I Improve How Brandes Runs on the GPU")
    box(s, 0.65, 1.55, 5.75, 0.72, "Brandes Defines the Exact Computation",
        "pale_blue", "gpu", 20, True, name="BrandesSideTitle")
    bullets(s, ["BFS from every source", "Reverse dependency accumulation", "Exact BC output"],
            x=0.85, y=2.48, w=5.20, gap=0.78, size=18, prefix="BrandesBullet")
    box(s, 6.88, 1.55, 5.80, 0.72, "I Improve the GPU Execution",
        "pale_orange", "pure", 20, True, name="GPUSideTitle")
    bullets(s, ["Batch many sources", "Assign one block to one source", "Switch BFS direction",
                "Use warp cooperation", "Use two CUDA streams"],
            x=7.08, y=2.35, w=5.12, gap=0.65, size=18, accent="pure", prefix="GPUBullet")
    box(s, 1.25, 6.05, 10.83, 0.82,
        "I keep the same BC definition. I change how the computation runs on the GPU.",
        "pale_teal", "um", 20, True, name="ContributionBoundary")


def slide_problem(prs):
    s = slide_frame(prs, "Exact All-Sources BC Is Computationally Expensive")
    steps = ["Source\nInitialization", "BFS\nTraversal", "Reverse Dependency\nAccumulation", "BC\nAccumulation"]
    fills = ["light", "pale_blue", "pale_purple", "pale_orange"]
    outlines = ["neutral", "gpu", "pathmerge", "pure"]
    widths = [2.55, 2.35, 3.30, 2.35]
    x = 0.75
    first_x = last_x = 0.0
    for i, (label, fill, outline, w) in enumerate(zip(steps, fills, outlines, widths)):
        box(s, x, 2.05, w, 1.35, label, fill, outline, FIGURE_MIN_PT, True, name=f"Step_{i+1}")
        if i == 0:
            first_x = x + w / 2
        if i == len(steps) - 1:
            last_x = x + w / 2
        else:
            arrow(s, x + w, 2.72, x + w + 0.30, 2.72, "gpu", 2.2, name=f"Flow_{i+1}")
        x += w + 0.30
    line(s, last_x, 3.40, last_x, 4.20, "neutral", 1.7, True, name="Loop_Down")
    line(s, last_x, 4.20, first_x, 4.20, "neutral", 1.7, True, name="Loop_Back")
    arrow(s, first_x, 4.20, first_x, 3.40, "neutral", 1.7, True, name="Loop_To_Source")
    text_box(s, 4.15, 4.30, 5.05, 0.45, "Repeated for every source vertex", FIGURE_MIN_PT, "neutral", True, name="LoopLabel")
    bullets(s, [
        "Exact BC repeats a BFS and a reverse dependency accumulation from every\nvertex, which costs O(|V|(|V|+|E|)).",
        "Frontier size and degree skew create irregular load that a fixed\nparallel granularity handles poorly.",
    ], y=4.95, gap=0.80, size=BODY_PT)
    note(s, "GPU parallelism is the natural response, but grouping more sources makes memory capacity the next constraint.", "gpu")


def slide_four_axes(prs):
    s = slide_frame(prs, "Performance Alone Is Not Enough")
    items = [
        ("Performance", "Is it faster than a tuned\nthird-party implementation?", "pale_blue", "gpu"),
        ("Component\nContributions", "Which execution\ncomponents contributed?", "pale_purple", "pathmerge"),
        ("Memory\nScalability", "How far does the feasible\nbatch range extend?", "pale_orange", "pure"),
        ("Numerical\nCorrectness", "Are the output BC vectors\nvalid?", "pale_green", "chunked"),
    ]
    x = 0.62
    for i, (head, caption, fill, outline) in enumerate(items):
        box(s, x, 1.75, 2.95, 1.65, head, fill, outline, 19, True, name=f"Axis_{i+1}")
        tx(s, x, 3.60, 2.95, 1.30, caption, FIGURE_MIN_PT, "neutral", False, PP_ALIGN.CENTER, f"AxisCaption_{i+1}")
        x += 3.13
    bullets(s, [
        "This work evaluates performance, component contribution, capacity, and\ncorrectness on one framework under consistent conditions.",
        "A runtime number without capacity and correctness evidence does not\nestablish an execution framework as sound.",
    ], y=5.05, gap=0.80)
    note(s, "These four axes correspond to the four research questions of this work.", "gpu")


def slide_framework(prs):
    s = slide_frame(prs, "The Proposal Is an Integrated GPU Execution Framework")
    box(s, 0.60, 2.55, 1.75, 1.20, "Source\nBatch", "light", "neutral", FIGURE_MIN_PT, True, name="SourceBatch")
    for idx, y in ((0, 1.62), (1, 3.92)):
        box(s, 2.85, y, 1.85, 1.05, f"Stream {idx}", "pale_blue", "gpu", FIGURE_MIN_PT, True, name=f"Stream_{idx}")
        arrow(s, 2.35, 3.15, 2.85, y + 0.52, "gpu", 2.0, name=f"Batch_to_Stream{idx}")
        box(s, 5.20, y, 2.05, 1.05, "One Block\nper Source", "light", "gpu", FIGURE_MIN_PT, True, name=f"BlockMap_{idx}")
        arrow(s, 4.70, y + 0.52, 5.20, y + 0.52, "gpu", 2.0, name=f"Stream{idx}_to_Block")
        box(s, 7.75, y, 1.95, 1.05, "Hybrid\nBFS", "pale_teal", "um", FIGURE_MIN_PT, True, name=f"HybridBFS_{idx}")
        arrow(s, 7.25, y + 0.52, 7.75, y + 0.52, "um", 2.0, name=f"Block{idx}_to_BFS")
        box(s, 10.20, y, 2.55, 1.05, "Warp-Cooperative\nAccumulation", "pale_purple", "pathmerge", FIGURE_MIN_PT, True, name=f"WarpAccum_{idx}")
        arrow(s, 9.70, y + 0.52, 10.20, y + 0.52, "pathmerge", 2.0, name=f"BFS{idx}_to_Accum")
    # Kept clear of the 2.67 in bottom edge of the stream-0 row above it.
    box(s, 5.55, 2.78, 3.30, 1.05, "Global BC\nAccumulation", "pale_orange", "pure", FIGURE_MIN_PT, True, name="GlobalBC")
    arrow(s, 11.45, 2.67, 8.85, 3.02, "pure", 2.0, name="Accum0_to_GlobalBC")
    arrow(s, 11.45, 3.92, 8.85, 3.58, "pure", 2.0, name="Accum1_to_GlobalBC")
    text_box(s, 0.60, 5.18, 12.15, 0.45,
             "Dual-Stream Execution: two streams share one framework with independent source-local buffers",
             FIGURE_MIN_PT, "neutral", True, name="FrameworkCaption")
    bullets(s, ["I combine four execution components.",
                "Each component solves a different GPU bottleneck."],
            y=5.65, gap=0.55, size=18)
    note(s, "I do not claim novelty for each component. I claim the integration and its evaluation.", "gpu", 6.78)


def slide_batching(prs, design):
    s = slide_frame(prs, "Source Batching Creates a Batch-Dependent Working Set")
    text_box(s, 0.60, 1.55, 3.0, 0.42, "Source vertices", FIGURE_MIN_PT, "neutral", True, name="SourceLabel")
    for i in range(6):
        circle(s, 0.62 + i * 0.62, 2.05, 0.50, str(i), "light", "neutral", FIGURE_MIN_PT, f"Source_{i}")
    arrow(s, 4.42, 2.30, 5.10, 2.30, "um", 2.2, name="Sources_to_Batch")
    box(s, 5.20, 1.72, 2.30, 1.15, "Requested\nBatch b512", "pale_teal", "um", FIGURE_MIN_PT, True, name="RequestedBatch")
    text_box(s, 7.70, 1.72, 2.6, 1.15, f"NS_eff = {design['NS_eff']}\n(streams)", FIGURE_MIN_PT, "gpu", True, name="NSeff")
    box(s, 10.45, 1.72, 2.30, 1.15, "Per-Source\nState", "pale_orange", "pure", FIGURE_MIN_PT, True, name="PerSourceState")
    arrow(s, 7.50, 2.30, 7.70, 2.30, "gpu", 2.0, name="Batch_to_NSeff")
    arrow(s, 10.25, 2.30, 10.45, 2.30, "pure", 2.0, name="NSeff_to_State")
    box(s, 2.35, 3.35, 8.60, 0.85,
        "Working set = NS_eff × EffectiveBatch × Per-Source State",
        "light", "pure", FIGURE_MIN_PT, True, name="WorkingSetFormula")
    box(s, 0.60, 4.55, 5.85, 1.05, "Input graph file\n(does not grow with batch)", "white", "neutral", FIGURE_MIN_PT, True, name="GraphFileBox")
    box(s, 6.90, 4.55, 5.85, 1.05, "Batch-dependent working set\n(grows with batch)", "pale_orange", "pure", FIGURE_MIN_PT, True, name="WorkingSetBox")
    text_box(s, 0.60, 5.67, 12.15, 0.45, NOT_PARTITION_SENTENCE,
             FIGURE_MIN_PT, "failure", True, name="NotPartitionNote")
    text_box(s, 0.60, 6.14, 12.15, 0.40,
             f"At b512 and NS_eff = {design['NS_eff']}: Estimated allocation ≈ 10.67 GB",
             FIGURE_MIN_PT, "pure", True, name="AllocationEstimate")
    text_box(s, 0.60, 6.56, 12.15, 0.64,
             f"{WORKING_SET_SENTENCE}\nThis is a code-derived estimate. It is not a measured footprint.",
             FIGURE_MIN_PT, "neutral", True, name="WorkingSetBoundary")


def slide_variants(prs):
    s = slide_frame(prs, "Three Memory-Management Variants Share One Framework")
    box(s, 1.35, 1.62, 10.65, 0.85, "Common GPU Execution Framework", "pale_blue", "gpu", 20, True, name="CommonFramework")
    variants = [
        (0.62, "GPU_Opt", "Unified Memory", "Managed allocation;\nhost/HBM placement", "pale_teal", "um"),
        (4.85, "GPU_Opt_Pure", "Device-Only Memory", "Device allocation;\ncapacity-bound", "pale_orange", "pure"),
        (9.08, "GPU_Opt_Pure_Chunked", "Source Sub-Batching", "Limits resident\nsource state", "pale_green", "chunked"),
    ]
    for i, (x, impl, label, desc, fill, outline) in enumerate(variants):
        arrow(s, x + 1.82, 2.47, x + 1.82, 3.00, "gpu", 2.0, name=f"Framework_to_{impl}")
        box(s, x, 3.00, 3.63, 2.30, "", fill, outline, FIGURE_MIN_PT, name=f"Variant_{impl}")
        text_box(s, x + 0.12, 3.18, 3.39, 0.45, impl, 18 if len(impl) < 16 else FIGURE_MIN_PT, "neutral", True, name=f"VariantName_{i}")
        text_box(s, x + 0.12, 3.78, 3.39, 0.45, label, 18, outline, True, name=f"VariantMode_{i}")
        text_box(s, x + 0.12, 4.35, 3.39, 0.80, desc, FIGURE_MIN_PT, "neutral", False, name=f"VariantDesc_{i}")
    bullets(s, ["All three use the same execution framework.",
                "They are not three separate proposals."], y=5.55, gap=0.56, size=18)
    note(s, "Chunked splits the source batch. It does not split the graph.", "gpu", 6.72)


def slide_eval_design(prs):
    s = slide_frame(prs, "I Separate Performance and Capacity Studies")
    rows = [
        ["Graph", "Nodes", "Edges", "Input File (MiB)", "Evaluation Purpose"],
        ["email-EuAll", "265,009", "364,481", "5.59", "Main performance (RQ1)"],
        ["roadNet-PA", "1,088,092", "1,541,898", "28.43", "Main performance (RQ1)"],
        ["roadNet-TX", "1,379,917", "1,921,660", "36.53", "Main performance (RQ1)"],
        ["roadNet-CA", "1,965,206", "2,766,607", "53.83", "Main performance (RQ1)"],
        ["325557 corrected_v1", "325,557", "3,216,152", "43.25", "Capacity tests"],
    ]
    table(s, 0.60, 1.62, 12.15, 3.30, rows, [2.85, 1.95, 1.95, 2.20, 3.20], FIGURE_MIN_PT, "EvaluationDesignTable")
    box(s, 0.60, 5.15, 12.15, 0.70, WORKING_SET_SENTENCE,
        "pale_orange", "pure", 18, True, name="CapacityCaveat")
    bullets(s, ["I use four graphs for main performance.",
                "I use corrected 325557 for capacity tests."],
            y=5.92, gap=0.46, size=18)
    note(s, "All runs use one GH200 GPU.", "gpu", 6.88)


def slide_runtime(prs, data):
    s = slide_frame(prs, "GPU_Opt Reduced Runtime on All Four Evaluated Graphs")
    values = [r["GPU_Opt_Median_s"] for r in data] + [r["PathMerge_Median_s"] for r in data]
    assert all(v > 0 for v in values), "a logarithmic value axis requires strictly positive data"
    assert min(values) >= 10 and max(values) <= 10000
    chart = add_chart(
        s, XL_CHART_TYPE.COLUMN_CLUSTERED, 0.60, 1.55, 12.15, 4.55,
        [GRAPH_SHORT[r["Graph"]] for r in data],
        [("GPU_Opt (fixed b512)", [r["GPU_Opt_Median_s"] for r in data]),
         ("Tuned PathMerge", [r["PathMerge_Median_s"] for r in data])],
        ["gpu", "pathmerge"], "Median runtime (s, log scale)", True, False)
    set_log_scale(chart.value_axis, 10, minimum=10, maximum=10000)
    bump_chart_fonts(chart)
    text_box(s, 0.60, 6.15, 12.15, 0.42,
             "Lower is better. Bars are raw-trial medians; 5 trials for email-EuAll, 3 for each road graph.",
             FIGURE_MIN_PT, "neutral", True, name="RuntimeCaption")
    note(s, "GPU_Opt uses fixed b512. PathMerge uses a tuned batch for each graph.", "gpu")


def slide_speedup(prs, data):
    s = slide_frame(prs, "GPU_Opt Achieved 1.31–3.17× Speedup")
    vals = [r["Speedup"] for r in data]
    chart = add_chart(
        s, XL_CHART_TYPE.COLUMN_CLUSTERED, 0.60, 1.55, 12.15, 4.35,
        [GRAPH_SHORT[r["Graph"]] for r in data],
        [("Speedup", vals), ("Parity 1.0×", [1.0] * len(vals))],
        ["gpu", "neutral"], "Speedup over tuned PathMerge (×)", True, True, '0.00"×"')
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 3.6
    series_to_line(chart, 1, "neutral", 1.8, True)
    bump_chart_fonts(chart)
    batches = "   ".join(f"{GRAPH_SHORT[r['Graph']]}: b{r['PathMerge_Tuned_Batch']}"
                         for r in load_runtime_batches(data))
    text_box(s, 0.60, 5.93, 12.15, 0.35, f"Tuned PathMerge batch per graph — {batches}",
             FIGURE_MIN_PT, "neutral", True, name="TunedBatchCaption")
    text_box(s, 0.60, 6.31, 12.15, 0.42, COMPARATOR_SENTENCE,
             FIGURE_MIN_PT, "failure", True, name="ComparatorCaveat")
    note(s, NO_GENERALIZE_SENTENCE, "failure", 6.78, name="NoGeneralizeNote")


def load_runtime_batches(speedup_rows):
    """Attach the tuned batch to the speedup rows without re-reading raw data."""
    runtime = {r["Graph"]: r for r in RUNTIME_CACHE}
    return [{**r, "PathMerge_Tuned_Batch": runtime[r["Graph"]]["PathMerge_Tuned_Batch"]} for r in speedup_rows]


def slide_ablation(prs, data):
    s = slide_frame(prs, "Hybrid BFS and Dual Streams Gave the Largest Observed Effects")
    lookup = {(r["Graph"], r["Factor"]): r["Main_Effect"] for r in data}
    graphs = ["325557_3216152_corrected_v1", "Synthetic-4 aggregate"]
    names = ["325557 corrected", "Synthetic-4 aggregate"]
    series = [
        ("H: Hybrid BFS", [lookup[(g, "H")] for g in graphs]),
        ("W: Warp-Cooperative", [lookup[(g, "W")] for g in graphs]),
        ("A: Dual-Stream", [lookup[(g, "A")] for g in graphs]),
        ("No effect (1.0x)", [1.0, 1.0]),
    ]
    chart = add_chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, 0.60, 1.55, 12.15, 4.35, names, series,
                      ["gpu", "chunked", "pathmerge", "neutral"], "Main-effect speedup (x)", True, True, '0.000"x"')
    chart.value_axis.minimum_scale = 0.8
    chart.value_axis.maximum_scale = 2.0
    series_to_line(chart, 3, "neutral", 1.5, True)
    bump_chart_fonts(chart)
    text_box(s, 0.60, 5.82, 12.15, 0.42,
             "Synthetic-4 is a mixed-checkpoint aggregate; the other three graphs come from an earlier checkpoint.",
             FIGURE_MIN_PT, "failure", True, name="MixedCheckpointNote")
    text_box(s, 0.60, 6.26, 12.15, 0.42,
             "Corrected 325557: H = 1.4767×   W = 1.1012×   A = 1.5563×",
             FIGURE_MIN_PT, "gpu", True, name="CorrectedEffects")
    note(s, "I did not run this analysis on roadNet graphs. I do not use it to explain the roadNet speedups.", "gpu", 6.70)


def slide_memory(prs, data):
    s = slide_frame(prs, "Memory Variants Expanded the Tested Batch Range")
    x0, x1 = 2.55, 12.35
    xmax = max(r["Requested_Batch"] for r in data) * 1.12
    # Row baselines are spaced so that the failure label hanging below one row and
    # the "Success bN" label rising from the row beneath it cannot collide: the
    # two labels sit 1.40 in apart with 0.22 in and 0.78 in offsets around their
    # own baselines, which leaves a clear gap at the 16 pt deck floor.
    rows = [("GPU_Opt_Pure", "Pure", "pure", 2.50), ("GPU_Opt", "UM", "um", 3.90), ("GPU_Opt_Pure_Chunked", "Chunked", "chunked", 5.30)]
    by_impl = {}
    for r in data:
        by_impl.setdefault(r["Implementation"], []).append(r)
    arrow(s, x0, 5.85, x1, 5.85, "neutral", 1.8, name="BatchAxis")
    for tick in [0] + sorted({r["Requested_Batch"] for r in data}):
        x = x0 + tick / xmax * (x1 - x0)
        line(s, x, 5.76, x, 5.94, "neutral", 1.0)
        text_box(s, x - 0.52, 5.98, 1.04, 0.34, str(tick), FIGURE_MIN_PT, "neutral", False, name=f"Tick_{tick}")
    text_box(s, 4.55, 6.38, 5.80, 0.38, "Requested batch (tested points only)", FIGURE_MIN_PT, "neutral", True, name="BatchAxisLabel")
    for impl, label, col, y in rows:
        text_box(s, 0.55, y - 0.22, 1.85, 0.45, label, 18, col, True, PP_ALIGN.RIGHT, name=f"Row_{label}")
        line(s, x0, y, x1, y, "neutral", 1.0, True)
        for r in by_impl[impl]:
            x = x0 + r["Requested_Batch"] / xmax * (x1 - x0)
            if r["Status"] == "Success":
                circle(s, x - 0.17, y - 0.17, 0.34, "", f"pale_{'orange' if col == 'pure' else 'teal' if col == 'um' else 'green'}", col, FIGURE_MIN_PT)
                text_box(s, x - 0.95, y - 0.78, 1.90, 0.55, f"Success b{r['Requested_Batch']}", FIGURE_MIN_PT, col, True, name=f"OK_{r['Config']}")
            else:
                line(s, x - 0.17, y - 0.17, x + 0.17, y + 0.17, "failure", 3.0)
                line(s, x - 0.17, y + 0.17, x + 0.17, y - 0.17, "failure", 3.0)
                txt = (f"b{r['Requested_Batch']}: CUDA OOM" if r["Status"] == "CUDA OOM"
                       else f"Host/cgroup OOM\nat b{r['Requested_Batch']}")
                fail_y = y + 0.05 if r["Status"] == "CUDA OOM" else y
                fail_h = 0.55 if r["Status"] == "CUDA OOM" else 0.62
                text_box(s, x - 1.35, fail_y, 2.70, fail_h, txt, FIGURE_MIN_PT,
                         "failure", True, name=f"FAIL_{r['Config']}")
    text_box(s, 0.60, 1.24, 12.15, 0.42,
             FAILURE_CATEGORICAL_SENTENCE,
             FIGURE_MIN_PT, "failure", True, name="MemoryCaption")
    note(s, CHUNKED_BOUND_SENTENCE, "failure", 6.85)


def slide_limits(prs):
    s = slide_frame(prs, "The Results Have Clear Limits", "failure")
    bullets(s, [
        "I use one GH200 GPU.",
        "I evaluate four graphs for main performance.",
        "I use one corrected graph for capacity tests.",
        "Trial counts are small.",
        "PathMerge is a retained third-party snapshot.",
    ], y=1.72, gap=0.78, accent="failure")
    box(s, 0.60, 6.05, 12.15, 0.95,
        SCOPE_SENTENCE,
        "pale_red", "failure", 18, True, name="ScopeStatement")


def slide_contributions(prs):
    s = slide_frame(prs, "Contributions")
    items = [
        ("1", "Integrated GPU\nFramework", "I integrate a GPU execution\nframework for exact BC.", "pale_blue", "gpu"),
        ("2", "External GPU\nComparison", "I compare it with a tuned\nexternal GPU implementation.", "pale_teal", "um"),
        ("3", "Component\nEffects", "I measure the effects of\nkey execution components.", "pale_purple", "pathmerge"),
        ("4", "Tested Memory\nCapacity", "I compare the tested capacity\nof three memory variants.", "pale_orange", "pure"),
    ]
    x = 0.62
    for num, head, caption, fill, outline in items:
        circle(s, x + 1.18, 1.68, 0.62, num, fill, outline, 20, f"ContribNum_{num}")
        box(s, x, 2.48, 2.98, 1.35, head, fill, outline, FIGURE_MIN_PT, True, name=f"Contribution_{num}")
        tx(s, x - 0.06, 3.98, 3.10, 1.45, caption, FIGURE_MIN_PT, "neutral", False, PP_ALIGN.CENTER, f"ContribCaption_{num}")
        x += 3.16
    note(s, "I do not present correctness as a main contribution.", "gpu", 5.75, name="NoteScope")
    note(s, "The conclusions apply only to the evaluated conditions.", "cpu", 6.45, name="NoteLimit")


def slide_conclusion(prs):
    s = slide_frame(prs, "The Framework Improved Performance and Expanded the Tested Batch Range")
    box(s, 0.60, 1.70, 12.15, 1.85,
        "I achieved 1.31–3.17× speedup on four evaluated graphs.",
        "pale_blue", "gpu", 18, True, name="ConclusionStatement")
    bullets(s, [
        "I used one fixed GPU_Opt batch.",
        "I used a tuned external comparator.",
        "Hybrid BFS and dual streams gave the largest observed effects.",
        "Memory variants changed the tested feasible batch range.",
        "The conclusions apply only to the evaluated conditions.",
    ], y=3.78, gap=0.52, size=18)
    tx(s, 0.60, 6.50, 12.15, 0.55, "Questions",
       BODY_PT, "gpu", True, PP_ALIGN.CENTER, "Questions")


# --- Backup deck ------------------------------------------------------------
def backup_environment(prs):
    s = slide_frame(prs, "Detailed Experimental Environment", "cpu")
    rows = [
        ["Component", "Specification"],
        ["GPU", "NVIDIA GH200 Grace Hopper Superchip (sm_90)"],
        ["Nominal HBM3", "96 GB"],
        ["Runtime-reported total / free at launch", "approx. 102.0 GB / approx. 101.4 GB (decimal)"],
        ["NVIDIA Driver", "595.58.03"],
        ["CUDA Toolkit (nvcc)", "release 13.0, V13.0.48"],
        ["Host C++ Compiler", "g++ (GCC) 11.4.1"],
        ["Nsight Systems", "2025.5.1.121"],
        ["HBM3 bandwidth (device-to-device)", "1818.6 GB/s (45.2% of theoretical)"],
        ["NVLink-C2C prefetch bandwidth", "177.7 GB/s (19.7% of theoretical)"],
        ["Memory-path resource configuration", "Host-memory-limited 100 GiB configuration"],
        ["Main-experiment aggregation / warmup", "Median of all recorded trials; no warmup, none discarded"],
    ]
    table(s, 0.60, 1.60, 12.15, 4.70, rows, [5.30, 6.85], FIGURE_MIN_PT, "EnvironmentTable")
    note(s, "Items that cannot be established independently from the records are left undetermined, not filled in.", "cpu", 6.50)


def backup_parameters(prs, runtime, design):
    s = slide_frame(prs, "Graph and Batch Parameters", "cpu")
    rows = [["Graph", "GPU_Opt\nBatch", "GPU_Opt\nMedian (s)", "PathMerge\nBatch",
             "PathMerge\nMedian (s)", "Trials"]]
    for r in runtime:
        rows.append([r["Graph"], f"b{r['GPU_Opt_Batch']}", f"{r['GPU_Opt_Median_s']:.2f}",
                     f"b{r['PathMerge_Tuned_Batch']}", f"{r['PathMerge_Median_s']:.2f}",
                     f"{r['GPU_Opt_N']} / {r['PathMerge_N']}"])
    parameter_table = table(
        s, 0.60, 1.60, 12.15, 2.95, rows,
        [2.55, 1.80, 2.20, 1.85, 2.20, 1.55],
        FIGURE_MIN_PT, "ParameterTable")
    # Assigning ``cell.text`` creates one paragraph per newline.  The generic
    # table helper styles the first paragraph, so explicitly style every header
    # paragraph on this two-line backup header.  This avoids renderer-dependent
    # fallback to black on the second line while preserving Graph/Trials
    # alignment, cell geometry, and all data.
    for c, cell in enumerate(parameter_table.table.rows[0].cells):
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        for para in cell.text_frame.paragraphs:
            para.alignment = alignment
            for run in para.runs:
                run.font.name = FONT
                run.font.size = Pt(FIGURE_MIN_PT)
                run.font.bold = True
                run.font.color.rgb = rgb(C["white"])
    box(s, 0.60, 4.70, 12.15, 1.15,
        f"Working set = NS_eff ({design['NS_eff']})  x  EffectiveBatch  x  Per-Source State (10,418,856 bytes)\n"
        "Hybrid BFS thresholds: alpha = %d, beta = %d" % (design["alpha"], design["beta"]),
        "light", "pure", FIGURE_MIN_PT, True, name="ParameterFormula")
    note(s, "Per-Source State is a code-derived allocation estimate, not measured memory usage.", "cpu", 5.98, name="NoteEstimate")
    note(s, "GPU_Opt uses a fixed b512 on every graph; only PathMerge is tuned per graph.", "cpu", 6.58, name="NoteTuning")


def backup_sweep(prs, data):
    s = slide_frame(prs, "PathMerge Batch-Size Sweep", "cpu")
    batches = sorted({r["Requested_Batch"] for r in data})
    series = []
    for graph in GRAPH_ORDER:
        lookup = {r["Requested_Batch"]: r["Median_Runtime_s"] for r in data if r["Graph"] == graph}
        series.append((GRAPH_SHORT[graph], [lookup.get(b) for b in batches]))
    chart = add_chart(s, XL_CHART_TYPE.LINE_MARKERS, 0.60, 1.55, 12.15, 4.60, [str(b) for b in batches],
                      series, ["gpu", "um", "pure", "pathmerge"], "Median runtime (s)", True, False)
    chart.category_axis.has_title = True
    chart.category_axis.axis_title.text_frame.text = "Requested batch"
    bump_chart_fonts(chart)
    clamp = next(r for r in data if r["Graph"] == "email-EuAll" and r["Clamped"] == "yes")
    text_box(s, 0.60, 6.20, 12.15, 0.42,
             f"Requested b{clamp['Requested_Batch']} is effective b{clamp['Effective_Batch']} for email-EuAll. "
             "Historical malformed 325557 is excluded.",
             FIGURE_MIN_PT, "failure", True, name="SweepCaption")
    note(s, "This sweep is what gives PathMerge its best per-graph batch size.", "cpu")


def backup_kernel(prs, data):
    s = slide_frame(prs, "Forced Block-vs-Shared Kernel Comparison", "cpu")
    lookup = {(r["Graph"], r["Kernel"]): r for r in data}
    graphs = ("roadNet-PA", "roadNet-TX")
    series = [("Shared kernel", [lookup[(g, "shared")]["Median_Runtime_s"] for g in graphs]),
              ("Block kernel", [lookup[(g, "block")]["Median_Runtime_s"] for g in graphs])]
    chart = add_chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, 0.60, 1.55, 12.15, 4.35, list(graphs), series,
                      ["cpu", "gpu"], "Median runtime (s)", False, True, "0")
    bump_chart_fonts(chart)
    # Some independent renderers duplicate the first category label beside a
    # native chart legend.  Slide 20 therefore uses only editable PowerPoint
    # shapes for its two-item legend and leaves the chart legend disabled.
    legend_y = 5.57
    for idx, (label, color, marker_x, text_x) in enumerate([
            ("Shared kernel", "cpu", 5.40, 5.68),
            ("Block kernel", "gpu", 7.30, 7.58)], 1):
        marker = s.shapes.add_shape(
            lib.MSO_SHAPE.RECTANGLE, Inches(marker_x), Inches(legend_y),
            Inches(0.18), Inches(0.18))
        marker.fill.solid()
        marker.fill.fore_color.rgb = rgb(C[color])
        marker.line.color.rgb = rgb(C[color])
        marker.name = f"KernelLegendMarker_{idx}"
        tx(s, text_x, legend_y - 0.07, 1.55, 0.34, label,
           FIGURE_MIN_PT, "neutral", False, PP_ALIGN.LEFT,
           f"KernelLegendLabel_{idx}", MSO_ANCHOR.MIDDLE)
    ratios = "   ".join(
        f"{g}: block is {lookup[(g, 'shared')]['Median_Runtime_s'] / lookup[(g, 'block')]['Median_Runtime_s']:.2f}× faster"
        for g in graphs)
    text_box(s, 0.60, 5.95, 12.15, 0.42, ratios, FIGURE_MIN_PT, "gpu", True, name="KernelRatios")
    text_box(s, 0.60, 6.40, 12.15, 0.42,
             "Measured only on roadNet-PA and roadNet-TX; no selector rule is inferred for unmeasured graphs.",
             FIGURE_MIN_PT, "failure", True, name="KernelCaveat")
    note(s, "The current implementation always uses the block kernel.", "cpu", 6.90)


def backup_phase(prs, data):
    s = slide_frame(prs, "Phase Breakdown and Profiling Scope", "cpu")
    lookup = {(r["Graph"], r["Phase"]): r["Median_Component_s"] for r in data}
    series = [(phase, [lookup[(g, phase)] for g in GRAPH_ORDER]) for phase in ("BFS", "Backward", "Other")]
    chart = add_chart(s, XL_CHART_TYPE.BAR_STACKED, 0.60, 1.55, 12.15, 4.55,
                      [GRAPH_SHORT[g] for g in GRAPH_ORDER], series, ["gpu", "pathmerge", "cpu"],
                      "Median component time (s)", True, False)
    bump_chart_fonts(chart)
    text_box(s, 0.60, 6.15, 12.15, 0.42,
             "Complete b512 wall-clock runs: Other = total - BFS - Backward. Not partial Nsight trace totals.",
             FIGURE_MIN_PT, "neutral", True, name="PhaseCaption")
    note(s, "The Nsight trace is a single trace of a single graph and does not generalize to all graphs.", "cpu")


def backup_correctness(prs, rows_tsv):
    s = slide_frame(prs, "Detailed Correctness Evidence", "cpu")
    # The graph moves to its own column so the reference/candidate names fit at the
    # 16 pt deck floor instead of needing a smaller exception size.
    rows = [["Tier", "Graph", "Reference", "Candidate", "Ref/Cand Batch", "Max Rel Err", "Result"]]
    short = {"325557_3216152_corrected_v1": "325557 corr.",
             "benchmark_7000_41459": "7000_41459",
             "benchmark_11023_62184": "11023_62184",
             "chain_200": "chain_200"}
    for r in rows_tsv:
        tier = "A" if r["EvidenceTier"] == "Independent CPU reference" else "B"
        rows.append([tier, short[r["Graph"]], r["Reference"], r["Candidate"],
                     f"{r['ReferenceBatch']} / {r['CandidateBatch']}",
                     r["MaxRelativeError"], r["ToleranceResult"]])
    table(s, 0.60, 1.58, 12.15, 4.55, rows, [0.60, 2.01, 2.72, 2.72, 1.85, 1.40, 0.85],
          FIGURE_MIN_PT, "CorrectnessDetailTable")
    text_box(s, 0.60, 6.12, 12.15, 0.40,
             "All 13 comparisons: MissingIndices = 0, MismatchedElements = 0, ByteIdentical = No.",
             FIGURE_MIN_PT, "neutral", True, name="CorrectnessDetailCaption")
    text_box(s, 0.60, 6.54, 12.15, 0.34, TIER_B_SENTENCE,
             FIGURE_MIN_PT, "failure", True, name="TierBBoundary")
    note(s, CORRECTNESS_ROLE_SENTENCE, "cpu", 6.92)


def backup_historical(prs):
    s = slide_frame(prs, "Historical Record of the Malformed Input", "failure")
    box(s, 0.60, 1.60, 12.15, 0.72,
        HISTORICAL_BANNER,
        "pale_red", "failure", 18, True, name="HistoricalBanner")
    rows = [
        ["Item", "Historical (old input)", "Current (corrected input)"],
        ["Input graph", "325557_3216152 (malformed)", "325557_3216152_corrected_v1 (valid)"],
        ["Defects found", "1-based IDs; 7 missing elements;\nout-of-range IDs", "Reconstructed by symmetry;\nvalidator added"],
        ["Correctness outcome", "CORE_FAIL; stress mismatch at 1e-6", "13 / 13 PASS under mixed tolerance"],
        ["Status in this thesis", "Retained as invalid-input evidence", "Sole basis of current conclusions"],
    ]
    table(s, 0.60, 2.55, 12.15, 3.10, rows, [3.05, 4.55, 4.55], FIGURE_MIN_PT, "HistoricalTable")
    bullets(s, ["Historical results are not used in any current conclusion."],
            y=5.78, gap=0.60, size=18, accent="failure")
    note(s, CORRECTNESS_ROLE_SENTENCE, "failure", 6.52)


# --- Speaker notes ----------------------------------------------------------
# Every slide carries a complete English script and a complete Japanese script.
# They are alternatives, not halves of one narration: either one alone can be
# read aloud for the whole talk, and the main deck sums to the same budget in
# both languages. Neither is a summary of the other.
class Note(NamedTuple):
    slide: int
    seconds: int
    en_script: str
    ja_script: str
    en_transition: str
    ja_transition: str
    limitations: tuple[tuple[str, str], ...]
    questions: tuple[tuple[str, str], ...]


NOTES = [
    Note(
        1, 25,
        "Thank you for the introduction. I will present the design and evaluation of a "
        "batch-based GPU execution framework for betweenness centrality on the GH200. "
        "The target is exact all-sources betweenness centrality, not an approximation, "
        "and every number I report was measured on a single GH200 GPU. I will cover the "
        "problem, the proposed framework, four evaluation axes, and the boundaries of "
        "the evidence.",
        "ご紹介ありがとうございます。本日は、GH200上の媒介中心性計算に向けたバッチ型GPU実行基盤の"
        "設計と評価について発表します。対象は近似ではなく厳密な全始点媒介中心性であり、報告する数値は"
        "すべてGH200一台で測定したものです。問題設定、提案基盤、四つの評価軸、根拠の限界の順で"
        "説明します。",
        "Let me begin with why exact all-sources betweenness centrality is expensive.",
        "まず、厳密な全始点媒介中心性がなぜ計算コストが高いのかから説明します。",
        (("All results come from a single GH200 GPU.",
          "すべての結果はGH200一台での測定に基づく。"),),
        (("Where does this sit relative to existing GPU BC work? It integrates known "
          "components and evaluates them under consistent conditions; it claims no novelty "
          "for any single component.",
          "既存のGPU BC研究との位置づけは。既知の要素を統合し一貫した条件で評価したものであり、"
          "個々の要素の初出性は主張しない。"),),
    ),
    Note(
        2, 55,
        "Betweenness centrality measures how often a vertex lies on shortest paths between "
        "other vertices. Computing it exactly uses Brandes' algorithm: from every source "
        "vertex we run a breadth-first search, then accumulate dependencies in reverse "
        "level order, and add the result into the global betweenness values. Because that "
        "entire loop repeats for every vertex, the cost is big O of |V| times, in "
        "parentheses, |V| plus |E|. Beyond the raw cost, the work is irregular. The "
        "breadth-first frontier changes size at every level, and skewed degree "
        "distributions mean a fixed parallel granularity either starves threads or "
        "serializes them. That combination of high cost and irregular load is what "
        "motivates a GPU execution framework rather than a single optimized kernel.",
        "媒介中心性は、ある頂点が他の頂点間の最短経路上にどれだけ現れるかを表す指標です。厳密に求める"
        "にはBrandesのアルゴリズムを用い、各始点で幅優先探索を行い、レベルの逆順に依存度を累積し、"
        "結果を全体の媒介中心性へ加算します。これを全頂点で繰り返すため、計算量はラージオー、"
        "|V|かける、括弧、|V|＋|E|となります。さらに負荷が不規則です。フロンティアサイズはレベルごとに"
        "変化し、次数分布の偏りにより、並列化の粒度を固定するとスレッドが遊ぶか直列化します。計算"
        "コストの高さと負荷の不規則性が、単一カーネルではなく実行基盤を必要とする理由です。",
        "Given that cost, GPU parallelism is the natural response, but it raises a second question.",
        "このコストを踏まえるとGPU並列化が自然な対応ですが、そこで第二の問題が生じます。",
        (("The target is exact computation, not approximation.",
          "対象は近似ではなく厳密計算である。"),),
        (("Why not use an approximation algorithm? The thesis targets exact all-sources BC, "
          "so approximate methods are out of scope.",
          "なぜ近似ではなく厳密計算か。本研究は厳密な全始点BCを対象としており、近似手法は範囲外である。"),),
    ),
    Note(
        3, 60,
        "Source batching, which groups many source vertices and processes them together, "
        "raises parallelism. But it also holds per-source state for every concurrent "
        "source, so memory capacity becomes a constraint that grows with the batch size. "
        "That means runtime alone cannot tell us whether an execution framework is sound. "
        "We need to know which execution components actually contributed, how far the "
        "feasible batch range extends before the run fails, and whether the betweenness "
        "vectors that come out are numerically valid. These four axes, performance, "
        "component contribution, memory scalability, and numerical correctness, are the "
        "four research questions of this work, and I evaluate all four on the same "
        "framework under consistent conditions. The other three axes are not decoration "
        "around the performance number; they are what makes it interpretable.",
        "始点をまとめて処理するsource batchingは並列性を高めます。しかし実行中の各始点の状態を保持"
        "するため、バッチサイズとともにメモリ容量が制約になります。つまり実行時間だけでは実行基盤"
        "として妥当か判断できません。どの実行要素が寄与したのか、失敗するまでバッチをどこまで大きく"
        "できるのか、出力されるBCベクトルが数値的に妥当かを知る必要があります。性能、要因寄与、"
        "メモリ容量、数値的正確性というこの四観点が本研究の四つのResearch Questionであり、同一基盤に"
        "対し一貫した条件で四つすべてを評価します。残る三観点は性能値の装飾ではなく、それを解釈"
        "可能にするものです。",
        "With those four axes fixed, let me show the framework itself.",
        "この四つの評価軸を踏まえて、実行基盤そのものを説明します。",
        (("A runtime number without capacity and correctness evidence does not establish a framework as sound.",
          "容量と正確性の根拠を欠いた実行時間だけでは、実行基盤としての妥当性は示せない。"),),
        (("How do the four axes relate? Capacity bounds which batch sizes performance can "
          "even be measured at, and correctness bounds whether any of those runs are meaningful.",
          "四つの観点はどう関係するか。容量はどのバッチで性能を測定できるかを規定し、正確性は"
          "その実行結果が意味を持つかを規定する。"),),
    ),
    Note(
        4, 85,
        "The proposal is a single GPU execution framework that integrates four components "
        "into one execution flow. First, block-based source assignment: each source vertex "
        "in the batch is mapped to one thread block, so the per-source traversal state "
        "stays block-local. Second, hybrid breadth-first search, which switches between "
        "top-down and bottom-up expansion according to frontier size, so that the wide "
        "levels in the middle of a traversal do not dominate. Third, warp-cooperative "
        "accumulation, which lets a warp cooperate on the reverse dependency accumulation "
        "instead of leaving one thread per vertex. Fourth, dual-stream execution, where two "
        "CUDA streams each hold independent source-local buffers and overlap initialization "
        "with computation. I want to be explicit about the claim here: none of these four "
        "techniques is claimed as new. The contribution is that they are integrated into "
        "one coherent framework and evaluated together under consistent conditions. The "
        "diagram shows the two streams running that same pipeline in parallel, each merging "
        "its partial results into the global betweenness array in the middle. The buffers "
        "are source-local, which is the property the overlap is designed around.",
        "提案手法は、四つの実行要素を一つの実行フローへ統合したGPU実行基盤です。第一に"
        "block-based source assignment。各始点を一スレッドブロックへ割り当て、探索状態をブロック内に"
        "閉じます。第二にHybrid BFS。フロンティアサイズに応じてtop-downとbottom-upを切り替え、"
        "中盤の広いレベルの支配を避けます。第三にWarp-Cooperative Accumulation。依存度累積を頂点"
        "あたり一スレッドではなくワープで協調処理します。第四にDual-Stream Execution。二つの"
        "CUDAストリームが独立した始点ローカルバッファを持ち初期化と計算を重ねます。四つとも新規性は"
        "主張しません。貢献はこれらを一貫した基盤へ統合し一貫した条件で評価した点です。図は二つの"
        "ストリームが同じパイプラインを並列実行し、中央の全体BC配列へ部分結果を統合する様子です。"
        "バッファは始点ローカル、重ね合わせはこれを前提とします。",
        "The batching that makes this framework fast is also what creates its capacity constraint.",
        "この基盤を高速にしているバッチ化は、同時に容量制約を生む要因でもあります。",
        (("No novelty is claimed for the individual components; the contribution is the integration and its evaluation.",
          "個々の要素技術の初出性は主張しない。貢献は統合とその評価である。"),),
        (("How does this differ from prior work? The components are known individually; "
          "the framework and the four-axis evaluation on one GH200 are what this work adds.",
          "既存研究とどこが異なるか。各要素は個別には既知であり、本研究が加えるのは統合された基盤と"
          "GH200上での四観点評価である。"),),
    ),
    Note(
        5, 75,
        "This slide defines where the capacity constraint actually comes from, because it is "
        "easy to get wrong. Source batching groups source vertices; it does not partition "
        "the graph. Every source in the batch still traverses the whole graph. What grows "
        "with the batch is the per-source state: the distance array, the sigma counts, the "
        "delta values, and the frontier structures, one copy per concurrent source. So the "
        "working set is the product of the effective number of streams, the effective batch "
        "size, and the per-source state. The consequence is the point I want you to take "
        "away: input graph size and the batch-dependent working set are different "
        "quantities. Even a small input graph can require a much larger batch-dependent "
        "working set. At b512 with two effective streams, the code-derived allocation "
        "estimate is about 10.67 gigabytes; this is an estimate, not a measured footprint. "
        "In the diagram, the box on the lower left is the input graph file, which does not "
        "grow with the batch, and the box on the lower right is the batch-dependent working "
        "set, which does.",
        "容量制約がどこから生じるかを定義します。誤解されやすい点です。source batchingは始点を"
        "まとめるもので、グラフを分割しません。各始点は依然として全体を探索します。"
        "バッチとともに増えるのは始点ごとの状態、すなわち距離配列、シグマ、デルタ、フロンティア"
        "構造であり、同時実行する始点の数だけ必要です。したがってworking setは、実効ストリーム数、"
        "実効バッチサイズ、始点ごとの状態量の積です。要点は、入力グラフのサイズとbatch依存の"
        "working setが別の量だということです。入力が小さくてもworking setははるかに大きくなり得ます。"
        "b512かつ実効ストリーム数2では、コード由来の割り当て推定値は約10.67ギガバイトです。これは"
        "推定値であり実測値ではありません。図の左下の箱が入力ファイルでバッチとともに増えず、"
        "右下の箱がbatch依存のworking setです。",
        "That distinction is what separates the three memory-management variants.",
        "この区別が、三つのメモリ管理方式を分ける基準になります。",
        (("Source batching is not graph partitioning; the input file size is not the working set.",
          "source batchingはgraph partitioningではなく、入力ファイルサイズはworking setではない。"),),
        (("Does a larger batch always run faster? No. Larger batches raise parallelism but "
          "also enlarge the working set, and past a point the run fails outright.",
          "バッチを大きくすれば必ず速くなるのか。ならない。並列性は上がるがworking setも増大し、"
          "ある点を超えると実行自体が失敗する。"),),
    ),
    Note(
        6, 60,
        "Because the working set is what binds, the three implementations differ only in how "
        "they manage memory. GPU_Opt uses unified memory, so allocations are managed and the "
        "runtime places pages between HBM and host memory. GPU_Opt_Pure uses explicit device "
        "allocation and is bounded by device capacity. GPU_Opt_Pure_Chunked sub-batches the "
        "source set so that only part of the source state is resident at a time. I want to "
        "be precise about two things. These are not three separate proposals; they are "
        "memory-management variants of one common framework, sharing the same kernels. And "
        "what Chunked subdivides is the source set, not the graph itself. The variants were "
        "designed to isolate differences in memory-management behavior while sharing the "
        "common execution framework, which is how the capacity comparison on the later "
        "slides should be read.",
        "制約はworking setであるため、三実装はメモリの扱い方だけが異なります。GPU_OptはUnified Memory"
        "を用い、ランタイムがHBMとホスト間にページを配置します。GPU_Opt_Pureは明示的な"
        "デバイス割り当てを用い、デバイス容量に制限されます。GPU_Opt_Pure_Chunkedは始点集合を分割し、"
        "始点状態の一部だけを常駐させます。二点重要です。これらは独立した三提案ではなく、同一カーネル"
        "を共有する共通基盤上のメモリ管理の違いです。またChunkedが分割するのは始点集合であり"
        "グラフではありません。共通基盤を保ったままメモリ管理の挙動差を切り分ける設計であり、"
        "容量比較もその観点で読みます。",
        "With the framework and its variants defined, let me turn to how they were evaluated.",
        "基盤とその変種を定義したので、次にどのように評価したかを説明します。",
        (("The three are variants of one framework, and Chunked subdivides sources, not the graph.",
          "三つは同一基盤の変種であり、Chunkedが分割するのは始点であってグラフではない。"),),
        (("Which variant should be used in practice? The comparison characterises capacity "
          "behaviour rather than ranking them; the right choice depends on the batch size needed.",
          "実運用ではどれを選ぶべきか。この比較は優劣を決めるものではなく容量特性を示すものであり、"
          "必要なバッチサイズに依存する。"),),
    ),
    Note(
        7, 55,
        "The evaluation deliberately separates two studies. Main performance uses four "
        "graphs: email-EuAll, and the three road networks roadNet-PA, roadNet-TX, and "
        "roadNet-CA. The memory and correctness study uses one graph, the corrected 325557 "
        "graph. The file sizes in this table, from about five to about fifty-four "
        "mebibytes, are the static size of the input on disk. As I said on the previous "
        "slide, that is not the batch-dependent working set, and I show both here so the "
        "two are not conflated. All measurements are on one GH200 GPU, reported values are "
        "medians over the recorded trials, and every speedup is a median-to-median ratio. "
        "The scope is deliberately small, and I will return to what that bounds at the end.",
        "評価は二つの研究を意図的に分離しています。主性能比較にはemail-EuAllと道路ネットワーク"
        "roadNet-PA、TX、CAの四グラフ、メモリと正確性の検証には修正版325557を用います。表の"
        "ファイルサイズ、およそ5から54メビバイトはディスク上の静的サイズです。前スライドの"
        "とおりbatch依存のworking setではなく、混同しないよう併記しています。測定はすべてGH200一台、"
        "報告値は記録した試行のmedian、speedupはmedian同士の比です。評価範囲は意図的に小さく、"
        "それが何を制約するかは最後に述べます。",
        "Let me start with the main performance result.",
        "まず主性能の結果から示します。",
        (("Input graph size is not the batch-dependent working set.",
          "入力グラフサイズはbatch依存のworking setではない。"),),
        (("Why only this many graphs? The scope was set by what could be measured and "
          "verified under consistent conditions; the conclusions are limited accordingly.",
          "なぜこのグラフ数なのか。一貫した条件で測定・検証できる範囲で設定しており、結論もその範囲に"
          "限定される。"),),
    ),
    Note(
        8, 75,
        "This is the main performance result. The blue bars are GPU_Opt at a fixed batch "
        "size of 512, and the purple bars are PathMerge, the external comparator, at the "
        "best batch size found for each graph by a separate sweep. Note the vertical axis is "
        "logarithmic. GPU_Opt reduced runtime on all four evaluated graphs. I want to be "
        "clear about how the comparison is set up, because it is deliberately unfavourable "
        "to the proposal: our side runs one fixed configuration across every graph, while "
        "the comparator is tuned per graph. The bars are medians of raw trials, five trials "
        "for email-EuAll and three for each road graph. Those are small samples, and I will "
        "return to that when I state the limitations. Reading the chart, the leftmost pair is "
        "email-EuAll, where the gap is widest, and the three pairs to the right are the road "
        "networks, where the gap is narrower but consistent. Because the axis is logarithmic, "
        "equal visual gaps mean equal ratios rather than equal differences in seconds.",
        "これが主性能の結果です。青い棒がバッチサイズ512に固定したGPU_Opt、紫の棒が外部比較対象の"
        "PathMergeで、後者は別途の掃引によりグラフごとに最良バッチを与えています。縦軸が対数目盛で"
        "ある点にご注意ください。GPU_Optは評価した四グラフすべてで実行時間を短縮しました。"
        "この比較設定は意図的に提案側へ不利です。提案側は全グラフで一つの固定設定を用いる一方、"
        "比較対象はグラフごとに調整されているためです。棒は生の試行のmedianで、試行数はemail-EuAllが"
        "5回、各道路グラフが3回です。小標本であり、限界の説明で改めて触れます。図の左端が"
        "email-EuAllで差が最も大きく、右側の三組が道路ネットワークで、差は小さいものの一貫して"
        "います。対数目盛であるため、見た目の間隔が等しいことは秒数差ではなく比が等しいことを"
        "意味します。",
        "Expressing the same data as a ratio gives the headline number.",
        "同じデータを比として表すと、主要な数値が得られます。",
        (("Trial counts are small: five for email-EuAll and three per road graph.",
          "試行数は小さく、email-EuAllが5回、各道路グラフが3回である。"),),
        (("Isn't the trial count too small? Yes, it is a real limitation; the medians are "
          "reported as such and no distributional claim is made.",
          "試行数が少ないのではないか。実際に限界であり、medianとして報告し分布に関する主張はしない。"),),
    ),
    Note(
        9, 90,
        "Expressed as a ratio, GPU_Opt achieved a 3.17 times speedup on email-EuAll, 1.31 "
        "times on roadNet-PA, 1.51 times on roadNet-TX, and 1.45 times on roadNet-CA. Each "
        "is a median-to-median ratio, and the dashed line marks parity at one. The tuned "
        "batch sizes for the comparator were b2048 for email-EuAll, b64 for roadNet-PA, b64 "
        "for roadNet-TX, and b32 for roadNet-CA. Now the important qualification. PathMerge is an "
        "evaluated third-party implementation and an external comparator, not ground truth. "
        "It is a snapshot we retained and measured; it is not the original authors' official "
        "implementation, and it is not a correctness reference. The result does not "
        "generalize to PathMerge implementations in general, and it does not generalize "
        "beyond these four graphs on this one GPU. The selected batch sizes differed "
        "substantially across the four graphs. This evaluation records that sensitivity but "
        "does not establish a single structural cause. Our side used b512 throughout, which "
        "is not the tuned setting for any of them.",
        "比として表すと、GPU_Optはemail-EuAllで3.17倍、PAで1.31倍、TXで1.51倍、CAで1.45倍の"
        "speedupを達成しました。いずれもmedian同士の比で、破線が等速の1.0倍です。比較対象の調整済み"
        "バッチはemail-EuAllがb2048、PAとTXがb64、CAがb32でした。重要な限定です。PathMergeは"
        "第三者実装かつ外部の比較対象であってground truthではなく、保存して測定したsnapshotで、"
        "公式実装でも正確性の参照でもありません。結果はPathMerge実装一般へも、この四グラフと一台の"
        "GPUを超えても一般化しません。"
        "採用バッチは四グラフ間で大きく異なりました。本評価はその感度を記録しますが、単一の構造的"
        "原因を確定しません。提案側はb512固定で、いずれのグラフの調整済み設定でもありません。",
        "The next question is which parts of the framework produced that gain.",
        "次の問いは、この向上が基盤のどの部分から生じたのかです。",
        (("PathMerge is an external comparator, not ground truth, and the result does not "
          "generalize to PathMerge in general.",
          "PathMergeは外部比較対象でありground truthではない。結果はPathMerge一般へ一般化しない。"),),
        (("Was the comparator tuned enough? A dedicated batch sweep, shown in the backup "
          "slides, selected its best batch per graph, while our side stayed fixed at b512.",
          "比較対象の調整は十分か。Backupに示す専用の掃引でグラフごとの最良バッチを選定しており、"
          "提案側はb512に固定している。"),),
    ),
    Note(
        10, 75,
        "To attribute the gain, I ran a factor decomposition over the three compile-time "
        "components: H for hybrid BFS, W for warp-cooperative accumulation, and A for "
        "dual-stream execution. On the corrected 325557 graph the main effects were 1.4767 "
        "for hybrid BFS, 1.1012 for warp-cooperative accumulation, and 1.5563 for "
        "dual-stream execution. On the synthetic-four aggregate they were 1.6787, 1.0661, "
        "and 1.3914. So hybrid BFS and dual-stream execution contributed most, while the "
        "warp-cooperative effect was smaller and differed between the two settings. Two caveats. The "
        "synthetic-four figure is an aggregate across a mixed set of checkpoints, which I "
        "note on the slide. And no factor decomposition was run on the road networks at all, "
        "so these numbers do not explain the main performance gap I just showed. Each bar is "
        "a main effect: the speedup attributable to enabling that one component, averaged "
        "across the full factorial of the three flags. Because these are main effects, they "
        "do not simply multiply together to give the overall speedup.",
        "要因分解は、三つのコンパイル時要素、HがHybrid BFS、WがWarp-Cooperative "
        "Accumulation、AがDual-Stream Executionについて行いました。修正版325557では主効果が"
        "H 1.4767倍、W 1.1012倍、A 1.5563倍、合成4グラフの集約では1.6787倍、1.0661倍、1.3914倍です。"
        "HとAの寄与が大きく、Wの効果は相対的に小さく条件により異なりました。注意点が二つ。合成4"
        "グラフの値はcheckpoint混在の集約であり、スライドにも明記しています。また道路網では"
        "要因分解を実施しておらず、これらの数値は主性能差を説明しません。各棒は主効果、すなわち三"
        "フラグの完全実施要因計画で平均した、その要素の有効化によるspeedupで、積算しても全体の"
        "speedupにはなりません。",
        "That covers performance. The next axis is capacity.",
        "性能についてはここまでです。次の観点は容量です。",
        (("No factor decomposition was run on roadNet, so these values do not explain the main performance gap.",
          "roadNetでは要因分解を実施しておらず、これらの値は主性能差を説明しない。"),),
        (("Why did the warp effect differ between the two settings? The measured values "
          "differ, but this evaluation does not establish a single cause, and no general "
          "rule is inferred from them.",
          "Warpの効果が二条件で異なったのはなぜか。測定値は異なるが、本評価は単一の原因を確定する"
          "ものではなく、これらから一般則も導かない。"),),
    ),
    Note(
        11, 80,
        "This is the capacity result on the corrected 325557 graph, one row per memory "
        "variant against requested batch size. Pure succeeded at b4096 and failed at b8192 "
        "with a CUDA device out-of-memory error. Unified memory succeeded at b10240 and "
        "failed at b12288, and that failure was different in kind: it was recorded as a host "
        "cgroup memory limit out-of-memory kill, exit 137, and not a device error. The two "
        "failure classes are reported separately and are not treated as one boundary. "
        "Chunked succeeded at b16384. Two things to read carefully. "
        "Failures are categorical outcomes, not zero-second runtimes, so a cross on this "
        "chart is a classification, not a measurement of zero. And Chunked succeeded at the "
        "tested upper bound of b16384; this does not imply unlimited capacity. That is the "
        "largest batch we tried, not a limit we found. Every point is a single targeted "
        "feasibility run. The point to take away is that the three rows failed at different "
        "batch sizes and with different failure classes, and exposing that difference is why "
        "the three memory-management variants exist.",
        "修正版325557の容量結果を、要求バッチサイズに対し方式ごとに示します。Pureはb4096で成功、"
        "b8192ではCUDAのデバイスout-of-memoryエラーで失敗しました。UMはb10240で成功、"
        "b12288で失敗しましたが、この失敗は種類が異なり、デバイス側エラーではなくホストのcgroup"
        "メモリ制限によるOOM kill、exit 137です。二つの失敗クラスは分けて報告し、単一境界としては"
        "扱いません。Chunkedはb16384で成功しました。注意点が二つ。失敗は区分された結果で0秒の実行"
        "時間ではなく、×印は分類で測定値0ではありません。またChunkedは試験上限で成功しま"
        "したが、容量が無制限とは言えません。これは試した最大のバッチであり限界ではありません。"
        "各点は1回のtargetedな実行可能性確認です。要点は三行が異なるバッチと失敗クラスで失敗した"
        "ことで、その差の可視化が三方式の理由です。",
        "That is capacity. The remaining axis is whether the output is numerically valid.",
        "容量については以上です。残る観点は、出力が数値的に妥当かどうかです。",
        (("Failures are categorical outcomes, not zero-second runtimes, and b16384 is a tested "
          "upper bound, not a demonstrated capacity limit.",
          "失敗は0秒の実行時間ではなく区分された結果であり、b16384は試験上限であって容量限界の"
          "証明ではない。"),),
        (("Does unified memory remove the memory constraint? No. It moved the boundary and "
          "changed the failure mode, but it still failed, at b12288, on the host side.",
          "UMを使えばメモリ制約は解消するのか。しない。境界と失敗の様態が変わっただけで、b12288で"
          "ホスト側で失敗している。"),),
    ),
    Note(
        12, 60,
        "Correctness was checked in two tiers. Tier A uses an independent Sequential CPU "
        "reference, comparing full betweenness vectors on three small graphs. Tier B "
        "compares implementation paths against each other on the corrected 325557 graph, "
        "ten comparisons. That gives thirteen comparisons in total, with zero missing "
        "indices and zero mismatched elements. All 13 comparisons passed the mixed "
        "tolerance, but none was byte-identical. Two qualifications matter here. Agreement "
        "under a tolerance is numerical agreement, not bitwise identity. The non-byte-identical "
        "results are consistent with different floating-point operation orders, but this "
        "evaluation does not establish a single cause. And Tier B compares implementation "
        "paths and is not an independent ground-truth evaluation, so the independent evidence "
        "is the three Tier A comparisons. The table breaks this down by tier, with the "
        "missing and mismatched columns reading zero on every row and the byte-identical "
        "column reading No on every row.",
        "正確性は二つのTierで検証しました。Tier Aは独立したSequential CPU参照を用い、小規模な三"
        "グラフでBCベクトル全体を比較します。Tier Bは修正版325557上で実装経路どうしを比較する10件"
        "です。合計13比較で、欠損インデックス0件、不一致要素0件。すべてmixed toleranceの下で合格"
        "しましたが、byte一致は一つもありません。二点の限定が重要です。許容誤差下の一致は"
        "数値的一致でありビット単位の同一性ではありません。byte非一致は浮動小数点演算順序の違いと"
        "整合しますが、本評価は単一の原因を確定しません。またTier Bは実装経路間の比較で、独立した"
        "ground truth評価ではありません。独立な根拠はTier Aの3比較です。",
        "Having covered all four axes, let me state where the evidence stops.",
        "四つの観点をすべて説明したので、次に根拠がどこで止まるかを述べます。",
        (("Tier B is not an independent ground-truth evaluation, and no comparison was byte-identical.",
          "Tier Bは独立したground truthとの評価ではなく、byte一致した比較は一つもない。"),),
        (("Is the absence of byte-identity a problem? It is not treated as a failure here "
          "because all comparisons passed the predefined mixed tolerance with zero mismatched "
          "elements. The exact cause of the byte difference is not independently established.",
          "byte一致でないことは問題か。ここでは失敗として扱っていない。全比較が事前に定めた"
          "mixed toleranceを不一致要素0で満たしたためである。byte差の正確な原因は独立には"
          "確定していない。"),),
    ),
    Note(
        13, 50,
        "Let me be explicit about the boundaries. The evaluation used a single GH200 GPU and "
        "does not generalize to other GPUs. Main performance covers four graphs, and the "
        "memory and correctness study uses corrected 325557 only. Trial counts are small, "
        "and each capacity boundary is a single targeted validation run. PathMerge is a "
        "retained third-party snapshot, so nothing here is a claim about PathMerge in "
        "general. And unified memory and Chunked have finite capacity limits; I do not claim "
        "they avoid out-of-memory failures. Stated together: the conclusions are limited to "
        "one GH200 GPU, the evaluated graphs, the retained implementation snapshots, and the "
        "recorded experimental conditions. These are not retractions; they mark the range "
        "within which the conclusions hold, and I would rather state them plainly than leave "
        "them to be inferred.",
        "根拠の境界を述べます。評価はGH200一台であり、他のGPUへ一般化しません。主性能は四グラフ、"
        "メモリと正確性の検証は修正版325557のみです。試行数は小さく、容量境界は各条件1回の"
        "targeted validationです。PathMergeは保存した第三者実装のsnapshotであり、PathMerge一般に"
        "関する主張ではありません。UMとChunkedの容量にも上限があり、OOM回避は主張しません。本研究の"
        "結論はGH200一台、評価したグラフ、保存した実装snapshot、記録された実験条件の範囲に限定"
        "されます。これらは結論の否定ではなく、成立範囲を示すものです。",
        "Within those boundaries, the contributions are as follows.",
        "この範囲の内側で、貢献は次のとおりです。",
        (("The limitations mark the range in which the conclusions hold; they are not retractions.",
          "限定は結論が成立する範囲を示すものであり、結論の撤回ではない。"),),
        (("What can be claimed under these limits? That on this hardware and these graphs, "
          "the integrated framework was faster and its capacity boundaries were characterised.",
          "この限定の下で何が言えるか。このハードウェアとこれらのグラフにおいて、統合基盤が高速で"
          "あり、その容量境界を特徴づけたということである。"),),
    ),
    Note(
        14, 30,
        "The contributions are four. First, the design and implementation of a GPU execution "
        "framework integrating existing components. Second, a performance evaluation against "
        "a tuned third-party comparator. Third, a component-level contribution analysis "
        "quantifying the H, W, and A effects including their graph dependence. Fourth, a "
        "boundary analysis that treats memory capacity and numerical correctness "
        "separately. Again, no novelty is claimed for the individual components. Each of the "
        "four is evaluated in the body of the thesis under the conditions I have described.",
        "貢献は四点です。第一に、既存の要素を統合したGPU実行基盤の設計と実装です。第二に、調整済みの"
        "第三者実装を比較対象とした性能評価です。第三に、H、W、Aの効果をグラフ依存性を含めて定量化した"
        "要素レベルの寄与分析です。第四に、メモリ容量と数値的正確性を分離して扱った境界分析です。"
        "繰り返しになりますが、個々の要素技術の初出性は主張しません。",
        "Let me close with the central conclusion.",
        "最後に中心となる結論を述べます。",
        (("Novelty is claimed for the integration and evaluation, not for the components.",
          "新規性は統合と評価について主張するものであり、個々の要素についてではない。"),),
        (("Which contribution is the most essential? The integration together with the "
          "four-axis evaluation, since no single component is claimed as new.",
          "最も本質的な貢献はどれか。個々の要素の新規性を主張しない以上、統合と四観点評価である。"),),
    ),
    Note(
        15, 25,
        "To conclude: the integrated block-based GPU implementation achieved a 1.31 to 3.17 "
        "times speedup over the tuned comparator on the four evaluated graphs, and the "
        "memory-path experiments clarified both the feasible batch ranges and the remaining "
        "numerical limitations. The gain came from multiple components of the framework "
        "rather than one technique. Thank you for your attention, and I am happy to take "
        "questions.",
        "結論です。統合されたblock-basedのGPU実装は、評価した四グラフで調整済みの比較対象より1.31倍"
        "から3.17倍高速であり、メモリ経路の実験は実行可能なバッチ範囲と数値的限界の双方を"
        "明らかにしました。この向上は単一技術ではなく複数要素から生じました。ご清聴ありがとう"
        "ございました。",
        "(End of the main talk; backup slides follow for questions.)",
        "（本編はここまでです。以降のBackupは質疑応答用です。）",
        (("The conclusion holds within the evaluated GPU, graphs, and implementation snapshots.",
          "結論は評価したGPU、グラフ、実装snapshotの範囲で成立する。"),),
        (("What is the future work? Broader graph coverage, more trials, and factor "
          "decomposition on the road networks, which was not performed here.",
          "今後の課題は何か。グラフ範囲の拡大、試行数の増加、そして本研究では未実施の道路"
          "ネットワークにおける要因分解である。"),),
    ),
    Note(
        16, 40,
        "This backup slide gives the full environment. The GPU is an NVIDIA GH200 Grace "
        "Hopper Superchip at sm_90 with a nominal 96 gigabytes of HBM3, driver 595.58.03, "
        "CUDA 13.0, and g++ 11.4.1. Measured device-to-device HBM3 bandwidth was 1818.6 "
        "gigabytes per second, and NVLink-C2C prefetch bandwidth was 177.7. The main "
        "experiments report the median of all recorded trials with no warmup and nothing "
        "discarded. Items that could not be established independently from the records were "
        "left undetermined rather than filled in.",
        "このBackupスライドは実験環境の全体を示します。GPUはNVIDIA GH200 Grace Hopper Superchipの"
        "sm_90で、公称HBM3容量は96ギガバイト、ドライバは595.58.03、CUDAは13.0、g++は11.4.1です。"
        "実測のデバイス間HBM3帯域は毎秒1818.6ギガバイト、NVLink-C2Cのプリフェッチ帯域は177.7でした。"
        "主実験は記録した全試行のmedianを報告し、warmupは行わず、破棄した試行もありません。記録から"
        "独立に確定できない項目は、補完せず未確定のまま扱っています。",
        "(Backup slide, shown on request.)",
        "（Backupスライド。質疑応答で必要に応じて表示する。）",
        (("Values that could not be independently established from the records are left undetermined.",
          "記録から独立に確定できない値は未確定のままとしている。"),),
        (("Why is the runtime-reported memory larger than 96 GB? The nominal capacity and "
          "runtime-reported values come from different records and units. They are shown "
          "separately without asserting that they are equivalent or identifying a single "
          "cause for the difference.",
          "実行時報告値が96GBより大きいのはなぜか。公称容量と実行時報告値は異なる記録と単位に"
          "由来する。両者は同値と断定せず、差の単一原因も特定せずに分けて示している。"),),
    ),
    Note(
        17, 40,
        "This backup slide lists the per-graph batch settings, medians, and trial counts, "
        "together with the working-set formula: the effective number of streams times the "
        "effective batch size times the per-source state, which is 10,418,856 bytes. The "
        "hybrid BFS alpha and beta thresholds are also given. Two points: the per-source "
        "state is an allocation estimate derived from the code, not a measured memory "
        "footprint, and GPU_Opt uses a fixed b512 on every graph while only PathMerge is "
        "tuned per graph.",
        "このBackupスライドは、グラフごとのバッチ設定、median、試行数に加えて、working setの算出式を"
        "示します。式は実効ストリーム数かける実効バッチサイズかける始点ごとの状態量であり、状態量は"
        "10,418,856バイトです。Hybrid BFSのalphaとbetaの閾値も併せて示しています。二点補足します。"
        "始点ごとの状態量はコードから導いた割り当ての推定値であり、実測のメモリ使用量ではありません。"
        "またGPU_Optは全グラフでb512に固定されており、グラフごとに調整しているのはPathMergeのみです。",
        "(Backup slide, shown on request.)",
        "（Backupスライド。質疑応答で必要に応じて表示する。）",
        (("Per-Source State is a code-derived allocation estimate, not measured memory usage.",
          "始点ごとの状態量はコード由来の割当推定であり、実測メモリ使用量ではない。"),),
        (("Why not tune GPU_Opt per graph as well? Leaving it fixed keeps the comparison "
          "unfavourable to the proposal rather than favourable.",
          "なぜGPU_Optもグラフごとに調整しないのか。固定のままとすることで、比較が提案側に有利に"
          "ならないようにしている。"),),
    ),
    Note(
        18, 40,
        "This backup slide shows the PathMerge batch sweep that justifies the tuned "
        "comparator settings. Median runtime is plotted against requested batch size for "
        "each of the four graphs, and the minimum of each curve is the batch reported on the "
        "speedup slide. One detail: for email-EuAll the requested batch is clamped to a "
        "smaller effective batch, which is noted on the slide. The historical malformed "
        "325557 graph is excluded from this sweep.",
        "このBackupスライドは、比較対象の調整設定を裏づけるPathMergeのバッチ掃引を示します。四つの"
        "グラフそれぞれについて、要求バッチサイズに対するmedian実行時間を示しており、各曲線の最小値が"
        "speedupのスライドで報告したバッチです。一点補足すると、email-EuAllでは要求バッチがより小さい"
        "実効バッチへクランプされており、その旨をスライドに注記しています。履歴的なmalformed版の"
        "325557グラフはこの掃引から除外しています。",
        "(Backup slide, shown on request.)",
        "（Backupスライド。質疑応答で必要に応じて表示する。）",
        (("The sweep is what gives PathMerge its best per-graph batch; the malformed graph is excluded.",
          "この掃引がPathMergeへグラフごとの最良バッチを与えている。malformed版グラフは除外している。"),),
        (("Could a batch outside the swept range be faster? Possibly; the sweep covers the "
          "tested points only and no claim is made beyond them.",
          "掃引範囲外のバッチがより速い可能性はあるか。あり得る。掃引は試験した点のみを対象とし、"
          "その外側については主張しない。"),),
    ),
    Note(
        19, 40,
        "This backup slide compares the block kernel against the shared-memory kernel by "
        "forcing each one, on roadNet-PA and roadNet-TX. The chart gives the median runtime "
        "for each kernel on each graph, with the per-graph ratio annotated beneath it. The "
        "block kernel was faster on both measured graphs. These measurements support the "
        "current block-kernel choice within the measured PA and TX scope, but they do not "
        "define a selector rule for unmeasured graphs.",
        "このBackupスライドは、blockカーネルとshared memoryカーネルをそれぞれ強制的に選択して、"
        "roadNet-PAとroadNet-TXで比較したものです。図は各グラフ・各カーネルのmedian実行時間を示し、"
        "下部にグラフごとの比を注記しています。測定した両グラフでblockカーネルの方が高速でした。"
        "この測定は、測定範囲であるPAとTXの内側では現行のblockカーネル選択を支持しますが、"
        "測定していないグラフに対する選択規則を定めるものではありません。",
        "(Backup slide, shown on request.)",
        "（Backupスライド。質疑応答で必要に応じて表示する。）",
        (("Measured on two graphs only; no selector rule is inferred for unmeasured graphs.",
          "測定は二グラフのみであり、未測定グラフへの選択規則は導かない。"),),
        (("Would the shared kernel ever win? It might on graphs not measured here; that "
          "was not tested, so no claim is made.",
          "sharedカーネルが勝つことはあるか。ここで測定していないグラフではあり得るが、未検証で"
          "あり主張しない。"),),
    ),
    Note(
        20, 40,
        "This backup slide breaks the runtime into BFS, backward accumulation, and other, "
        "for each graph. These come from complete b512 wall-clock runs, where other is the "
        "total minus BFS minus backward; they are not partial Nsight trace totals. Where "
        "Nsight profiling is referenced elsewhere, that is a single trace of a single graph "
        "and does not generalize to all graphs.",
        "このBackupスライドは、各グラフの実行時間をBFS、逆方向の累積、その他へ分解したものです。"
        "これらはb512での完全な実行の実時間に基づいており、その他は総時間からBFSと逆方向累積を"
        "引いた値です。部分的なNsightトレースの合計ではありません。他所でNsightプロファイルに"
        "言及している箇所については、単一グラフの単一トレースであり、全グラフへ一般化しません。",
        "(Backup slide, shown on request.)",
        "（Backupスライド。質疑応答で必要に応じて表示する。）",
        (("These are complete wall-clock runs, not partial Nsight trace totals.",
          "これは完全実行の実時間であり、部分的なNsightトレースの合計ではない。"),),
        (("Why not report Nsight totals directly? A single partial trace would not account "
          "for the whole run, so the wall-clock decomposition is reported instead.",
          "なぜNsightの合計を直接報告しないのか。単一の部分トレースでは実行全体を説明できないため、"
          "実時間による分解を報告している。"),),
    ),
    Note(
        21, 40,
        "This backup slide lists all thirteen correctness comparisons individually: the "
        "tier, the graph, the reference and candidate implementations, their batch sizes, "
        "the maximum relative error, and the tolerance result. Across all thirteen, missing "
        "indices are zero, mismatched elements are zero, and byte-identical is No. The three "
        "Tier A rows are the independent evidence; Tier B compares implementation paths and "
        "is not an independent ground-truth evaluation.",
        "このBackupスライドは、13件の正確性比較を個別に列挙したものです。Tier、グラフ、参照実装と"
        "候補実装、それぞれのバッチサイズ、最大相対誤差、許容判定を示しています。13件すべてにおいて、"
        "欠損インデックスは0、不一致要素は0、byte一致はNoです。Tier Aの3行が独立な根拠であり、"
        "Tier Bは実装経路どうしの比較であって、独立したground truthとの評価ではありません。",
        "(Backup slide, shown on request.)",
        "（Backupスライド。質疑応答で必要に応じて表示する。）",
        (("Tier B is implementation-path consistency, not comparison against independent ground truth.",
          "Tier Bは実装経路間の整合であり、独立したground truthとの比較ではない。"),),
        (("Why is Tier A limited to small graphs? The independent CPU reference is "
          "sequential, so full-vector comparison is only tractable at small scale.",
          "なぜTier Aは小規模グラフに限られるのか。独立参照がSequential CPU実装であり、"
          "ベクトル全体の比較は小規模でのみ現実的だからである。"),),
    ),
    Note(
        22, 40,
        "This backup slide is a historical record and is deliberately separated from the "
        "current results. An earlier version of the 325557 graph was malformed: it used "
        "1-based vertex identifiers, was missing seven elements, and contained out-of-range "
        "identifiers. On that input, correctness reported CORE_FAIL with a stress mismatch "
        "at the 1e-6 level. The graph was reconstructed by symmetry, a validator was added, "
        "and on the corrected input all thirteen comparisons pass. The old results are "
        "retained as evidence about invalid input, and they are not used in any current "
        "conclusion. The banner on this slide says so explicitly.",
        "このBackupスライドは履歴的な記録であり、現在の結果とは意図的に分離しています。325557グラフの"
        "旧版はmalformedでした。頂点IDが1始まりで、要素が7個欠落しており、範囲外のIDを含んでいました。"
        "この入力では正確性検証がCORE_FAILとなり、1e-6の水準でstress不一致が生じていました。その後、"
        "対称性からグラフを再構成し、検証器を追加した結果、修正版の入力では13比較すべてが合格して"
        "います。旧結果は不正な入力に関する記録として保存していますが、現在のいかなる結論にも用いて"
        "いません。このスライドのバナーはその点を明示しています。",
        "(Backup slide, shown on request.)",
        "（Backupスライド。質疑応答で必要に応じて表示する。）",
        (("Historical evidence only; it is not part of the current results.",
          "履歴的な記録にすぎず、現在の結果の一部ではない。"),),
        (("Why keep the malformed-input results at all? Detecting the defect and "
          "re-validating on the corrected input is itself part of the reproducibility record.",
          "malformed入力の結果をなぜ残すのか。不整合を検出し修正版で再検証した経緯自体が"
          "再現性の記録だからである。"),),
    ),
]

# Gate V1.2 uses concise bilingual scripts. Timing is intentionally left to the
# user, so the legacy V1.1 scripts above are replaced before any output is built.
def _v12_note(slide, en, ja, en_transition, ja_transition, limitation_en,
              limitation_ja, question_en, question_ja):
    return Note(slide, 0, en, ja, en_transition, ja_transition,
                ((limitation_en, limitation_ja),), ((question_en, question_ja),))


NOTES = [
    _v12_note(1,
        "The main result comes first. GPU_Opt achieved 1.31 to 3.17 times speedup on four evaluated graphs. The comparison used one GH200 GPU and median runtimes.",
        "最初に主結果を示します。GPU_Optは評価した4グラフで1.31倍から3.17倍の高速化を達成しました。比較はGH200一台と実行時間の中央値を用いました。",
        "I will outline the talk.", "発表の流れを示します。",
        "The result applies only to the evaluated conditions.", "結果は評価した条件にのみ適用されます。",
        "What is the comparator? PathMerge is a tuned third-party external GPU implementation.",
        "比較対象は何か。PathMergeは調整済みの第三者GPU実装です。"),
    _v12_note(2,
        "I first explain BC and Brandes. I then show the GPU design, evaluation, results, and limits.",
        "最初にBCとBrandesを説明します。その後、GPU設計、評価、結果、限界を示します。",
        "First, what does BC measure?", "まず、BCが何を測るかを説明します。",
        "The agenda does not fix the speaking time.", "このAgendaは発表時間を固定しません。",
        "Why start with results? The audience can see the outcome before the technical details.",
        "なぜ結果から始めるのか。技術詳細の前に成果を共有できるからです。"),
    _v12_note(3,
        "BC measures how often a vertex lies on shortest paths. In this path graph, C has BC four. C connects the two sides most often.",
        "BCは頂点が最短経路上に現れる頻度を表します。このパスグラフではCのBCが4で、両側を最も多く結びます。",
        "Next, I explain the exact BC algorithm.", "次に、厳密BCアルゴリズムを説明します。",
        "This graph is a simple teaching example.", "このグラフは説明用の簡単な例です。",
        "Why is C important? More shortest paths pass through C than through any other vertex.",
        "なぜCが重要か。他のどの頂点よりも多くの最短経路がCを通るからです。"),
    _v12_note(4,
        "Direct path enumeration is expensive. Brandes reuses shortest-path information from each BFS. For unweighted graphs, exact BC costs O of V times V plus E.",
        "経路を直接列挙すると高コストです。Brandesは各BFSの最短経路情報を再利用します。無重みグラフの厳密BCはO(|V|(|V|+|E|))です。",
        "The next slide separates Brandes from this study's contribution.", "次のスライドでBrandesと本研究の貢献を分けます。",
        "I do not improve Brandes' asymptotic complexity.", "本研究はBrandesの漸近計算量を改善しません。",
        "Does this study propose a new BC algorithm? No. It keeps Brandes' mathematical algorithm.",
        "本研究は新しいBCアルゴリズムを提案するのか。いいえ。Brandesの数理アルゴリズムを維持します。"),
    _v12_note(5,
        "Brandes defines the exact computation. I improve its GPU execution. I batch sources, map blocks, switch BFS direction, cooperate within warps, and use two streams.",
        "Brandesが厳密計算を定義します。本研究はGPU実行を改善します。始点のバッチ化、block割当、BFS方向切替、warp協調、2 streamを用います。",
        "I now show the integrated framework.", "次に統合実行基盤を示します。",
        "The BC definition and Brandes equations stay unchanged.", "BCの定義とBrandesの数式は変えません。",
        "What is new? The contribution is the integrated GPU execution framework and its evaluation.",
        "何が新しいのか。統合GPU実行基盤とその評価が貢献です。"),
    _v12_note(6,
        "The framework combines four execution components. Each component targets a different GPU bottleneck. I claim the integration and evaluation, not each component's first use.",
        "実行基盤は4つの実行要素を組み合わせます。各要素は異なるGPUボトルネックを対象にします。個別要素の初出ではなく統合と評価を主張します。",
        "Source batching also creates a memory constraint.", "始点バッチ化はメモリ制約も生みます。",
        "Individual component novelty is not claimed.", "個々の要素の新規性は主張しません。",
        "Why integrate the components? Their bottlenecks occur in one shared execution flow.",
        "なぜ統合するのか。各ボトルネックが同じ実行フロー内に現れるからです。"),
    _v12_note(7,
        "Source batching groups source vertices. It does not split the graph. At b512 and two effective streams, the code-derived allocation estimate is about 10.67 GB.",
        "始点バッチ化は始点頂点をまとめます。グラフを分割しません。b512と実効2 streamでは、コード由来の割当推定は約10.67 GBです。",
        "Three variants manage this working set differently.", "3つのvariantはこのworking setを異なる方法で管理します。",
        "The estimate is not a measured memory footprint.", "この推定値は実測メモリ使用量ではありません。",
        "Is source batching graph partitioning? No. The input graph stays fixed.",
        "始点バッチ化はgraph partitioningか。いいえ。入力グラフは固定です。"),
    _v12_note(8,
        "GPU_Opt uses Unified Memory. GPU_Opt_Pure uses device memory. GPU_Opt_Pure_Chunked splits the source batch. All variants share one execution framework.",
        "GPU_OptはUnified Memory、GPU_Opt_Pureはdevice memoryを使います。GPU_Opt_Pure_Chunkedは始点バッチを分割します。全variantは同じ実行基盤を共有します。",
        "I next separate the evaluation scopes.", "次に評価範囲を分けます。",
        "These are memory variants, not separate proposals.", "これらはメモリvariantであり、別々の提案ではありません。",
        "Does Chunked split the graph? No. It splits the source batch.",
        "Chunkedはグラフを分割するのか。いいえ。始点バッチを分割します。"),
    _v12_note(9,
        "I use four graphs for main performance. I use corrected 325557 for capacity tests. Every run uses one GH200 GPU.",
        "主性能評価には4グラフを使います。容量試験にはcorrected 325557を使います。すべての実行はGH200一台です。",
        "The first result is runtime.", "最初の結果は実行時間です。",
        "Input file size is not working-set size.", "入力ファイルサイズはworking setサイズではありません。",
        "Why separate the studies? Performance and capacity use different graph scopes.",
        "なぜ試験を分けるのか。性能と容量ではグラフ範囲が異なるからです。"),
    _v12_note(10,
        "GPU_Opt reduced median runtime on all four graphs. GPU_Opt uses fixed b512. PathMerge uses a tuned batch for each graph. The chart uses a log scale.",
        "GPU_Optは4グラフすべてで実行時間中央値を短縮しました。GPU_Optはb512固定です。PathMergeは各グラフで調整したバッチを使います。縦軸は対数です。",
        "The next slide shows the speedup ratios.", "次に高速化率を示します。",
        "Trial counts are small, so I report medians.", "試行数が少ないため中央値を報告します。",
        "Why use a log scale? Runtime spans a wide range across the graphs.",
        "なぜ対数軸を使うのか。グラフ間で実行時間の範囲が広いからです。"),
    _v12_note(11,
        "GPU_Opt achieved 3.17, 1.31, 1.51, and 1.45 times speedup. PathMerge is an external comparator, not ground truth.",
        "GPU_Optは3.17倍、1.31倍、1.51倍、1.45倍の高速化を達成しました。PathMergeは外部比較対象でありground truthではありません。",
        "I next examine the execution components.", "次に実行要素の効果を確認します。",
        "I do not generalize this result to PathMerge in general.", "この結果をPathMerge一般へは一般化しません。",
        "Were both batch sizes tuned? No. GPU_Opt stayed at b512; PathMerge was tuned per graph.",
        "両方のバッチを調整したのか。いいえ。GPU_Optはb512固定で、PathMergeは各グラフで調整しました。"),
    _v12_note(12,
        "On corrected 325557, H was 1.4767 times, W was 1.1012 times, and A was 1.5563 times. Hybrid BFS and dual streams had the largest observed effects.",
        "corrected 325557ではHが1.4767倍、Wが1.1012倍、Aが1.5563倍でした。Hybrid BFSとdual streamsが最大の観測効果を示しました。",
        "The next result concerns tested memory capacity.", "次は試験したメモリ容量の結果です。",
        "The aggregate mixes checkpoints. No roadNet factor analysis was run.", "aggregateはcheckpointが混在します。roadNetの要因分析は未実施です。",
        "Do these factors explain roadNet speedups? No. This analysis did not use roadNet graphs.",
        "これらの要因はroadNetの高速化を説明するのか。いいえ。この分析ではroadNetを使っていません。"),
    _v12_note(13,
        "Pure passed b4096 and failed at b8192 with CUDA OOM. UM passed b10240 and ended at b12288 with a host OOM kill. Chunked passed b16384.",
        "Pureはb4096に成功しb8192でCUDA OOMとなりました。UMはb10240に成功しb12288でhost OOM killとなりました。Chunkedはb16384に成功しました。",
        "These results have clear limits.", "これらの結果には明確な限界があります。",
        "A tested upper bound is not an unlimited capacity claim.", "試験上限は無制限の容量を意味しません。",
        "Are failures zero-second runtimes? No. They are categorical failure outcomes.",
        "失敗は0秒の実行時間か。いいえ。カテゴリとしての失敗結果です。"),
    _v12_note(14,
        "The evidence uses one GH200, four main graphs, one capacity graph, and small trial counts. I do not generalize beyond these conditions.",
        "根拠はGH200一台、主性能4グラフ、容量1グラフ、少数試行に基づきます。これらの条件を超えて一般化しません。",
        "Within these limits, I summarize four contributions.", "この限界の範囲で4つの貢献をまとめます。",
        "PathMerge is a retained third-party snapshot.", "PathMergeは保存された第三者snapshotです。",
        "Does the result cover other GPUs? No. Only one GH200 was evaluated.",
        "他のGPUも対象か。いいえ。評価したGPUはGH200一台です。"),
    _v12_note(15,
        "I contribute an integrated exact-BC GPU framework, an external comparison, component measurements, and a tested capacity comparison. Correctness remains required validation, not a main contribution.",
        "貢献は厳密BCの統合GPU基盤、外部比較、要素効果の測定、試験容量の比較です。正確性は必要な検証であり、主要貢献ではありません。",
        "I will close with the result and scope.", "最後に結果と適用範囲を述べます。",
        "No new BC algorithm is claimed.", "新しいBCアルゴリズムは主張しません。",
        "What is the central contribution? The integration and evaluation of GPU execution methods.",
        "中心的な貢献は何か。GPU実行方法の統合と評価です。"),
    _v12_note(16,
        "The framework achieved 1.31 to 3.17 times speedup. Hybrid BFS and dual streams had the largest observed effects. Memory variants changed the tested feasible batch range.",
        "実行基盤は1.31倍から3.17倍の高速化を達成しました。Hybrid BFSとdual streamsが最大の観測効果を示しました。メモリvariantは試験上の実行可能バッチ範囲を変えました。",
        "Questions are welcome. Backup evidence follows.", "質問を受けます。以降はBackup資料です。",
        "The conclusions apply only to the evaluated conditions.", "結論は評価した条件にのみ適用されます。",
        "What should be tested next? More GPUs, graphs, trials, and roadNet factor analysis.",
        "次に何を試験すべきか。GPU、グラフ、試行数の拡大とroadNetの要因分析です。"),
    _v12_note(17,
        "This backup lists the hardware, software, bandwidth records, and aggregation method. The main experiments use one GH200 GPU.",
        "このBackupはhardware、software、帯域記録、集計方法を示します。主実験はGH200一台を使います。",
        "The next backup lists graph and batch parameters.", "次のBackupはグラフとバッチの設定です。",
        "Undetermined items stay undetermined.", "未確定項目は未確定のままです。",
        "Why are two memory capacity values shown? They come from different records and units.",
        "なぜメモリ容量値が2つあるのか。異なる記録と単位に由来するからです。"),
    _v12_note(18,
        "This backup lists graph batches, medians, trial counts, and the working-set formula. The per-source state is code-derived, not measured.",
        "このBackupはグラフ別バッチ、中央値、試行数、working set式を示します。始点ごとの状態量はコード由来で、実測ではありません。",
        "The next backup shows PathMerge tuning.", "次のBackupはPathMergeの調整を示します。",
        "GPU_Opt stays at b512 on every main graph.", "GPU_Optは主性能の全グラフでb512固定です。",
        "Why not tune GPU_Opt? The fixed setting avoids a proposal-favoring comparison.",
        "なぜGPU_Optを調整しないのか。提案側に有利な比較を避けるためです。"),
    _v12_note(19,
        "This sweep selects the tested PathMerge batch for each graph. The historical malformed graph is excluded.",
        "このsweepは各グラフのPathMergeバッチを選びます。履歴的なmalformed graphは除外します。",
        "The next backup compares two kernels.", "次のBackupは2つのkernelを比較します。",
        "The sweep supports only the tested batch points.", "sweepが裏づけるのは試験点だけです。",
        "Could another batch be faster? It is possible outside the tested points.",
        "別のバッチがより速い可能性はあるか。試験点の外ではあり得ます。"),
    _v12_note(20,
        "This backup forces block and shared kernels on PA and TX. The block kernel was faster on both tested graphs.",
        "このBackupはPAとTXでblock kernelとshared kernelを強制比較します。試験した両グラフでblock kernelが高速でした。",
        "The next backup shows phase timing.", "次のBackupはphase timingを示します。",
        "No selector rule is inferred for unmeasured graphs.", "未測定グラフのselector ruleは導きません。",
        "Can shared memory win elsewhere? It was not tested on other graphs.",
        "他のグラフでshared memoryが勝つか。他のグラフでは試験していません。"),
    _v12_note(21,
        "This backup separates BFS, backward accumulation, and other time. The values come from complete b512 wall-clock runs.",
        "このBackupはBFS、backward accumulation、otherの時間を分けます。値は完全なb512 wall-clock runから得ています。",
        "The next backup gives detailed correctness validation.", "次のBackupは詳細な正確性検証です。",
        "The phase values are not partial Nsight totals.", "phase値は部分的なNsight合計ではありません。",
        "Why use wall-clock data? A partial trace does not cover the complete run.",
        "なぜwall-clock dataを使うのか。部分traceは完全な実行を覆わないからです。"),
    _v12_note(22,
        "Correctness is required validation, not the main research result. All 13 comparisons passed mixed tolerance. Tier A has three independent CPU comparisons. Tier B has ten path comparisons.",
        "正確性は必要な検証であり、主要な研究成果ではありません。13比較はすべてmixed toleranceに合格しました。Tier Aは独立CPU比較3件、Tier Bは実装経路比較10件です。",
        "The final backup preserves the malformed-input history.", "最後のBackupはmalformed inputの履歴を保存します。",
        "Tier B is not an independent ground-truth evaluation.", "Tier Bは独立したground truth評価ではありません。",
        "Was correctness checked? Yes. The detailed table retains all validation data.",
        "正確性は確認したのか。はい。詳細表にすべての検証データを保持しています。"),
    _v12_note(23,
        "This backup preserves the malformed-input record. The old input failed validation. The corrected input passed all 13 comparisons. Current conclusions use only the corrected input.",
        "このBackupはmalformed inputの記録を保存します。旧入力は検証に失敗しました。修正入力は13比較すべてに合格しました。現在の結論は修正入力だけを使います。",
        "This ends the backup material.", "Backup資料は以上です。",
        "Historical evidence is not part of the current results.", "履歴的根拠は現在の結果の一部ではありません。",
        "Why keep the old result? It documents input validation and correction.",
        "なぜ旧結果を残すのか。入力検証と修正の経緯を記録するためです。"),
]

MAIN_SLIDES = 16
BACKUP_SLIDES = 7

# Canonical one-line titles, index 0 = slide 1. ``slide_frame`` may break a title
# across two lines to fit the band, but never changes its wording.
SLIDE_TITLES = [
    THESIS_TITLE_FLAT,
    "Agenda",
    "Betweenness Centrality Finds Important Bridge Vertices",
    "Brandes' Algorithm Reduces the Cost of Exact BC",
    "I Improve How Brandes Runs on the GPU",
    "The Proposal Is an Integrated GPU Execution Framework",
    "Source Batching Creates a Batch-Dependent Working Set",
    "Three Memory-Management Variants Share One Framework",
    "I Separate Performance and Capacity Studies",
    "GPU_Opt Reduced Runtime on All Four Evaluated Graphs",
    "GPU_Opt Achieved 1.31–3.17× Speedup",
    "Hybrid BFS and Dual Streams Gave the Largest Observed Effects",
    "Memory Variants Expanded the Tested Batch Range",
    "The Results Have Clear Limits",
    "Contributions",
    "The Framework Improved Performance and Expanded the Tested Batch Range",
    "Detailed Experimental Environment",
    "Graph and Batch Parameters",
    "PathMerge Batch-Size Sweep",
    "Forced Block-vs-Shared Kernel Comparison",
    "Phase Breakdown and Profiling Scope",
    "Detailed Correctness Evidence",
    "Historical Record of the Malformed Input",
]


# --- Spoken-length estimation ----------------------------------------------
# Both scripts must be independently deliverable in the slide's budget, so each
# is estimated in its own units: English at a presentation pace in words per
# second, Japanese in characters per second. The bands are wide because these
# are planning aids, not measurements; they exist to catch a script that is a
# stub or that has drifted to double its slot.
EN_WORDS_PER_SEC = 2.5
JA_CHARS_PER_SEC = 5.5
PACE_BAND = (0.55, 1.60)


def en_seconds(text: str) -> float:
    return len(text.split()) / EN_WORDS_PER_SEC


def ja_seconds(text: str) -> float:
    return len("".join(text.split())) / JA_CHARS_PER_SEC


def notes_pacing() -> dict:
    """Estimated spoken length of each script, plus the main-deck totals."""
    rows = []
    for n in NOTES:
        rows.append({
            "Slide": n.slide,
            "Declared": n.seconds,
            "EnglishEst": round(en_seconds(n.en_script), 1),
            "JapaneseEst": round(ja_seconds(n.ja_script), 1),
        })
    main = [r for r in rows if r["Slide"] <= MAIN_SLIDES]
    return {
        "rows": rows,
        "declared_total": sum(r["Declared"] for r in main),
        "english_total": round(sum(r["EnglishEst"] for r in main)),
        "japanese_total": round(sum(r["JapaneseEst"] for r in main)),
        "backup_declared_total": sum(r["Declared"] for r in rows if r["Slide"] > MAIN_SLIDES),
    }


def validate_notes_content() -> None:
    """Both scripts exist and are each usable without a fixed time budget."""
    assert len(NOTES) == MAIN_SLIDES + BACKUP_SLIDES, len(NOTES)
    assert [n.slide for n in NOTES] == list(range(1, 24))
    for n in NOTES:
        assert n.en_script.strip() and n.ja_script.strip(), f"slide {n.slide}: empty script"
        assert not CJK_RE.search(n.en_script), f"slide {n.slide}: Japanese inside the English script"
        assert CJK_RE.search(n.ja_script), f"slide {n.slide}: Japanese script has no Japanese"
        assert re.search(r"[A-Za-z]", n.en_script), f"slide {n.slide}: English script has no Latin text"
        assert n.limitations and n.questions, f"slide {n.slide}: missing limitation or question"
        for en_text, ja_text in n.limitations + n.questions:
            assert not CJK_RE.search(en_text), f"slide {n.slide}: Japanese in an English entry"
            assert CJK_RE.search(ja_text), f"slide {n.slide}: Japanese entry has no Japanese"


def attach_notes(prs):
    """Embed both scripts in the notes pane under stable bilingual headings."""
    for n in NOTES:
        tf = prs.slides[n.slide - 1].notes_slide.notes_text_frame
        lines = [
            f"Slide {n.slide} — {SLIDE_TITLES[n.slide - 1]}",
            "Duration: [USER TO EDIT]",
            "",
            "[English Script]",
            n.en_script,
            "",
            "[日本語説明]",
            n.ja_script,
            "",
            "[English Transition]",
            n.en_transition,
            "",
            "[日本語トランジション]",
            n.ja_transition,
            "",
            "[Limitations / 限定]",
        ]
        for en_text, ja_text in n.limitations:
            lines += [f"- {en_text}", f"- {ja_text}"]
        lines += ["", "[Expected Questions / 想定質問]"]
        for en_text, ja_text in n.questions:
            lines += [f"- {en_text}", f"- {ja_text}"]
        tf.text = "\n".join(lines)


SLIDE_FIGURES = {
    3: "native BC graph", 6: "F04", 7: "F05", 8: "F08", 10: "F09",
    11: "F10", 12: "F12", 13: "F13", 19: "F11", 20: "F14", 21: "F15",
}

SLIDE_PURPOSE = {
    1: "Lead with the measured result and scope",
    2: "Preview the results-first narrative",
    3: "Explain BC with an editable five-node graph",
    4: "Explain how Brandes reduces exact-BC cost",
    5: "Separate Brandes from this study's GPU contribution",
    6: "Present the integrated GPU execution framework",
    7: "Explain the batch-dependent working set",
    8: "Show three memory variants of one framework",
    9: "Separate main performance from capacity tests",
    10: "Report runtime reduction on four graphs",
    11: "Report speedup over the tuned comparator",
    12: "Report observed component effects",
    13: "Report tested memory-capacity outcomes",
    14: "State limits without correctness as a main result",
    15: "Summarize four research contributions",
    16: "Close with the result, capacity range, and scope",
    17: "Backup: full hardware and software environment",
    18: "Backup: graph, batch, and working-set parameters",
    19: "Backup: PathMerge batch sweep justifying tuning",
    20: "Backup: forced block-vs-shared kernel comparison",
    21: "Backup: phase breakdown and profiling scope",
    22: "Backup: required correctness validation detail",
    23: "Backup: historical malformed-input evidence",
}


# --- Documents --------------------------------------------------------------
def write_manifest(stats):
    header = ["SlideNumber", "Section", "SlideTitle", "NarrativePurpose", "FigureID",
              "ObjectType", "Editable", "PlannedSeconds", "CanonicalSource"]
    kinds = {st["Slide"]: st for st in stats["slide_stats"]}
    sources = {
        10: "result/tables/thesis/T2_main_performance.tsv;raw_data/main_performance/proposed_variants/*",
        11: "result/main_performance/proposed_vs_pathmerge/comparison.tsv;raw_data/main_performance/proposed_variants/*",
        12: "result/tables/thesis/T3_ablation_summary.tsv;raw_data/ablation/*;raw_data/corrected_325557/*",
        13: "result/tables/thesis/T4_memory_scalability.tsv;result/memory_scalability/corrected_325557/feasibility_boundary.tsv",
        9: "result/tables/thesis/T1_graph_metadata.tsv;result/datasets/graph_catalog.tsv",
        17: "result/tables/thesis/T6_experimental_environment.tsv",
        18: "result/tables/thesis/T2_main_performance.tsv;docs/thesis/writing/japanese/appendix_a_experimental_parameters.md",
        19: "raw_data/tuning/pathmerge/*/pathmerge_bc/job_multi_20260710/pathmerge_sweep_results.tsv",
        20: "raw_data/tuning/kernel_selection/*",
        21: "raw_data/main_performance/proposed_variants/*/_run/job_2357334_20260711/phase_timing.log",
        22: "result/tables/thesis/T5_correctness_summary.tsv",
        23: "result/provenance/GRAPH_325557_INTEGRITY_AUDIT.md;result/datasets/graph_catalog.tsv",
    }
    with MANIFEST_PATH.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f, delimiter="\t", lineterminator="\n")
        w.writerow(header)
        for n in range(1, 24):
            st = kinds[n]
            kind = []
            if st["Charts"]:
                kind.append("native chart")
            if st["Tables"]:
                kind.append("native table")
            if st["Shapes"] > st["Charts"] + st["Tables"]:
                kind.append("native shapes and text")
            w.writerow([
                n,
                "Main" if n <= MAIN_SLIDES else "Backup",
                SLIDE_TITLES[n - 1],
                SLIDE_PURPOSE[n],
                SLIDE_FIGURES.get(n, "NA"),
                "; ".join(kind),
                "yes",
                "[USER TO EDIT]",
                sources.get(n, "docs/thesis/writing/japanese/*"),
            ])


def write_notes_document():
    lines = [
        "# Speaker Notes (Bilingual)", "",
        "Presentation timing is intentionally not fixed.",
        "The user will adjust slide selection and speaking time after rehearsal.", "",
        "スライド面はすべて英語である。英語原稿と日本語説明は同じ主張と限定を述べる。", "",
        "## Main", "",
    ]
    for n in NOTES:
        if n.slide == MAIN_SLIDES + 1:
            lines += ["## Backup", "", "Backup は質疑応答用である。", ""]
        lines += [
            f"## Slide {n.slide} — {SLIDE_TITLES[n.slide - 1]}", "",
            "### Duration", "", "[USER TO EDIT]", "",
            "### English Script", "", n.en_script, "",
            "### 日本語説明", "", n.ja_script, "",
            "### English Transition", "", n.en_transition, "",
            "### 日本語トランジション", "", n.ja_transition, "",
            "### Limitations to State", "",
        ]
        for en_text, ja_text in n.limitations:
            lines += [f"- {en_text}", f"- {ja_text}"]
        lines += ["", "### Expected Questions", ""]
        for en_text, ja_text in n.questions:
            lines += [f"- {en_text}", f"- {ja_text}"]
        lines += [""]
    NOTES_PATH.write_text("\n".join(lines), encoding="utf-8")
    return None

    # Legacy V1.1 timed-note writer retained below for provenance only.
    pacing = notes_pacing()
    total = pacing["declared_total"]
    lines = [
        "# Speaker Notes (Bilingual)",
        "",
        "スライド面はすべて英語である。日本語はこのノートにのみ存在する。",
        "",
        "各スライドは **完全な英語スクリプト** と **完全な日本語説明** の両方を持つ。"
        "この二つは一本の読み上げ原稿の前半・後半ではなく、発表言語に応じて **どちらか一方だけを"
        "読む** ための代替である。両方を続けて読むことは想定していない。",
        "",
        f"本編 {MAIN_SLIDES} 枚の想定時間合計は {total} 秒（{total / 60:.1f} 分）である。"
        f"英語版のみを読んだ場合の推定合計は約 {pacing['english_total']} 秒、"
        f"日本語版のみを読んだ場合の推定合計は約 {pacing['japanese_total']} 秒であり、"
        "いずれか一方だけで発表が成立する。",
        "",
        f"推定は英語 {EN_WORDS_PER_SEC} words/秒、日本語 {JA_CHARS_PER_SEC} 文字/秒の"
        "発表ペースによる計画値であり、実測ではない。",
        "",
        "発表時間は固定せず、リハーサル後にユーザーが調整する。",
        "",
        "同じ内容は PPTX のノートペインにも埋め込まれており、スライド面には表示されない。",
        "",
        "## Main",
        "",
    ]
    for n in NOTES:
        if n.slide == MAIN_SLIDES + 1:
            lines += ["## Backup", "",
                      "Backup は本編では使用せず、質疑応答時にのみ参照する。"
                      f"想定時間の合計 {pacing['backup_declared_total']} 秒は本編 {total} 秒には含まれない。",
                      ""]
        lines += [
            f"## Slide {n.slide} — {SLIDE_TITLES[n.slide - 1]}",
            "",
            "### Duration",
            "",
            f"{n.seconds} seconds",
            "",
            "### English Script",
            "",
            n.en_script,
            "",
            "### 日本語説明",
            "",
            n.ja_script,
            "",
            "### English Transition",
            "",
            n.en_transition,
            "",
            "### 日本語トランジション",
            "",
            n.ja_transition,
            "",
            "### Limitations to State",
            "",
        ]
        for en_text, ja_text in n.limitations:
            lines += [f"- {en_text}", f"- {ja_text}"]
        lines += ["", "### Expected Questions", ""]
        for en_text, ja_text in n.questions:
            lines += [f"- {en_text}", f"- {ja_text}"]
        lines += [""]
    NOTES_PATH.write_text("\n".join(lines), encoding="utf-8")
    return total


def write_plan_document(total_seconds):
    rows = "\n".join(
        f"| {n} | {'Main' if n <= MAIN_SLIDES else 'Backup'} | {SLIDE_TITLES[n - 1]} | "
        f"{SLIDE_PURPOSE[n]} | {SLIDE_FIGURES.get(n, '—')} |"
        for n in range(1, 24))
    PLAN_PATH.write_text(f"""# Presentation Plan — Gate V1.2.1

## Narrative

> Results → Agenda → BC → Brandes → GPU Execution → Evaluation → Results → Limits → Contributions → Conclusion

The main result appears on Slide 1: GPU_Opt achieved 1.31–3.17× speedup on four evaluated graphs.

Presentation timing is intentionally not fixed.
The user will adjust slide selection and speaking time after rehearsal.

## Slide map

| # | Section | Title | Narrative purpose | Figure |
|---:|---|---|---|---|
{rows}

## Claim boundaries

- Brandes reduces the algorithmic cost of exact BC.
- This study improves how Brandes runs on the GPU.
- This study does not change Brandes' mathematical algorithm or asymptotic complexity.
- Source batching groups sources. It does not split the graph.
- PathMerge is a third-party external comparator, not ground truth.
- Correctness is required validation, not a main research result.
- Detailed correctness and malformed-input evidence remain in Backup Slides 22–23.
- All visible slide text is English. Speaker notes contain English and Japanese.
""", encoding="utf-8")
    return

    # Legacy V1.1 timed plan retained below for provenance only.
    pacing = notes_pacing()
    rows = "\n".join(
        f"| {n} | {'Main' if n <= MAIN_SLIDES else 'Backup'} | {SLIDE_TITLES[n - 1]} "
        f"| {SLIDE_PURPOSE[n]} | {SLIDE_FIGURES.get(n, '—')} |"
        for n in range(1, 23))
    PLAN_PATH.write_text(f"""# Presentation Plan

## 1. Narrative

本発表は次の流れで構成する。

> Problem → Proposed Framework → Evaluation → Evidence → Limitations → Conclusion

聴衆が発表終了時に理解すべき中心メッセージは次の1文である。

> 固定b512のblock-based GPU_Optは、評価した4グラフにおいて、グラフごとに調整した第三者実装PathMergeより1.31〜3.17倍高速だった。性能向上は統合GPU実行基盤の複数要素から生じ、UM・Pure・Chunkedの比較では性能だけでなくbatch-dependent working setに対する容量特性と数値整合性も明らかになった。

## 2. Time budget

本編 {MAIN_SLIDES} スライド、Backup {BACKUP_SLIDES} スライド。本編の想定時間合計は {total_seconds} 秒（{total_seconds / 60:.1f} 分）である。

英語版のみを読んだ場合の推定合計は約 {pacing['english_total']} 秒、日本語版のみを読んだ場合の推定合計は約 {pacing['japanese_total']} 秒である。二つのスクリプトは代替であり、合算して読むことは想定していない。

発表時間は固定しない。リハーサル後にユーザーがスライド選択と発表時間を調整する。

## 3. Slide map

| # | Section | Title | Narrative purpose | Figure |
|---:|---|---|---|---|
{rows}

## 4. Language and design rules

- **スライド面に表示される文字はすべて英語**とする。タイトル、本文、bullet、callout、caption、footnote、表ヘッダ、表セル、chart のタイトル・軸・凡例・データラベル、diagram のノード、注釈、Backup スライドを含む。
- 日本語が存在してよいのは speaker notes のみである。スライド面に仮名・漢字を残さない。
- speaker notes は各スライドに **完全な英語スクリプト** と **完全な日本語説明** の両方を持つ。二つは代替であり、一方だけで発表が成立する。
- 表示書体は Arial に統一する（Yu Gothic 依存はスライド面から除去済み）。
- 16:9、白背景。配色は編集可能図ライブラリ（`docs/thesis/figures/editable/`）から継承する。
- deck title 30 pt 以上、slide title {TITLE_PT} pt 以上、本文 {BODY_PT} pt 以上、図表内文字 {FIGURE_MIN_PT} pt 以上、footnote {FOOTNOTE_PT} pt 以上（footnote の使用は最小限）。
- 英語化で行が長くなる場合は文を短くする。font size を下げて収めない。
- 1 スライド 1 メッセージ、bullet は 5 項目以内。
- 結果スライドのタイトルは値の羅列ではなく「何が分かったか」を述べる。

## 5. Claim boundaries enforced in this deck

- PathMerge は評価用に保存した第三者実装の external comparator であり、原著者公式実装でも ground truth でもない。
- graph file size と batch-dependent working set を明確に分離する。
- source batching は graph partitioning ではない。
- UM・Chunked が無制限に OOM を回避するとは記載しない。試験上限を明示する。
- 旧 malformed 入力の結果は Backup B7 に履歴的記録として分離し、現在の結論に混入させない。
- 評価していない GPU・グラフ・PathMerge 一般へ一般化しない。
""", encoding="utf-8")


def write_readme(stats, total_seconds):
    README_PATH.write_text(f"""# Master's Thesis Presentation (Gate V1.2.1)

Results-first, fully editable PowerPoint deck for the master's thesis presentation.

## Files

- `editable/master_thesis_presentation_v1.pptx`: {MAIN_SLIDES} main slides and {BACKUP_SLIDES} backup slides.
- `presentation_plan.md`: results-first narrative and slide map.
- `speaker_notes_bilingual.md`: English and Japanese notes for all 23 slides.
- `PRESENTATION_MANIFEST.tsv`: slide relationships, object types, and sources.
- `../../../scripts/generate_thesis_presentation.py`: deterministic generator and validator.

## Presentation timing

Presentation timing is intentionally not fixed.
The user will adjust slide selection and speaking time after rehearsal.

Every notes entry uses `Duration: [USER TO EDIT]` in the PowerPoint notes pane.

## Correctness role

Correctness is treated as required validation. It is not presented as the main research result.
Detailed correctness evidence and the historical malformed input remain in Backup Slides 22–23.
The Tier A/B data and all 13 comparisons are unchanged.

## Editability

- Editable objects: {stats['objects']}
- Native charts: {stats['charts']}
- Native tables: {stats['tables']}
- Raster pictures: {stats['pictures']}
- Raster-only slides: {stats['raster_only_slides']}

All diagrams use native PowerPoint shapes and connectors. Charts retain embedded workbooks.

## Regeneration

```bash
cd thesis_bc_project
PYTHONPATH=<deps> python3 scripts/generate_thesis_presentation.py
```

The generator checks canonical measured values, slide language, font floors, text fit,
object bounds, editable objects, notes, chart identifiers, and raster use.
""", encoding="utf-8")
    return

    # Legacy V1.1 README writer retained below for provenance only.
    pacing = notes_pacing()
    README_PATH.write_text(f"""# Master's Thesis Presentation (Gate V1.1 draft)

修士論文発表用スライドである。**スライド面に表示される文字はすべて英語**であり、
日本語は speaker notes にのみ存在する。正本は論文本文（`../writing/japanese/`）と
編集可能図ライブラリ（`../figures/editable/`）であり、本ディレクトリはそれらから
決定的に生成される。

## Files

- `editable/master_thesis_presentation_v1.pptx`: 16:9、白背景の発表スライド。本編 {MAIN_SLIDES} 枚 + Backup {BACKUP_SLIDES} 枚。表示文字はすべて英語。
- `presentation_plan.md`: narrative、時間配分、slide map、言語・デザイン規則、主張の限定。
- `speaker_notes_bilingual.md`: スライドごとの想定時間・英語スクリプト・日本語説明・transition・限定・想定質問。
- `PRESENTATION_MANIFEST.tsv`: スライド番号、section、タイトル、narrative purpose、図番号、オブジェクト種別、canonical source。
- `../../../scripts/generate_thesis_presentation.py`: 再生成器。

## Presentation length

発表時間は固定しない。リハーサル後にユーザーがスライド選択と発表時間を調整する。

## Bilingual speaker notes

各スライドは完全な英語スクリプトと完全な日本語説明の両方を持つ。二つは**代替**であり、
発表言語に応じてどちらか一方だけを読む。両方を続けて読むことは想定していない。

- 英語版のみを読んだ場合の本編推定合計: 約 {pacing['english_total']} 秒
- 日本語版のみを読んだ場合の本編推定合計: 約 {pacing['japanese_total']} 秒
- Backup {BACKUP_SLIDES} 枚（合計 {pacing['backup_declared_total']} 秒）は本編合計に含まない

推定は英語 {EN_WORDS_PER_SEC} words/秒、日本語 {JA_CHARS_PER_SEC} 文字/秒による計画値であり、実測ではない。

## Regeneration

```bash
cd thesis_bc_project
PYTHONPATH=<deps> python3 scripts/generate_thesis_presentation.py
```

必要 library は `python-pptx`, `openpyxl`, `lxml`。生成器は
`scripts/generate_editable_figure_library.py` の helper・palette・data loader・
package normalizer を import して再利用するため、図ライブラリと配色・書体・数値が
乖離しない。canonical 値の照合に失敗した場合は生成を停止する。

PPTX は zip entry timestamp と `docProps/core.xml` の日時を固定値へ正規化するため、
同一 interpreter での再生成は byte-identical になる。

## Editability

- raster-only slide: {stats['raster_only_slides']}
- 埋め込み raster 画像: {stats['pictures']}
- native chart: {stats['charts']}（embedded workbook あり）
- native table: {stats['tables']}
- 編集可能オブジェクト総数: {stats['objects']}

すべての図は native shape / connector / text box / chart / table であり、
PowerPoint 上で個別に選択・移動・再着色・文字編集できる。chart は右クリック
**Edit Data** で embedded workbook を編集できる。図を PNG として貼り付けた箇所はない。

## Self-validation

生成器は出力後に次を検査し、違反があれば停止する。

- スライド面（`ppt/slides`・`ppt/charts`・`ppt/diagrams`・`ppt/embeddings`）に仮名・漢字が 1 文字も存在しないこと。notes slide は検査対象外。
- 表示文字が ASCII、または明示的に許可した約物（{"".join(sorted(VISIBLE_NON_ASCII))}）のみであること。
- 全 23 notes slide に `[English Script]` と `[日本語説明]` の双方があり、英語側に日本語が混入せず、日本語側に日本語が存在すること。
- 各スクリプトが宣言した想定秒数に対して妥当な長さであること（一方が他方の要約になっていないこと）。
- 全 shape の各行が、指定 font size で box 内幅に収まること（途中改行の防止）。
- 表示文字（表セルを含む）が全スライドで {FIGURE_MIN_PT} pt 以上であること。
- 主張の限定文（PathMerge・working set・容量上限・Tier A/B・適用範囲）が字句どおり残存すること。
- 表セルの文字が列幅に収まり、行の自動伸長による枠外はみ出しを起こさないこと。
- slide 境界外の object が存在しないこと。
- raster 画像が存在しないこと。
- chart の `c:axId` / `c:crossAx` が正の整数で、chart 間で衝突しないこと。
- 主要値（speedup、tuned batch、H/W/A、容量境界、correctness tier 件数）が canonical source と一致すること。

## Human review still required

本環境に PowerPoint / LibreOffice renderer は存在しないため、OOXML・幾何・text-fit
検査のみを実施しており、**実際のレンダリング目視は未実施**である。実機 PowerPoint で
次を確認すること。

- Slide 1 の `[TO BE FILLED]`（Name・Affiliation・Supervisor・Date）の記入。
- chart の data label と凡例の実描画（Slide 8, 9, 10, 18, 19, 20）。
- 2 行に分割した長いタイトル（Slide 9, 12, 15）の折り返し位置と title rule との間隔。
- Slide 11 の feasibility marker と注記の重なり。
- Slide 12 および Backup（Slide 16-22）の表の可読性。
- notes pane に英語・日本語の双方が正しく表示されること。
""", encoding="utf-8")


# --- Validation -------------------------------------------------------------
def iter_table_texts(shape):
    if not shape.has_table:
        return
    for row in shape.table.rows:
        for cell in row.cells:
            yield cell.text


def check_visible(problems, slide_no, where, text):
    """Audience-visible text is English: no CJK, and no stray non-ASCII."""
    if CJK_RE.search(text):
        problems["visible_japanese"].append((slide_no, where, "".join(sorted(set(CJK_RE.findall(text))))))
    stray = {c for c in text if ord(c) > 127 and c not in VISIBLE_NON_ASCII}
    if stray:
        problems["visible_non_ascii"].append((slide_no, where, "".join(sorted(stray))))


def font_floor(shape_name: str) -> float:
    """Minimum point size for a shape, by the role its name encodes."""
    if shape_name == "ThesisTitle":
        return 30
    if shape_name == "Title":
        return TITLE_PT
    if "Bullet_" in shape_name:
        return FIGURE_MIN_PT
    # Captions, notes, and figure/table interiors. The deck floor already sits
    # above the 14 pt footnote minimum, so it is applied uniformly here.
    return FIGURE_MIN_PT


def scan_package_language(zf, problems):
    """No Japanese anywhere the audience can see; notes slides are excluded.

    The shape pass cannot see chart categories, series names, axis titles, or the
    workbook a chart carries for Edit Data, so the saved package is re-read and
    every visible part is scanned as raw XML.
    """
    visible = sorted(n for n in zf.namelist()
                     if (n.startswith(("ppt/slides/", "ppt/charts/", "ppt/diagrams/"))
                         and n.endswith(".xml"))
                     or (n.startswith("ppt/embeddings/") and n.endswith(".xlsx")))
    assert not any(n.startswith("ppt/notesSlides/") for n in visible), visible
    for name in visible:
        blob = zf.read(name)
        if name.endswith(".xlsx"):
            # An embedded workbook keeps its strings inside its own zip.
            # An embedded workbook keeps its cell text inside its own zip, and
            # only the content parts are audience-visible. Every workbook
            # openpyxl writes also carries the stock Office theme, whose
            # East-Asian font *fallback declarations* ("MS PGothic" and friends)
            # are typeface names in attributes that never render as text. Rather
            # than scan past them, they are asserted to be the only residue.
            with zipfile.ZipFile(io.BytesIO(blob)) as book:
                text = "\n".join(
                    book.read(inner).decode("utf-8", "ignore") for inner in book.namelist()
                    if inner == "xl/sharedStrings.xml" or inner.startswith("xl/worksheets/"))
                for inner in book.namelist():
                    if not inner.endswith(".xml"):
                        continue
                    stripped = re.sub(r'typeface="[^"]*"', "", book.read(inner).decode("utf-8", "ignore"))
                    assert not CJK_RE.search(stripped), f"{name}/{inner}: CJK outside a typeface attribute"
        else:
            text = blob.decode("utf-8")
        found = set(CJK_RE.findall(text))
        if found:
            problems["visible_japanese"].append((name, "package scan", "".join(sorted(found))))
    return visible


def check_required_sentences(prs):
    """The mandated claim-boundary wordings must survive verbatim on their slide."""
    for slide_no, sentences in REQUIRED_SENTENCES.items():
        blob = []
        for shape in prs.slides[slide_no - 1].shapes:
            if getattr(shape, "has_table", False):
                blob += [c.text for row in shape.table.rows for c in row.cells]
            elif getattr(shape, "has_text_frame", False):
                blob.append(shape.text)
        # Line breaks exist only to fit a shape and never change the wording.
        flat = " ".join(" ".join(t.split()) for t in blob)
        for sentence in sentences:
            assert " ".join(sentence.split()) in flat, f"slide {slide_no}: missing verbatim: {sentence}"


def check_notes_bilingual(prs):
    """Every slide's notes carry both complete scripts under explicit headings."""
    report = []
    for idx, slide in enumerate(prs.slides, 1):
        text = slide.notes_slide.notes_text_frame.text
        title = SLIDE_TITLES[idx - 1]
        assert text.count("[English Script]") == 1, f"slide {idx}: English Script heading"
        assert text.count("[日本語説明]") == 1, f"slide {idx}: Japanese heading"
        english = text.split("[English Script]", 1)[1].split("[日本語説明]", 1)[0].strip()
        japanese = text.split("[日本語説明]", 1)[1].split("[English Transition]", 1)[0].strip()
        assert f"Slide {idx} — {title}" in text, f"slide {idx}: notes title mismatch"
        assert "Duration: [USER TO EDIT]" in text, f"slide {idx}: editable duration missing"
        assert english and japanese, f"slide {idx}: empty script"
        assert re.search(r"[A-Za-z]", english), f"slide {idx}: English script has no Latin text"
        assert not CJK_RE.search(english), f"slide {idx}: Japanese inside the English script"
        assert CJK_RE.search(japanese), f"slide {idx}: Japanese script has no Japanese"
        for heading in ("[English Transition]", "[日本語トランジション]",
                        "[Limitations / 限定]", "[Expected Questions / 想定質問]"):
            assert heading in text, f"slide {idx}: missing {heading}"
        report.append({"Slide": idx, "EnglishChars": len(english), "JapaneseChars": len(japanese)})
    return report


def validate(prs_path, expect_slides, values):
    prs = Presentation(prs_path)
    assert len(prs.slides) == expect_slides, len(prs.slides)
    problems = {"visible_japanese": [], "visible_non_ascii": [], "bounds": [], "font": [],
                "overflow": [], "raster": [], "duplicate_names": []}
    slide_stats = []
    totals = {"objects": 0, "charts": 0, "tables": 0, "pictures": 0, "text": 0}
    for sn, slide in enumerate(prs.slides, 1):
        charts = tables = pictures = texts = 0
        for shape in slide.shapes:
            totals["objects"] += 1
            if shape.left < 0 or shape.top < 0 or \
               shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
                problems["bounds"].append((sn, shape.name))
            if getattr(shape, "has_chart", False):
                charts += 1
            if getattr(shape, "has_table", False):
                tables += 1
                tbl = shape.table
                for r_i, row in enumerate(tbl.rows):
                    for c_i, cell in enumerate(row.cells):
                        text = cell.text
                        check_visible(problems, sn, f"{shape.name} r{r_i}c{c_i}", text)
                        runs = [x for p in cell.text_frame.paragraphs for x in p.runs if x.text]
                        if not runs:
                            continue
                        size = max(x.font.size.pt for x in runs if x.font.size)
                        bold = any(x.font.bold for x in runs)
                        if size < FIGURE_MIN_PT:
                            problems["font"].append((sn, shape.name, f"r{r_i}c{c_i}", size, FIGURE_MIN_PT))
                        # 0.14 in of left+right cell margin is set in table().
                        avail = tbl.columns[c_i].width / 914400 - 0.14
                        need = text_block_in(text, size, bold)[0] - TEXT_INSET_IN
                        if need > avail + 1e-6:
                            problems["overflow"].append(
                                (sn, shape.name, f"cell r{r_i}c{c_i}", round(need, 3), round(avail, 3), text[:40]))
            if shape.shape_type == 13:
                pictures += 1
                problems["raster"].append((sn, shape.name))
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                texts += 1
                check_visible(problems, sn, shape.name, shape.text)
                runs = [r for p in shape.text_frame.paragraphs for r in p.runs if r.text]
                if runs:
                    size = max(r.font.size.pt for r in runs if r.font.size)
                    bold = any(r.font.bold for r in runs)
                    floor = font_floor(shape.name)
                    if size < floor:
                        problems["font"].append((sn, shape.name, size, floor))
                    need_w, need_h = text_block_in(shape.text, size, bold)
                    have_w = shape.width / 914400
                    have_h = shape.height / 914400
                    if need_w > have_w + 1e-6:
                        problems["overflow"].append((sn, shape.name, "width", round(need_w, 3), round(have_w, 3), shape.text[:60]))
                    if need_h > have_h + 1e-6:
                        problems["overflow"].append((sn, shape.name, "height", round(need_h, 3), round(have_h, 3), shape.text[:60]))
        assert texts > 0, f"slide {sn}: no native text"
        assert len(slide.shapes) - pictures > 1, f"slide {sn}: not a multi-object editable slide"
        # Duplicate names make a shape ambiguous to select in the PowerPoint
        # selection pane, which defeats the point of shipping an editable deck.
        named = [sh.name for sh in slide.shapes]
        problems["duplicate_names"] += [(sn, n) for n in set(named)
                                        if named.count(n) > 1 and not n.startswith("Straight Connector")]
        slide_stats.append({"Slide": sn, "Shapes": len(slide.shapes), "TextObjects": texts,
                            "Charts": charts, "Tables": tables, "Pictures": pictures})
        totals["charts"] += charts
        totals["tables"] += tables
        totals["pictures"] += pictures
        totals["text"] += texts

    # Gate V1.2.1 regressions: the complete Slide 18 header must be styled,
    # including the second paragraph created by each explicit newline.
    slide18 = prs.slides[17]
    parameter_tables = [shape for shape in slide18.shapes
                        if getattr(shape, "has_table", False) and shape.name == "ParameterTable"]
    assert len(parameter_tables) == 1, "slide 18: ParameterTable missing or duplicated"
    for c, cell in enumerate(parameter_tables[0].table.rows[0].cells):
        expected_alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        assert cell.vertical_anchor == MSO_ANCHOR.MIDDLE, f"slide 18 header c{c}: vertical alignment"
        for para in cell.text_frame.paragraphs:
            assert para.alignment == expected_alignment, f"slide 18 header c{c}: horizontal alignment"
            runs = [run for run in para.runs if run.text]
            assert runs, f"slide 18 header c{c}: empty paragraph"
            for run in runs:
                assert run.font.name == FONT, f"slide 18 header c{c}: font {run.font.name}"
                assert run.font.size and run.font.size.pt == FIGURE_MIN_PT, \
                    f"slide 18 header c{c}: size {run.font.size}"
                assert run.font.bold is True, f"slide 18 header c{c}: not bold"
                assert str(run.font.color.rgb) == C["white"], \
                    f"slide 18 header c{c}: color {run.font.color.rgb}"

    # Slide 20 deliberately has no native chart legend.  Its replacement is a
    # two-item collection of editable native markers and text boxes.
    slide20 = prs.slides[19]
    kernel_charts = [shape.chart for shape in slide20.shapes if getattr(shape, "has_chart", False)]
    assert len(kernel_charts) == 1, "slide 20: expected one chart"
    assert not kernel_charts[0].has_legend, "slide 20: native legend must be disabled"
    marker_names = sorted(shape.name for shape in slide20.shapes
                          if shape.name.startswith("KernelLegendMarker_"))
    label_shapes = sorted((shape for shape in slide20.shapes
                           if shape.name.startswith("KernelLegendLabel_")), key=lambda shape: shape.name)
    assert marker_names == ["KernelLegendMarker_1", "KernelLegendMarker_2"], marker_names
    assert [shape.text for shape in label_shapes] == ["Shared kernel", "Block kernel"], \
        [shape.text for shape in label_shapes]

    with zipfile.ZipFile(prs_path) as zf:
        names = zf.namelist()
        raster = [n for n in names if n.startswith("ppt/media/")
                  and n.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff"))]
        chart_xml = [n for n in names if re.fullmatch(r"ppt/charts/chart\d+\.xml", n)]
        embeddings = [n for n in names if n.startswith("ppt/embeddings/")]
        problems["raster"] += raster
        # Chart interiors carry category, series, and axis-title text.
        for name in chart_xml:
            text = zf.read(name).decode("utf-8")
            for chunk in re.findall(r"<a:t>([^<]*)</a:t>", text) + re.findall(r"<c:v>([^<]*)</c:v>", text):
                check_visible(problems, name, "chart interior", chunk)
        visible_parts = scan_package_language(zf, problems)
        axis_ids = {}
        for name in chart_xml:
            text = zf.read(name).decode("utf-8")
            ids = [int(v) for v in re.findall(r'<c:axId val="(-?\d+)"/>', text)]
            cross = [int(v) for v in re.findall(r'<c:crossAx val="(-?\d+)"/>', text)]
            assert all(v >= 0 for v in ids + cross), f"{name}: negative axis id"
            assert set(cross) <= set(ids), f"{name}: dangling crossAx"
            axis_ids[name] = sorted(set(ids))
        flat = [v for ids in axis_ids.values() for v in ids]
        assert len(flat) == len(set(flat)), "axis ids collide across charts"
        assert len(chart_xml) == totals["charts"]
        assert len(embeddings) >= totals["charts"]

    for key, found in problems.items():
        assert not found, f"{key}: {found}"

    check_required_sentences(prs)
    notes_report = check_notes_bilingual(prs)

    # Canonical values must survive into the rendered deck.
    assert values["speedups"] == [3.17, 1.31, 1.51, 1.45], values["speedups"]
    assert values["tuned_batches"] == [2048, 64, 64, 32], values["tuned_batches"]
    assert values["hwa_corrected"] == [1.4767, 1.1012, 1.5563], values["hwa_corrected"]
    assert values["hwa_synth"] == [1.6787, 1.0661, 1.3914], values["hwa_synth"]
    assert values["memory"] == [
        ("pure_b4096", "Success"), ("pure_b8192", "CUDA OOM"),
        ("um_b10240", "Success"), ("um_b12288", "Cgroup host-memory OOM kill"),
        ("chunked_b16384", "Success")], values["memory"]
    assert values["tier_a"] == 3 and values["tier_b"] == 10 and values["tier_total"] == 13
    assert values["mismatch_total"] == 0 and values["missing_total"] == 0
    assert values["byte_identical"] == {"No"}
    assert values["kernel"] == [
        ("roadNet-PA", 1063.712326, 701.573311, 1.52),
        ("roadNet-TX", 1639.164633, 984.587390, 1.66),
    ], values["kernel"]

    return {"slides": len(prs.slides), "objects": totals["objects"], "charts": totals["charts"],
            "tables": totals["tables"], "pictures": totals["pictures"], "text_objects": totals["text"],
            "raster_only_slides": 0, "slide_stats": slide_stats, "axis_ids": axis_ids,
            "visible_parts": visible_parts, "notes_report": notes_report}


RUNTIME_CACHE: list[dict] = []


def collect_values(data, t5_rows):
    speedups = [round(r["Speedup"], 2) for r in data["speedup"]]
    tuned = [r["PathMerge_Tuned_Batch"] for r in data["runtime"]]
    lookup = {(r["Graph"], r["Factor"]): r["Main_Effect"] for r in data["ablation"]}
    corrected = [round(lookup[("325557_3216152_corrected_v1", f)], 4) for f in ("H", "W", "A")]
    synth = [round(lookup[("Synthetic-4 aggregate", f)], 4) for f in ("H", "W", "A")]
    tier_a = sum(1 for r in t5_rows if r["EvidenceTier"] == "Independent CPU reference")
    tier_b = sum(1 for r in t5_rows if r["EvidenceTier"] == "Cross-implementation consistency")
    kernel_lookup = {(r["Graph"], r["Kernel"]): r["Median_Runtime_s"] for r in data["kernel"]}
    kernel = [
        (graph, kernel_lookup[(graph, "shared")], kernel_lookup[(graph, "block")],
         round(kernel_lookup[(graph, "shared")] / kernel_lookup[(graph, "block")], 2))
        for graph in ("roadNet-PA", "roadNet-TX")
    ]
    return {
        "speedups": speedups,
        "tuned_batches": tuned,
        "hwa_corrected": corrected,
        "hwa_synth": synth,
        "memory": [(r["Config"], r["Status"]) for r in data["memory"]],
        "tier_a": tier_a,
        "tier_b": tier_b,
        "tier_total": len(t5_rows),
        "mismatch_total": sum(int(r["MismatchedElements"]) for r in t5_rows),
        "missing_total": sum(int(r["MissingIndices"]) for r in t5_rows),
        "byte_identical": {r["ByteIdentical"] for r in t5_rows},
        "kernel": kernel,
    }


def build(prs_path, data, t5_rows):
    prs = Presentation()
    prs.slide_width = Inches(lib.SLIDE_W)
    prs.slide_height = Inches(lib.SLIDE_H)
    slide_title(prs)
    slide_agenda(prs)
    slide_bc_graph(prs)
    slide_brandes(prs)
    slide_gpu_contribution(prs)
    slide_framework(prs)
    slide_batching(prs, data["design"])
    slide_variants(prs)
    slide_eval_design(prs)
    slide_runtime(prs, data["runtime"])
    slide_speedup(prs, data["speedup"])
    slide_ablation(prs, data["ablation"])
    slide_memory(prs, data["memory"])
    slide_limits(prs)
    slide_contributions(prs)
    slide_conclusion(prs)
    backup_environment(prs)
    backup_parameters(prs, data["runtime"], data["design"])
    backup_sweep(prs, data["sweep"])
    backup_kernel(prs, data["kernel"])
    backup_phase(prs, data["phase"])
    backup_correctness(prs, t5_rows)
    backup_historical(prs)
    assert len(prs.slides) == MAIN_SLIDES + BACKUP_SLIDES
    attach_notes(prs)
    index = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_chart", False):
                normalize_axis_ids(shape.chart, AXIS_ID_BASE + 500 + index * 10)
                index += 1
    prs.save(prs_path)


def main():
    EDITABLE.mkdir(parents=True, exist_ok=True)
    validate_notes_content()
    data = load_all()
    RUNTIME_CACHE.extend(data["runtime"])
    t5_rows = lib.read_tsv("result/tables/thesis/T5_correctness_summary.tsv")
    values = collect_values(data, t5_rows)

    build(PPTX_PATH, data, t5_rows)
    normalize_in_place(PPTX_PATH)
    stats = validate(PPTX_PATH, 23, values)

    total_seconds = write_notes_document()
    write_plan_document(total_seconds)
    write_manifest(stats)
    write_readme(stats, total_seconds)

    pacing = notes_pacing()
    print("Gate V1.2.1 presentation generated and validated")
    print(f"main_slides\t{MAIN_SLIDES}")
    print(f"backup_slides\t{BACKUP_SLIDES}")
    for key in ("slides", "objects", "text_objects", "charts", "tables", "pictures", "raster_only_slides"):
        print(f"{key}\t{stats[key]}")
    print("presentation_timing\t[USER TO EDIT]")
    print(f"visible_japanese_chars\t0")
    print(f"visible_parts_scanned\t{len(stats['visible_parts'])}")
    print(f"english_script_slides\t{len(stats['notes_report'])}/23")
    print(f"japanese_script_slides\t{len(stats['notes_report'])}/23")
    print(f"notes_path\t{NOTES_PATH.name}")
    print(f"speedups\t{values['speedups']}")
    print(f"tuned_batches\t{values['tuned_batches']}")
    print(f"hwa_corrected\t{values['hwa_corrected']}")
    print(f"hwa_synth\t{values['hwa_synth']}")
    print(f"correctness\ttierA={values['tier_a']} tierB={values['tier_b']} total={values['tier_total']} "
          f"mismatch={values['mismatch_total']} missing={values['missing_total']} byte={sorted(values['byte_identical'])}")
    for name, ids in stats["axis_ids"].items():
        print(f"axis_ids\t{name}\t{ids[0]},{ids[1]}")
    for row in stats["slide_stats"]:
        print("slide\t" + "\t".join(f"{k}={v}" for k, v in row.items()))


if __name__ == "__main__":
    try:
        main()
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        raise
